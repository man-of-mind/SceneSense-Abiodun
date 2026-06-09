"""Build the 24-slide SceneSense Agent intro deck.

Run: python3 build_scenesense_intro_deck.py
Output: SceneSense_Agent_Intro_Deck.pptx in the same folder.

Story arc:
  [0] Title
  [1] Outline/Agenda
  --- Section 1: Context & Motivation ---
  [2] Introduction  – the physical-AI era
  [3] Motivation    – why perception alone isn't enough
  --- Section 2: The problem ---
  [4] The problem: occlusion
  [5] Cooperative perception concept
  [6] Use cases
  [7] What to share? (sharing options)
  --- Section 3: Split inferencing ---
  [8]  What is split inferencing?
  [9]  Why split inferencing?
  [10] Methodology / architecture (OD + SEG)
  [11] Compression knobs
  [12] Baseline measurements (real data)
  --- Section 4: SceneSense Agent ---
  [13] Bridge: single-UE → multi-UE
  [14] SceneSense Agent — core idea
  [15] Hypothesis and goal
  [16] Things we have to get right
  --- Section 5: Architecture breakdown ---
  [17] Big picture architecture
  [18] Breakdown: UE-side RL agent
  [19] Breakdown: Edge AI server + Shared Spatial Map
  [20] Breakdown: Map-sharing agent + guardrails
  --- Section 6: Plan ---
  [21] Evaluation campaigns
  [22] Timeline
  [23] Risks, open questions, next steps
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# -------------------- Palette --------------------
NAVY       = RGBColor(0x0B, 0x3D, 0x91)
ID_BLUE    = RGBColor(0x00, 0xA9, 0xE0)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xEF, 0xEF, 0xEF)
SOFT_BG    = RGBColor(0xF7, 0xF9, 0xFC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE     = RGBColor(0xEA, 0x8A, 0x1F)
GREEN      = RGBColor(0x2E, 0x8B, 0x57)
RED        = RGBColor(0xC0, 0x39, 0x2B)
PURPLE     = RGBColor(0x6A, 0x4C, 0x93)
TEAL       = RGBColor(0x00, 0x7E, 0x7F)

# Section colours (used on the outline slide and section dividers)
SECTION_COLORS = {
    1: ID_BLUE,
    2: ORANGE,
    3: GREEN,
    4: PURPLE,
    5: RED,
    6: TEAL,
}
SECTION_LABELS = {
    1: "Context & Motivation",
    2: "The Problem",
    3: "Split Inferencing",
    4: "SceneSense Agent",
    5: "Architecture",
    6: "Plan",
}

# -------------------- Presentation setup --------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ================================================================
#  Core helpers
# ================================================================

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_text_frame(tf, lines, default_size=14, default_color=DARK_GRAY,
                     word_wrap=True):
    tf.word_wrap = word_wrap
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(ln, dict):
            p.alignment = ln.get("align", PP_ALIGN.LEFT)
            if "space_after" in ln:
                p.space_after = Pt(ln["space_after"])
            if "space_before" in ln:
                p.space_before = Pt(ln["space_before"])
            run = p.add_run()
            run.text = ln.get("text", "")
            run.font.name = ln.get("font", "Calibri")
            run.font.size = Pt(ln.get("size", default_size))
            run.font.bold = ln.get("bold", False)
            run.font.italic = ln.get("italic", False)
            run.font.color.rgb = ln.get("color", default_color)
        else:
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = ln
            run.font.name = "Calibri"
            run.font.size = Pt(default_size)
            run.font.color.rgb = default_color


def add_text_box(slide, left, top, width, height,
                 text, font_size=14, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(9000)
    tf.margin_bottom = Emu(9000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rich_tb(slide, left, top, width, height, lines,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    _fill_text_frame(tf, lines)
    return tb


def add_rect(slide, left, top, width, height,
             fill=WHITE, line=NAVY, line_width=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    return s


def set_shape_text(shape, text, size=12, bold=True, color=DARK_GRAY,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_label(slide, left, top, width, height, text,
              size=12, bold=True, color=DARK_GRAY, italic=False,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(18000)
    tf.margin_right = Emu(18000)
    tf.margin_top = Emu(9000)
    tf.margin_bottom = Emu(9000)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, line_width=2.0):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(line_width)
    ln = connector.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return connector


def add_bullets(slide, left, top, width, height, items,
                size=14, color=DARK_GRAY):
    """items: list of str OR (head_str, body_str) tuples."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(9000)
    tf.margin_bottom = Emu(9000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        if isinstance(item, tuple):
            head, body = item
            run = p.add_run()
            run.text = "•  "
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
            rh = p.add_run()
            rh.text = head + " "
            rh.font.bold = True
            rh.font.size = Pt(size)
            rh.font.color.rgb = NAVY
            rh.font.name = "Calibri"
            rb = p.add_run()
            rb.text = body
            rb.font.size = Pt(size)
            rb.font.color.rgb = color
            rb.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = "•  " + item
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return tb


def add_chrome(slide, title, subtitle=None, section=None):
    """Standard title bar + optional section tag + footer."""
    # Top accent strip
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = (
        SECTION_COLORS.get(section, ID_BLUE) if section else ID_BLUE)
    strip.line.fill.background()
    # Section badge (top-right)
    if section is not None:
        sc = SECTION_COLORS[section]
        badge = add_rect(slide, Inches(11.0), Inches(0.1),
                         Inches(2.1), Inches(0.4),
                         fill=sc, line=None,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        set_shape_text(badge,
                       f"§{section}  {SECTION_LABELS[section]}",
                       size=9, bold=True, color=WHITE)
    # Title
    title_h = Inches(0.55) if subtitle else Inches(0.7)
    add_text_box(slide, Inches(0.4), Inches(0.15), Inches(10.4), title_h,
                 title, font_size=26, bold=True, color=NAVY)
    if subtitle:
        add_text_box(slide, Inches(0.4), Inches(0.72), Inches(12.5),
                     Inches(0.42), subtitle,
                     font_size=13, italic=True, color=MED_GRAY)
    # Footer
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.32), SLIDE_W, Inches(0.03))
    footer.fill.solid()
    footer.fill.fore_color.rgb = SECTION_COLORS.get(section, ID_BLUE) if section else ID_BLUE
    footer.line.fill.background()
    add_text_box(slide, Inches(0.4), Inches(7.18), Inches(12.5),
                 Inches(0.25),
                 "©2026 InterDigital, Inc.  ·  SceneSense Agent — intern presentation",
                 font_size=9, italic=True, color=MED_GRAY)


# ================================================================
#  SLIDE 1  — Title
# ================================================================
def slide_title():
    s = add_blank_slide()
    bg = add_rect(s, 0, 0, SLIDE_W, SLIDE_H,
                  fill=NAVY, line=None, shape=MSO_SHAPE.RECTANGLE)
    # Bottom accent bar
    acc = add_rect(s, 0, Inches(6.05), SLIDE_W, Inches(0.22),
                   fill=ID_BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)

    add_text_box(s, Inches(0.7), Inches(1.2), Inches(12.0), Inches(1.2),
                 "SceneSense Agent",
                 font_size=56, bold=True, color=WHITE)
    add_text_box(s, Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.65),
                 "Agent-Controlled Split Inference for Network-aware",
                 font_size=23, color=ID_BLUE)
    add_text_box(s, Inches(0.7), Inches(2.85), Inches(12.0), Inches(0.65),
                 "Cooperative Perception over Shared Spatial Maps",
                 font_size=23, color=ID_BLUE)

    add_text_box(s, Inches(0.7), Inches(4.2), Inches(12.0), Inches(0.55),
                 "Abiodun Ganiyu  ·  IDCC × NEU 6-Month Internship",
                 font_size=20, bold=True, color=WHITE)
    add_text_box(s, Inches(0.7), Inches(4.72), Inches(12.0), Inches(0.5),
                 "Supervised by Subhramoy Mohanti  ·  June 2026",
                 font_size=16, color=WHITE)
    add_text_box(s, Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.35),
                 "©2026 InterDigital, Inc. All Rights Reserved.",
                 font_size=10, italic=True, color=WHITE)
    return s


# ================================================================
#  SLIDE 2  — Outline
# ================================================================
def slide_outline():
    s = add_blank_slide()
    add_chrome(s, "Outline")
    add_text_box(s, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.4),
                 "What this presentation covers — and in what order",
                 font_size=15, italic=True, color=MED_GRAY)

    sections = [
        (1, "Context & Motivation",
         ["The physical-AI era",
          "Why perception alone isn't enough",
          "Motivation for cooperative perception"]),
        (2, "The Problem",
         ["Occlusion as a concrete failure mode",
          "Cooperative perception concept",
          "Use cases and deployment contexts"]),
        (3, "Split Inferencing",
         ["What it is and why it's the right transport",
          "Architecture (OD + SEG)",
          "Compression knobs and baseline measurements"]),
        (4, "SceneSense Agent",
         ["Core idea and research hypothesis",
          "From single-UE to multi-UE coordination",
          "Things we have to get right"]),
        (5, "Architecture Breakdown",
         ["UE-side RL agent",
          "Edge AI server + Shared Spatial Map",
          "Map-sharing agent and guardrails"]),
        (6, "Plan",
         ["Evaluation campaigns",
          "Six-month timeline",
          "Risks and next steps"]),
    ]

    cw = Inches(4.18)
    ch = Inches(2.65)
    gx = Inches(0.13)
    gy = Inches(0.1)
    x0 = Inches(0.4)
    y0 = Inches(1.55)

    for i, (sec, head, items) in enumerate(sections):
        r = i // 3
        c = i % 3
        left = x0 + (cw + gx) * c
        top  = y0 + (ch + gy) * r
        col  = SECTION_COLORS[sec]

        card = add_rect(s, left, top, cw, ch,
                        fill=WHITE, line=col, line_width=2)
        band = add_rect(s, left, top, cw, Inches(0.52),
                        fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        set_shape_text(band, f"§{sec}  {head}", size=13, bold=True,
                       color=WHITE)
        for j, item in enumerate(items):
            add_text_box(s,
                         left + Inches(0.22),
                         top + Inches(0.62) + Inches(j * 0.63),
                         cw - Inches(0.44), Inches(0.6),
                         f"  ›  {item}", font_size=12, color=DARK_GRAY)
    return s


# ================================================================
#  SLIDE 3  — Introduction: the physical-AI era
# ================================================================
def slide_introduction():
    s = add_blank_slide()
    add_chrome(s,
               "Introduction — the physical-AI era",
               "Machines that need to understand the physical world in real time",
               section=1)

    # Left column: context text
    add_rich_tb(s, Inches(0.5), Inches(1.45), Inches(7.1), Inches(5.8),
                [
                    {"text": "The next wave of 5G/6G services",
                     "size": 19, "bold": True, "color": NAVY,
                     "space_after": 8},
                    {"text": "3GPP Release 18+ is bringing a new class of applications to the mobile network: physical-AI services. These are systems that perceive, reason about, and act in the physical world — in real time, over a shared wireless infrastructure.",
                     "size": 13, "color": DARK_GRAY, "space_after": 16},

                    {"text": "Three converging trends",
                     "size": 17, "bold": True, "color": NAVY,
                     "space_after": 6},
                    {"text": "•  Connected and autonomous vehicles (CAVs) — share perception to navigate safely",
                     "size": 13, "color": DARK_GRAY, "space_after": 4},
                    {"text": "•  Infrastructure sensing — roadside cameras and smart intersections observe hazards before any vehicle can",
                     "size": 13, "color": DARK_GRAY, "space_after": 4},
                    {"text": "•  Semantic machine communication — instead of streaming pixels, machines exchange task-relevant representations (features, masks, object evidence)",
                     "size": 13, "color": DARK_GRAY, "space_after": 16},

                    {"text": "The challenge",
                     "size": 17, "bold": True, "color": NAVY,
                     "space_after": 6},
                    {"text": "These services will not scale if every machine streams everything all the time. They also cannot safely compress away vulnerable objects — pedestrians, cyclists, children — just because the link is congested.",
                     "size": 13, "color": DARK_GRAY, "italic": True},
                ])

    # Right column: three stat-style boxes
    stats = [
        ("~70%", "of fatal urban collisions involve a vehicle failing to perceive a pedestrian obscured from its direct viewpoint", NAVY),
        ("~4.6×", "more bytes required for segmentation features vs detection — static compression strategies cannot serve both tasks", ORANGE),
        ("5QI 85–88", "3GPP targets for split AI/ML: 5–10 ms PDB, 255B–1.1 KB MDBV — today's uncompressed feature tensors exceed this by 100–1600×", ID_BLUE),
    ]
    ry = Inches(1.55)
    for head, body, col in stats:
        box = add_rect(s, Inches(7.85), ry, Inches(5.1), Inches(1.7),
                       fill=SOFT_BG, line=col, line_width=2)
        add_text_box(s, Inches(7.95), ry + Inches(0.08),
                     Inches(1.35), Inches(1.55),
                     head, font_size=34, bold=True, color=col,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(s, Inches(9.35), ry + Inches(0.12),
                     Inches(3.5), Inches(1.45),
                     body, font_size=11, color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
        ry += Inches(1.85)
    return s


# ================================================================
#  SLIDE 4  — Motivation
# ================================================================
def slide_motivation():
    s = add_blank_slide()
    add_chrome(s,
               "Motivation — why existing approaches fall short",
               "Three gaps that SceneSense Agent is designed to close",
               section=1)

    gaps = [
        (
            "Gap 1 — Single-sensor perception has hard limits",
            "No amount of local compute or model quality can detect an object that is completely occluded from the sensor. A vehicle behind a parked truck cannot see a child stepping off the kerb. Infrastructure sensors and nearby vehicles *already have* that information — it just isn't shared.",
            "SOTA today:",
            "Most cooperative-perception research (CoDriving, Coopernaut, Where2comm) treats the network as a black box and does not optimize for the shared wireless resource.",
            ID_BLUE,
        ),
        (
            "Gap 2 — Static compression knobs waste accuracy or bandwidth",
            "Choosing AE channels, ROI threshold, or quantization level once and leaving it fixed produces: (a) over-spending when the scene is simple or the link is good; (b) destroying task-critical accuracy when the scene is complex or the link is bad.",
            "SOTA today:",
            "Rate–distortion sweeps show the Pareto frontier, but no real-time agent selects the right operating point from scene + network + model state.",
            ORANGE,
        ),
        (
            "Gap 3 — Perception and the network do not talk to each other",
            "The 5G scheduler decides resource allocation independently of whether the current frame contains a pedestrian about to be hit. The perception system ships features independently of whether the link is about to drop them anyway. This mismatch wastes both spectrum and safety.",
            "SCAN-AI (Mohanti et al.) closes this for single-UE video uplink. The multi-UE cooperative extension — where vehicles compete for shared radio resources while trying to fuse each other's perception — is the open problem.",
            "",
            GREEN,
        ),
    ]

    y = Inches(1.45)
    h = Inches(1.75)
    gap_y = Inches(0.12)
    for i, (head, body, label, detail, col) in enumerate(gaps):
        top = y + (h + gap_y) * i
        card = add_rect(s, Inches(0.45), top, Inches(12.45), h,
                        fill=WHITE, line=col, line_width=1.5)
        # Left number band
        nb = add_rect(s, Inches(0.45), top, Inches(0.55), h,
                      fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        set_shape_text(nb, str(i + 1), size=24, bold=True, color=WHITE)
        # Heading
        add_text_box(s, Inches(1.1), top + Inches(0.08),
                     Inches(11.6), Inches(0.4),
                     head, font_size=15, bold=True, color=NAVY)
        # Body
        add_text_box(s, Inches(1.1), top + Inches(0.47),
                     Inches(7.5), h - Inches(0.55),
                     body, font_size=12, color=DARK_GRAY)
        # SOTA tag
        if label:
            add_text_box(s, Inches(8.8), top + Inches(0.47),
                         Inches(0.95), Inches(0.4),
                         label, font_size=10, bold=True, color=col)
            add_text_box(s, Inches(8.8), top + Inches(0.82),
                         Inches(3.95), h - Inches(0.9),
                         detail, font_size=11, italic=True,
                         color=DARK_GRAY)

    # Bottom tagline
    add_text_box(s, Inches(0.45), Inches(6.8), Inches(12.45), Inches(0.4),
                 "SceneSense Agent addresses all three: learned policies + task-precision guardrails + a shared physical-world map.",
                 font_size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    return s


# ================================================================
#  SLIDE 5  — The problem: occlusion
# ================================================================
def slide_problem_occlusion():
    s = add_blank_slide()
    add_chrome(s,
               "The problem — a single car cannot see everything",
               "Perception is bounded by occlusion and viewpoint, not just compute",
               section=2)

    # Left: road sketch
    sx = Inches(0.5)
    sy = Inches(1.5)
    sw = Inches(6.5)
    sh = Inches(5.3)
    canvas = add_rect(s, sx, sy, sw, sh, fill=SOFT_BG,
                      line=LIGHT_GRAY, shape=MSO_SHAPE.RECTANGLE)
    # Road
    road = add_rect(s, sx + Inches(0.3), sy + Inches(2.8),
                    sw - Inches(0.6), Inches(1.15),
                    fill=RGBColor(0x55, 0x55, 0x55), line=None,
                    shape=MSO_SHAPE.RECTANGLE)
    # Lane dashes
    for k in range(8):
        seg = add_rect(s, sx + Inches(0.5 + k * 0.75), sy + Inches(3.3),
                       Inches(0.4), Inches(0.07),
                       fill=WHITE, line=None, shape=MSO_SHAPE.RECTANGLE)
    # Sidewalk
    sw_box = add_rect(s, sx + Inches(0.3), sy + Inches(1.85),
                      sw - Inches(0.6), Inches(0.9),
                      fill=RGBColor(0xCC, 0xCC, 0xCC), line=None,
                      shape=MSO_SHAPE.RECTANGLE)
    # Ego car
    ego = add_rect(s, sx + Inches(0.6), sy + Inches(2.9),
                   Inches(1.15), Inches(0.65),
                   fill=ID_BLUE, line=NAVY,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    set_shape_text(ego, "EGO", size=12, color=WHITE)

    # Parked truck (occluder)
    truck = add_rect(s, sx + Inches(2.7), sy + Inches(2.05),
                     Inches(1.6), Inches(1.55),
                     fill=DARK_GRAY, line=DARK_GRAY,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    set_shape_text(truck, "PARKED\nTRUCK", size=12, color=WHITE)
    add_text_box(s, sx + Inches(2.3), sy + Inches(0.15),
                 Inches(2.5), Inches(0.4),
                 "occluder", font_size=11, bold=True, color=DARK_GRAY,
                 align=PP_ALIGN.CENTER)

    # Pedestrian (hidden)
    ped = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              sx + Inches(3.1), sy + Inches(1.35),
                              Inches(0.35), Inches(0.55))
    ped.fill.solid()
    ped.fill.fore_color.rgb = RED
    ped.line.fill.background()
    add_text_box(s, sx + Inches(2.65), sy + Inches(0.85),
                 Inches(1.2), Inches(0.4),
                 "pedestrian\n(hidden!)", font_size=11, bold=True,
                 color=RED, align=PP_ALIGN.CENTER)

    # Ego's limited view cone
    add_text_box(s, sx + Inches(1.85), sy + Inches(3.05),
                 Inches(0.9), Inches(0.55),
                 "ego\nview →", font_size=10, italic=True, color=ID_BLUE,
                 align=PP_ALIGN.CENTER)

    # ✗ symbol
    add_text_box(s, sx + Inches(2.1), sy + Inches(2.2),
                 Inches(0.55), Inches(0.55),
                 "✗", font_size=28, bold=True, color=RED,
                 align=PP_ALIGN.CENTER)

    # Other car on right (can see)
    car2 = add_rect(s, sx + Inches(4.8), sy + Inches(3.0),
                    Inches(1.1), Inches(0.6),
                    fill=ORANGE, line=DARK_GRAY,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    set_shape_text(car2, "OTHER UE", size=10, color=WHITE)
    add_text_box(s, sx + Inches(4.0), sy + Inches(1.65),
                 Inches(2.3), Inches(0.45),
                 "← other UE can see it!",
                 font_size=11, italic=True, color=GREEN,
                 align=PP_ALIGN.CENTER)

    # Right: bullet points
    add_text_box(s, Inches(7.3), Inches(1.5), Inches(5.6), Inches(0.5),
                 "What's happening here", font_size=19, bold=True,
                 color=NAVY)
    add_bullets(s, Inches(7.3), Inches(2.05), Inches(5.6), Inches(5.0),
                [
                    ("Ego drives forward.", "Cannot detect the pedestrian — they are fully hidden behind the parked truck."),
                    ("Pedestrian steps into road.", "By the time they become visible, it may be too late to stop."),
                    ("The other UE can see them.", "It has a direct line of sight — and its feature stream already contains the pedestrian's evidence."),
                    ("This scenario is repeatable.", "Parked trucks at intersections, children between cars, cyclists emerging from driveways — all the same failure mode."),
                    ("Cooperative perception solves it.", "If the other UE's perception feeds a shared spatial map, the ego can be warned before the collision occurs."),
                ], size=13)
    return s


# ================================================================
#  SLIDE 6  — Cooperative perception concept
# ================================================================
def slide_coop_concept():
    s = add_blank_slide()
    add_chrome(s,
               "Cooperative perception — many partial views, one shared world model",
               "Each vehicle contributes what it can see; the map aggregates everything",
               section=2)

    ue_colors = [ID_BLUE, ORANGE, GREEN]
    # UEs on left, arrows to shared network, shared map on right
    for i in range(3):
        cy = Inches(1.8 + i * 1.6)
        car = add_rect(s, Inches(0.5), cy, Inches(1.5), Inches(0.9),
                       fill=ue_colors[i], line=DARK_GRAY,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        set_shape_text(car, f"Vehicle {i+1}\nUE {i+1}", size=12, color=WHITE)
        add_text_box(s, Inches(0.5), cy + Inches(1.0), Inches(1.5),
                     Inches(0.35),
                     ["sees cars/kerb ahead", "side-angle of crossing", "rear view + overtaking"][i],
                     font_size=9, italic=True, color=MED_GRAY,
                     align=PP_ALIGN.CENTER)
        add_arrow(s, Inches(2.0), cy + Inches(0.45),
                  Inches(3.9), Inches(3.75),
                  color=ue_colors[i], line_width=2.25)

    # gNB
    gnb = add_rect(s, Inches(3.9), Inches(3.3), Inches(1.7), Inches(0.95),
                   fill=WHITE, line=NAVY, line_width=2)
    set_shape_text(gnb, "OAI gNB\n5G uplink", size=12, bold=True, color=NAVY)

    # Arrow to edge server
    add_arrow(s, Inches(5.6), Inches(3.75), Inches(6.4), Inches(3.75),
              color=NAVY, line_width=2.5)

    # Edge fusion box
    fuse = add_rect(s, Inches(6.4), Inches(2.85), Inches(3.0), Inches(1.85),
                    fill=SOFT_BG, line=NAVY, line_width=2)
    set_shape_text(fuse, "Edge AI server\n(back-half inference\n+ confidence tagging)",
                   size=12, bold=True, color=NAVY)

    # Arrow to shared map
    add_arrow(s, Inches(9.4), Inches(3.75), Inches(10.1), Inches(3.75),
              color=ORANGE, line_width=2.5)

    # Shared spatial map
    smap = add_rect(s, Inches(10.1), Inches(2.65), Inches(3.0), Inches(2.25),
                    fill=SOFT_BG, line=ORANGE, line_width=2)
    add_text_box(s, Inches(10.1), Inches(2.68), Inches(3.0), Inches(0.45),
                 "Shared Spatial Map", font_size=14, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(10.2), Inches(3.15), Inches(2.8), Inches(1.7),
                 "class · pose · confidence\nfreshness · occlusion state\nprovenance (which UE)",
                 font_size=11, italic=True, color=DARK_GRAY,
                 align=PP_ALIGN.CENTER)

    # Downlink arrows (map updates to vehicles)
    for i in range(3):
        cy = Inches(2.25 + i * 1.6)
        add_arrow(s, Inches(10.1), Inches(3.75),
                  Inches(2.05), cy + Inches(0.45),
                  color=MED_GRAY, line_width=1.5)

    # Bottom callout
    add_text_box(s, Inches(0.5), Inches(6.55), Inches(12.35), Inches(0.6),
                 "Key insight: the shared map is not raw video and not final detections — it is task-useful evidence that any vehicle can act on, even if it never had direct line of sight.",
                 font_size=13, italic=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    return s


# ================================================================
#  SLIDE 7  — Use cases
# ================================================================
def slide_use_cases():
    s = add_blank_slide()
    add_chrome(s,
               "Where this matters — use cases",
               "Same pattern: spatially separated sensors, partial views, shared evidence map",
               section=2)

    cases = [
        ("Connected vehicles",
         "Cars share blind-spot perception around occluding trucks and intersections. Ego reacts before it can even see the hazard.",
         ID_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Smart intersections",
         "Roadside cameras observe pedestrians / cyclists before vehicles arrive. Warnings shared to approaching traffic.",
         ORANGE, MSO_SHAPE.RECTANGLE),
        ("Drone-assisted sensing",
         "Overhead view from UAVs shared with ground vehicles for occluded terrain and obstacles.",
         GREEN, MSO_SHAPE.OCTAGON),
        ("Warehouse robots",
         "Fleet shares object locations under shelving and partial occlusions. Collision-free paths computed with shared map.",
         PURPLE, MSO_SHAPE.PENTAGON),
        ("Industrial safety",
         "Machines share safety-zone awareness; humans entering hazardous areas trigger shared stop commands.",
         RED, MSO_SHAPE.HEXAGON),
        ("Multi-camera security",
         "Distributed cameras fuse detections of persons/vehicles across a site into one coherent spatial model.",
         TEAL, MSO_SHAPE.DIAMOND),
    ]
    cw = Inches(4.15)
    ch = Inches(2.65)
    gx = Inches(0.14)
    gy = Inches(0.15)
    x0 = Inches(0.45)
    y0 = Inches(1.48)
    for idx, (head, body, col, ish) in enumerate(cases):
        r = idx // 3
        c = idx % 3
        lx = x0 + (cw + gx) * c
        ty = y0 + (ch + gy) * r
        add_rect(s, lx, ty, cw, ch, fill=WHITE,
                 line=LIGHT_GRAY, line_width=1.5)
        icon = s.shapes.add_shape(ish, lx + Inches(0.22), ty + Inches(0.22),
                                   Inches(0.65), Inches(0.65))
        icon.fill.solid()
        icon.fill.fore_color.rgb = col
        icon.line.fill.background()
        add_text_box(s, lx + Inches(1.05), ty + Inches(0.22),
                     cw - Inches(1.15), Inches(0.55),
                     head, font_size=15, bold=True, color=NAVY)
        add_text_box(s, lx + Inches(0.22), ty + Inches(1.0),
                     cw - Inches(0.44), ch - Inches(1.1),
                     body, font_size=12, color=DARK_GRAY)

    add_text_box(s, Inches(0.45), Inches(7.05), Inches(12.45), Inches(0.35),
                 "Our research focus: connected vehicles + smart intersections. The same approach extends to all six.",
                 font_size=11, italic=True, color=MED_GRAY,
                 align=PP_ALIGN.CENTER)
    return s


# ================================================================
#  SLIDE 8  — What to share?
# ================================================================
def slide_sharing_options():
    s = add_blank_slide()
    add_chrome(s,
               "What should the vehicles share?",
               "Three obvious options — only one hits the sweet spot for safety and scalability",
               section=2)

    opts = [
        ("Raw video / camera frames",
         ["Highest fidelity for fusion",
          "Huge bandwidth (Mbps per stream, per UE)",
          "Bursty; occupies entire uplink at fleet scale"],
         "Too big. Cannot fit shared 5G spectrum.",
         RED, "✗"),
        ("Final detections / boxes only",
         ["Tiny payload",
          "Easy to standardize",
          "Receiver cannot re-fuse, estimate confidence, or recover context"],
         "Too thin. Information is already gone.",
         ORANGE, "−"),
        ("Intermediate feature tensors",
         ["Already compact, task-aligned summary",
          "Rich enough for the server head to re-run and fuse",
          "Compressible: AE / ROI / quantization cut 2–10× further"],
         "Sweet spot for safety-critical cooperative perception over 5G.",
         GREEN, "✓"),
    ]
    x = Inches(0.45)
    w = Inches(4.18)
    h = Inches(5.5)
    gap = Inches(0.13)
    y = Inches(1.48)
    for i, (head, bullets, verdict, col, mark) in enumerate(opts):
        left = x + (w + gap) * i
        add_rect(s, left, y, w, h, fill=WHITE, line=col, line_width=1.5)
        band = add_rect(s, left, y, w, Inches(0.7),
                        fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        set_shape_text(band, f"{mark}  {head}", size=14, bold=True,
                       color=WHITE)
        add_bullets(s, left, y + Inches(0.85), w, Inches(3.3),
                    bullets, size=13)
        vb = add_rect(s, left + Inches(0.2), y + Inches(4.3),
                      w - Inches(0.4), Inches(1.0),
                      fill=SOFT_BG, line=col,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        set_shape_text(vb, verdict, size=12, bold=True, color=col)
    return s


# ================================================================
#  SLIDE 9  — What is split inferencing?
# ================================================================
def slide_split_what():
    s = add_blank_slide()
    add_chrome(s,
               "What is split inferencing?",
               "Run the front half on the UE, ship features, run the back half on the edge server",
               section=3)

    dy = Inches(1.55)
    dh = Inches(3.0)

    # UE box
    ue = add_rect(s, Inches(0.55), dy, Inches(3.6), dh,
                  fill=SOFT_BG, line=NAVY, line_width=2)
    add_text_box(s, Inches(0.55), dy + Inches(0.05), Inches(3.6),
                 Inches(0.42),
                 "UE (vehicle / roadside pole)", font_size=13, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
    # Camera
    cam = add_rect(s, Inches(0.8), dy + Inches(0.65),
                   Inches(1.1), Inches(0.75),
                   fill=WHITE, line=DARK_GRAY)
    set_shape_text(cam, "Camera\n+ Radar", size=11)
    # Backbone
    bb = add_rect(s, Inches(2.35), dy + Inches(0.65),
                  Inches(1.55), Inches(1.1),
                  fill=ID_BLUE, line=NAVY)
    set_shape_text(bb, "Front half\n(Backbone)\nMobileNetV3\n/ ResNet",
                   size=10, color=WHITE)
    add_arrow(s, Inches(1.9), dy + Inches(1.05),
              Inches(2.35), dy + Inches(1.05), color=DARK_GRAY)
    # Split point
    sp = add_rect(s, Inches(2.35), dy + Inches(2.0),
                  Inches(1.55), Inches(0.5),
                  fill=ORANGE, line=DARK_GRAY)
    set_shape_text(sp, "✦  SPLIT POINT  ✦", size=10, color=WHITE)
    add_arrow(s, Inches(3.12), dy + Inches(1.75),
              Inches(3.12), dy + Inches(2.0), color=DARK_GRAY)
    # Feature annotation
    add_text_box(s, Inches(0.6), dy + Inches(2.6), Inches(3.5),
                 Inches(0.35),
                 "feature tensor  [C, H, W]  →  compact, task-rich",
                 font_size=10, italic=True, color=NAVY,
                 align=PP_ALIGN.CENTER)

    # Compression block
    comp = add_rect(s, Inches(4.5), dy + Inches(0.65),
                    Inches(1.5), Inches(1.85),
                    fill=WHITE, line=ORANGE, line_width=1.5)
    set_shape_text(comp, "Compress\nAE / ROI /\nQuantize /\nEntropy", size=11)
    add_arrow(s, Inches(4.15), dy + Inches(1.3),
              Inches(4.5), dy + Inches(1.3), color=DARK_GRAY)

    # Network cloud
    net = add_rect(s, Inches(6.35), dy + Inches(0.55),
                   Inches(1.5), Inches(2.0),
                   fill=SOFT_BG, line=NAVY, line_width=2,
                   shape=MSO_SHAPE.CLOUD)
    set_shape_text(net, "5G uplink\n(UDP)", size=11, bold=True, color=NAVY)
    add_arrow(s, Inches(6.0), dy + Inches(1.45),
              Inches(6.35), dy + Inches(1.45), color=NAVY)

    # Edge server
    srv = add_rect(s, Inches(8.15), dy, Inches(4.75), dh,
                   fill=SOFT_BG, line=NAVY, line_width=2)
    add_text_box(s, Inches(8.15), dy + Inches(0.05), Inches(4.75),
                 Inches(0.42),
                 "Edge AI server", font_size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    dq = add_rect(s, Inches(8.4), dy + Inches(0.65),
                  Inches(1.45), Inches(0.9),
                  fill=ORANGE, line=DARK_GRAY)
    set_shape_text(dq, "Decompress\nDe-quant", size=10, color=WHITE)
    hd = add_rect(s, Inches(10.1), dy + Inches(0.65),
                  Inches(1.45), Inches(0.9),
                  fill=NAVY, line=NAVY)
    set_shape_text(hd, "Back half\n(OD / SEG\nhead)", size=10, color=WHITE)
    out = add_rect(s, Inches(11.7), dy + Inches(0.65),
                   Inches(1.1), Inches(0.9),
                   fill=WHITE, line=DARK_GRAY)
    set_shape_text(out, "Boxes /\nMasks", size=10)
    add_arrow(s, Inches(7.85), dy + Inches(1.45),
              Inches(8.4), dy + Inches(1.1), color=NAVY)
    add_arrow(s, Inches(9.85), dy + Inches(1.1),
              Inches(10.1), dy + Inches(1.1), color=DARK_GRAY)
    add_arrow(s, Inches(11.55), dy + Inches(1.1),
              Inches(11.7), dy + Inches(1.1), color=DARK_GRAY)
    # Return
    add_text_box(s, Inches(8.4), dy + Inches(1.85), Inches(4.4),
                 Inches(0.42),
                 "result returned to UE for local control + display",
                 font_size=10, italic=True, color=MED_GRAY,
                 align=PP_ALIGN.CENTER)

    # Bottom bullets
    add_bullets(s, Inches(0.55), Inches(4.9), Inches(12.35), Inches(2.0),
                [
                    ("What is shipped.", "Not raw pixels. Not final boxes. A compact feature tensor — what the backbone 'saw' before the task-specific head."),
                    ("Why it works.", "Backbone features are high-information, low-redundancy. The head can run on the server without ever seeing the original image."),
                    ("What is configurable.", "Where to cut, how to compress (AE channels, ROI gating, quantization bits), how often to send."),
                ], size=14)
    return s


# ================================================================
#  SLIDE 10  — Why split inferencing?
# ================================================================
def slide_split_why():
    s = add_blank_slide()
    add_chrome(s,
               "Why split inferencing?",
               "Four reasons it beats raw video and bare detections for cooperative perception",
               section=3)

    reasons = [
        ("Bandwidth savings", "After quantization + entropy coding, feature tensors are 1–2 orders of magnitude smaller than raw video at the same task accuracy. OD baseline: ~87 KB/frame; SEG: ~410 KB/frame — vs ~900 KB/frame uncompressed at 720p.", ID_BLUE),
        ("Edge compute reuse", "The UE runs only the backbone (fast, low-power). The edge server runs the heavy head once for many UEs — fleet economics work because the expensive compute is amortized.", ORANGE),
        ("Privacy posture", "Raw pixels never leave the device. Only task-relevant feature activations cross the radio link. Regulatory and user-trust considerations are easier to meet.", GREEN),
        ("Fusion-ready representations", "Feature tensors carry more context than detections: spatial activations, confidence gradients, multi-scale cues. The server can fuse across UEs and reason about occlusion before committing to a detection.", PURPLE),
    ]
    w = Inches(6.2)
    h = Inches(2.5)
    gx = Inches(0.18)
    gy = Inches(0.2)
    x0 = Inches(0.45)
    y0 = Inches(1.55)
    for i, (head, body, col) in enumerate(reasons):
        r = i // 2
        c = i % 2
        left = x0 + (w + gx) * c
        top  = y0 + (h + gy) * r
        card = add_rect(s, left, top, w, h,
                        fill=WHITE, line=col, line_width=1.5)
        strip = add_rect(s, left, top, Inches(0.2), h,
                         fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  left + Inches(0.44), top + Inches(0.28),
                                  Inches(0.55), Inches(0.55))
        num.fill.solid()
        num.fill.fore_color.rgb = col
        num.line.fill.background()
        set_shape_text(num, str(i + 1), size=17, bold=True, color=WHITE)
        add_text_box(s, left + Inches(1.14), top + Inches(0.22),
                     w - Inches(1.25), Inches(0.55),
                     head, font_size=17, bold=True, color=NAVY)
        add_text_box(s, left + Inches(1.14), top + Inches(0.78),
                     w - Inches(1.25), h - Inches(0.88),
                     body, font_size=12, color=DARK_GRAY)
    add_text_box(s, Inches(0.45), Inches(6.85), Inches(12.45), Inches(0.4),
                 "Net effect: a transport-friendly, task-aligned intermediate representation that preserves task utility while fitting inside cellular budgets.",
                 font_size=12, italic=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    return s


# ================================================================
#  SLIDE 11  — Methodology / split architecture
# ================================================================
def slide_methodology():
    s = add_blank_slide()
    add_chrome(s,
               "Methodology — OD and SEG share the same UE pipeline",
               "Two task heads, same split-point and compression stack, very different payload sizes",
               section=3)

    dy = Inches(1.55)
    dh = Inches(4.3)

    # UE box
    ue = add_rect(s, Inches(0.5), dy, Inches(3.5), dh,
                  fill=SOFT_BG, line=NAVY, line_width=2)
    add_text_box(s, Inches(0.5), dy + Inches(0.05), Inches(3.5),
                 Inches(0.4), "UE side",
                 font_size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    # Sensors
    cam = add_rect(s, Inches(0.7), dy + Inches(0.6),
                   Inches(1.05), Inches(0.68), fill=WHITE, line=DARK_GRAY)
    set_shape_text(cam, "Camera", size=11)
    rad = add_rect(s, Inches(0.7), dy + Inches(1.4),
                   Inches(1.05), Inches(0.55), fill=WHITE, line=DARK_GRAY)
    set_shape_text(rad, "Radar\n(fusion opt.)", size=10)
    # Backbone
    bb = add_rect(s, Inches(2.1), dy + Inches(0.75),
                  Inches(1.65), Inches(1.1),
                  fill=ID_BLUE, line=NAVY)
    set_shape_text(bb, "Backbone\nMobileNetV3 /\nResNet+FPN",
                   size=10, color=WHITE)
    add_arrow(s, Inches(1.75), dy + Inches(0.95),
              Inches(2.1), dy + Inches(1.2), color=DARK_GRAY)
    add_arrow(s, Inches(1.75), dy + Inches(1.68),
              Inches(2.1), dy + Inches(1.45), color=DARK_GRAY)
    # Split point
    sp = add_rect(s, Inches(2.1), dy + Inches(2.1),
                  Inches(1.65), Inches(0.45),
                  fill=ORANGE, line=DARK_GRAY)
    set_shape_text(sp, "✦ SPLIT POINT ✦", size=10, color=WHITE)
    add_arrow(s, Inches(2.92), dy + Inches(1.85),
              Inches(2.92), dy + Inches(2.1), color=DARK_GRAY)
    # Compression
    comp = add_rect(s, Inches(2.1), dy + Inches(2.75),
                    Inches(1.65), Inches(1.4),
                    fill=WHITE, line=ORANGE, line_width=1.5)
    set_shape_text(comp, "AE / ROI\nQuantize\nEntropy code", size=10)
    add_arrow(s, Inches(2.92), dy + Inches(2.55),
              Inches(2.92), dy + Inches(2.75), color=DARK_GRAY)

    # Arrow to network
    add_arrow(s, Inches(3.75), dy + Inches(3.5),
              Inches(4.55), dy + Inches(2.15),
              color=NAVY, line_width=2)

    # Network cloud
    net = add_rect(s, Inches(4.55), dy + Inches(1.65),
                   Inches(1.6), Inches(1.3),
                   fill=SOFT_BG, line=NAVY, line_width=2,
                   shape=MSO_SHAPE.CLOUD)
    set_shape_text(net, "OAI 5G\nUDP", size=11, bold=True, color=NAVY)

    add_arrow(s, Inches(6.15), dy + Inches(2.3),
              Inches(6.85), dy + Inches(2.3),
              color=NAVY, line_width=2)

    # Server box
    srv = add_rect(s, Inches(6.85), dy, Inches(6.1), dh,
                   fill=SOFT_BG, line=NAVY, line_width=2)
    add_text_box(s, Inches(6.85), dy + Inches(0.05), Inches(6.1),
                 Inches(0.4), "Edge AI server",
                 font_size=13, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    dq = add_rect(s, Inches(7.1), dy + Inches(1.85),
                  Inches(1.45), Inches(0.95),
                  fill=ORANGE, line=DARK_GRAY)
    set_shape_text(dq, "Decompress\nDe-quant", size=10, color=WHITE)
    # OD head branch
    od = add_rect(s, Inches(8.85), dy + Inches(0.75),
                  Inches(1.75), Inches(1.05),
                  fill=NAVY, line=NAVY)
    set_shape_text(od, "OD head\nFaster R-CNN", size=11, color=WHITE)
    od_out = add_rect(s, Inches(10.85), dy + Inches(0.88),
                      Inches(1.9), Inches(0.8),
                      fill=WHITE, line=NAVY)
    set_shape_text(od_out, "Boxes + classes\n~87 KB/frame", size=10)
    # SEG head branch
    seg = add_rect(s, Inches(8.85), dy + Inches(2.85),
                   Inches(1.75), Inches(1.05),
                   fill=NAVY, line=NAVY)
    set_shape_text(seg, "SEG head\nLR-ASPP", size=11, color=WHITE)
    seg_out = add_rect(s, Inches(10.85), dy + Inches(2.95),
                       Inches(1.9), Inches(0.8),
                       fill=WHITE, line=NAVY)
    set_shape_text(seg_out, "Per-pixel mask\n~410 KB/frame", size=10)
    add_arrow(s, Inches(8.55), dy + Inches(2.3),
              Inches(8.85), dy + Inches(1.3), color=DARK_GRAY)
    add_arrow(s, Inches(8.55), dy + Inches(2.3),
              Inches(8.85), dy + Inches(3.35), color=DARK_GRAY)
    add_arrow(s, Inches(10.6), dy + Inches(1.28),
              Inches(10.85), dy + Inches(1.28), color=DARK_GRAY)
    add_arrow(s, Inches(10.6), dy + Inches(3.38),
              Inches(10.85), dy + Inches(3.38), color=DARK_GRAY)
    add_arrow(s, Inches(7.85), dy + Inches(2.3),
              Inches(8.55), dy + Inches(2.3), color=DARK_GRAY)

    # Bottom note
    add_bullets(s, Inches(0.5), Inches(6.1), Inches(12.4), Inches(1.1),
                [
                    ("Same front-half.", "Both tasks use the same backbone → split point → compress → UDP pipeline."),
                    ("Different payload.", "OD tensor is ~87 KB fixed shape; SEG tensor is ~4.6× larger. Same link, very different sizing problem."),
                ], size=13)
    return s


# ================================================================
#  SLIDE 12  — Compression knobs
# ================================================================
def slide_compression_knobs():
    s = add_blank_slide()
    add_chrome(s,
               "The compression knobs at the split point",
               "What the RL agent will control — today they are static, hand-picked settings",
               section=3)

    knobs = [
        ("AE channels  (128 / 64 / 32)",
         "Bottleneck auto-encoder replaces feature channels with a learned latent. Lower channels → fewer bytes, more reconstruction error. RD-AE 128 already cuts OD payload 57% vs baseline at similar mAP.",
         ID_BLUE),
        ("ROI threshold  (0.1 / 0.3 / 0.5)",
         "Saliency/objectness gate zeros out cells below the threshold before shipping. Higher threshold → smaller payload. ROI 0.4 reduces OD bytes 66% but may drop weak positives in crowded scenes.",
         ORANGE),
        ("Quantization  (8 / 6 / 4 bit)",
         "Per-channel uint8 → uint4. uint4 cuts SEG payload ~47% vs baseline while holding foreground IoU at 0.509 — but the floor matters: below 4-bit, accuracy falls sharply.",
         GREEN),
        ("Frame send / skip",
         "Send every Nth frame; UE uses a stale result on skipped frames. Saves bytes and RTT under congestion; increases object-state staleness, which is hazardous for fast-moving objects.",
         PURPLE),
        ("Redundancy  (FEC / dup)",
         "Add parity or duplicate chunks when measured packet loss is high. Costs extra bytes up front; protects feature integrity under bad channel and avoids costly full-frame retransmit.",
         TEAL),
    ]
    x = Inches(0.45)
    y = Inches(1.5)
    w = Inches(12.45)
    eh = Inches(0.98)
    gap = Inches(0.1)
    for i, (head, body, col) in enumerate(knobs):
        top = y + (eh + gap) * i
        add_rect(s, x, top, w, eh, fill=WHITE, line=col, line_width=1.5)
        cb = add_rect(s, x, top, Inches(2.85), eh,
                      fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        set_shape_text(cb, head, size=12, bold=True, color=WHITE)
        add_text_box(s, x + Inches(3.05), top + Inches(0.08),
                     w - Inches(3.2), eh - Inches(0.12),
                     body, font_size=12, color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
    return s


# ================================================================
#  SLIDE 13  — Baseline measurements
# ================================================================
def slide_measurements():
    s = add_blank_slide()
    add_chrome(s,
               "Baseline measurements — what we observed on our testbed",
               "Real numbers from OD + SEG split inference over OAI 5G (single-UE)",
               section=3)

    # Table on left
    rows = [
        ["Config",       "p50 payload", "p95 payload", "p50 RTT",  "Chunks"],
        ["OD baseline",  "~86.9 KB",    "~88.1 KB",    "~14 ms",   "2"],
        ["OD ROI 0.4",   "~30.4 KB",    "~55.0 KB",    "~13 ms",   "1"],
        ["OD RD-AE 128", "~36.6 KB",    "~38.2 KB",    "~13 ms",   "1"],
        ["SEG baseline", "~409.7 KB",   "~425.2 KB",   "~41 ms",   "7–8"],
        ["SEG uint4",    "~217.6 KB",   "~222.9 KB",   "~33 ms",   "4"],
        ["SEG ROI 0.1",  "~392.6 KB",   "~415.5 KB",   "~42 ms",   "7–8"],
    ]
    col_w = [Inches(2.4), Inches(1.4), Inches(1.4), Inches(1.0), Inches(1.2)]
    rh = Inches(0.55)
    tx = Inches(0.5)
    ty = Inches(1.55)
    for ri, row in enumerate(rows):
        cx = tx
        for ci, val in enumerate(row):
            cell = add_rect(s, cx, ty, col_w[ci], rh,
                            fill=(NAVY if ri == 0 else
                                  (SOFT_BG if ri % 2 else WHITE)),
                            line=LIGHT_GRAY, shape=MSO_SHAPE.RECTANGLE)
            set_shape_text(cell, val,
                           size=(12 if ri == 0 else 11),
                           bold=(ri == 0),
                           color=(WHITE if ri == 0 else DARK_GRAY))
            cx += col_w[ci]
        ty += rh

    # Right: takeaways
    add_text_box(s, Inches(8.2), Inches(1.55), Inches(4.8), Inches(0.5),
                 "What this tells us", font_size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(8.2), Inches(2.1), Inches(4.8), Inches(4.2),
                [
                    ("SEG is ~4.6× OD.", "Not the same problem. A link sized for OD is under-provisioned for SEG."),
                    ("Static knobs help.", "AE / ROI / uint4 cut 1.5–4× bytes — with different accuracy cost per knob."),
                    ("Knobs are NOT equivalent.", "Same byte budget reached by different routes gives different AP / mIoU."),
                    ("5QI burst is the hard wall.", "Even compressed SEG exceeds MDBV for 5QI 89/90 burst volume — traffic profile, not just latency, is the bottleneck."),
                ], size=13)

    # Bottom hook
    hookbox = add_rect(s, Inches(0.5), Inches(6.2), Inches(12.4),
                       Inches(0.9), fill=SOFT_BG, line=NAVY, line_width=1.5)
    add_text_box(s, Inches(0.7), Inches(6.28), Inches(12.0), Inches(0.75),
                 "These are STATIC choices. The agent's job is to pick the right knob combination for this frame, this scene, this link — without being told in advance which one to use.",
                 font_size=13, bold=True, italic=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


# ================================================================
#  SLIDE 14  — Bridge: single → multi-UE
# ================================================================
def slide_bridge():
    s = add_blank_slide()
    add_chrome(s,
               "From single-UE to multi-UE cooperative perception",
               "The same pipeline, scaled — and now the network is shared and the decisions must be coordinated",
               section=4)

    lx = Inches(0.5)
    ly = Inches(1.55)
    lw = Inches(6.0)
    lh = Inches(5.1)
    add_rect(s, lx, ly, lw, lh, fill=SOFT_BG,
             line=LIGHT_GRAY, shape=MSO_SHAPE.RECTANGLE)
    add_text_box(s, lx, ly + Inches(0.05), lw, Inches(0.4),
                 "Today (what we already built)",
                 font_size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    car = add_rect(s, lx + Inches(0.3), ly + Inches(2.0),
                   Inches(1.3), Inches(0.8),
                   fill=ID_BLUE, line=NAVY,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    set_shape_text(car, "UE", size=12, color=WHITE)
    gnb = add_rect(s, lx + Inches(2.2), ly + Inches(2.0),
                   Inches(1.3), Inches(0.8),
                   fill=WHITE, line=NAVY)
    set_shape_text(gnb, "gNB", size=12, color=NAVY)
    srv = add_rect(s, lx + Inches(4.1), ly + Inches(2.0),
                   Inches(1.6), Inches(0.8),
                   fill=NAVY, line=NAVY)
    set_shape_text(srv, "Edge head", size=12, color=WHITE)
    add_arrow(s, lx + Inches(1.6), ly + Inches(2.4),
              lx + Inches(2.2), ly + Inches(2.4), color=NAVY)
    add_arrow(s, lx + Inches(3.5), ly + Inches(2.4),
              lx + Inches(4.1), ly + Inches(2.4), color=NAVY)
    add_bullets(s, lx + Inches(0.3), ly + Inches(3.2),
                lw - Inches(0.6), Inches(1.8),
                [
                    "One uplink path, one task",
                    "Compression decisions fully local and isolated",
                    "Network state relatively stable",
                    "Baseline characterised (AI-traffic deck)",
                ], size=12)

    rx = Inches(6.85)
    ry = Inches(1.55)
    rw = Inches(6.05)
    rh = Inches(5.1)
    add_rect(s, rx, ry, rw, rh, fill=SOFT_BG,
             line=ID_BLUE, line_width=2, shape=MSO_SHAPE.RECTANGLE)
    add_text_box(s, rx, ry + Inches(0.05), rw, Inches(0.4),
                 "Where we are going",
                 font_size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    cc = [ID_BLUE, ORANGE, GREEN]
    for i in range(3):
        cy = ry + Inches(0.7 + i * 1.15)
        ci = add_rect(s, rx + Inches(0.2), cy, Inches(1.1),
                      Inches(0.65), fill=cc[i], line=DARK_GRAY,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        set_shape_text(ci, f"UE {i+1}", size=11, color=WHITE)
        add_arrow(s, rx + Inches(1.3), cy + Inches(0.32),
                  rx + Inches(2.4), ry + Inches(2.2),
                  color=cc[i], line_width=2)
    gnb2 = add_rect(s, rx + Inches(2.4), ry + Inches(1.8),
                    Inches(1.1), Inches(0.85),
                    fill=WHITE, line=NAVY, line_width=2)
    set_shape_text(gnb2, "gNB\n(shared)", size=10, color=NAVY)
    fuse = add_rect(s, rx + Inches(3.85), ry + Inches(1.45),
                    Inches(1.95), Inches(1.6),
                    fill=NAVY, line=NAVY)
    set_shape_text(fuse, "Edge fusion\n+\nShared\nspatial map", size=11,
                   color=WHITE)
    add_arrow(s, rx + Inches(3.5), ry + Inches(2.22),
              rx + Inches(3.85), ry + Inches(2.22), color=NAVY, line_width=2)
    add_bullets(s, rx + Inches(0.3), ry + Inches(3.4),
                rw - Inches(0.6), Inches(1.8),
                [
                    "Shared radio pool → UEs compete for bandwidth",
                    "Different scenes → each UE needs a different budget",
                    "Fusion needs confidence + provenance, not just boxes",
                    "Compression decisions must be coordinated and network-aware",
                ], size=12)

    # Divider arrow in the middle top
    add_text_box(s, Inches(5.8), Inches(2.6), Inches(1.8), Inches(0.5),
                 "scale +\ncoordinate", font_size=11, bold=True,
                 color=ID_BLUE, align=PP_ALIGN.CENTER)
    return s


# ================================================================
#  SLIDE 15  — SceneSense Agent core idea
# ================================================================
def slide_scenesense_idea():
    s = add_blank_slide()
    add_chrome(s,
               "SceneSense Agent — the core idea",
               "Learn what to share, when to share it, and how much to spend — subject to hard safety floors",
               section=4)

    # Left: motivating quote
    add_rich_tb(s, Inches(0.5), Inches(1.5), Inches(6.65), Inches(5.55),
                [
                    {"text": '"', "size": 60, "color": ID_BLUE, "bold": True},
                    {"text": "Can a machine learn, in real time, what visual information is worth sending over a busy wireless link so that safety-critical objects are still understood, bandwidth is not wasted, and nearby autonomous machines can be warned about hazards they cannot directly see?",
                     "size": 16, "italic": True, "color": NAVY, "space_after": 14},
                    {"text": "  "},
                    {"text": "— SceneSense Agent research proposal, May 2026",
                     "size": 11, "italic": True, "color": MED_GRAY},
                    {"text": "  ", "size": 14, "space_after": 10},
                    {"text": "This project's answer:", "size": 14, "bold": True,
                     "color": NAVY, "space_after": 6},
                    {"text": "Yes — through a constrained reinforcement-learning agent that uses scene content, link health, and model confidence to pick a compression policy, protected by deterministic guardrails that enforce safety-class recall and task precision.",
                     "size": 13, "color": DARK_GRAY},
                ])

    # Right: what we are doing — 5 cards
    add_text_box(s, Inches(7.4), Inches(1.5), Inches(5.5), Inches(0.5),
                 "What we are building", font_size=18, bold=True, color=NAVY)
    items = [
        ("Learn policies, not static knobs.",
         "RL agent picks AE / ROI / quant / scheduling / FEC from live scene + network + model state.", ID_BLUE),
        ("Guard safety classes.",
         "Hard cap on AP, mIoU, foreground IoU, and pedestrian / cyclist recall — cannot be traded away for bytes.", RED),
        ("Feed a shared spatial map.",
         "Accepted outputs populate a confidence-tagged map for occluded-object alerts to other vehicles.", ORANGE),
        ("Demonstrate the value.",
         "An AV uses the shared map to avoid a collision it could never detect with its own sensors alone.", GREEN),
    ]
    iy = Inches(2.1)
    for head, body, col in items:
        box = add_rect(s, Inches(7.4), iy, Inches(5.5), Inches(1.18),
                       fill=WHITE, line=col, line_width=1.5)
        strip = add_rect(s, Inches(7.4), iy, Inches(0.2), Inches(1.18),
                         fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        add_text_box(s, Inches(7.75), iy + Inches(0.08),
                     Inches(5.0), Inches(0.4),
                     head, font_size=13, bold=True, color=NAVY)
        add_text_box(s, Inches(7.75), iy + Inches(0.45),
                     Inches(5.0), Inches(0.65),
                     body, font_size=12, color=DARK_GRAY)
        iy += Inches(1.28)
    return s


# ================================================================
#  SLIDE 16  — Hypothesis and goal
# ================================================================
def slide_hypothesis():
    s = add_blank_slide()
    add_chrome(s,
               "Hypothesis and goal",
               "Concrete, testable claims the 6-month plan is built around",
               section=4)

    add_text_box(s, Inches(0.5), Inches(1.55), Inches(7.0), Inches(0.48),
                 "Hypotheses", font_size=19, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(2.05), Inches(7.0), Inches(5.0),
                [
                    ("H1 — Controller outperforms static knobs.",
                     "A constrained RL agent beats the best single static policy at the same task precision."),
                    ("H2 — Guardrails prevent safety failures.",
                     "No accepted action drops AP / mIoU / pedestrian / cyclist recall below configured floors."),
                    ("H3 — Outputs are useful physical-AI signals.",
                     "Confidence tags, foreground regions, and freshness metadata feed a trustworthy spatial map."),
                    ("H4 — Shared map reduces collision risk.",
                     "A learned map-sharing policy lets a vehicle avoid a hazard it cannot directly perceive."),
                ], size=14)

    add_text_box(s, Inches(7.8), Inches(1.55), Inches(5.2), Inches(0.48),
                 "First implementation goal", font_size=19, bold=True,
                 color=NAVY)
    goal = add_rect(s, Inches(7.8), Inches(2.1), Inches(5.2), Inches(2.9),
                    fill=SOFT_BG, line=ID_BLUE, line_width=2)
    set_shape_text(goal,
                   "Learned split-model control\nunder task-precision guardrails",
                   size=18, bold=True, color=NAVY)
    add_text_box(s, Inches(7.8), Inches(5.05), Inches(5.2), Inches(1.1),
                 "Task utility is the first-class objective — not bytes, not latency. Bytes and latency are costs, subject to guardrail floors, not the thing we optimise.",
                 font_size=12, italic=True, color=DARK_GRAY)
    add_text_box(s, Inches(7.8), Inches(6.0), Inches(5.2), Inches(0.8),
                 "Scope: months 1–3 prove the controller; months 4–6 close the spatial-map loop.",
                 font_size=12, italic=True, color=MED_GRAY)
    return s


# ================================================================
#  SLIDE 17  — Things we have to get right
# ================================================================
def slide_things_to_consider():
    s = add_blank_slide()
    add_chrome(s,
               "Things we have to get right",
               "Each axis is both a research question and an engineering challenge",
               section=4)

    items = [
        ("Scene awareness",
         "Represent 'what is happening' compactly enough to be an RL input: density, foreground fraction, object scale, occlusion ratio.",
         ID_BLUE),
        ("Network awareness",
         "Which 5G signals are available in real time? RTT, packet loss, queue delay, CQI, scheduling grants. OAI exposes some but not all cleanly.",
         ORANGE),
        ("Model awareness",
         "How confident is the model on this frame? Per-class confidence, uncertainty, foreground IoU proxy — without running a separate eval pipeline.",
         GREEN),
        ("Task-precision guardrails",
         "Reject any action that drops AP / mIoU below configured floor or increases pedestrian / small-object misses. Must be deterministic and cheap.",
         RED),
        ("Coordination scope",
         "Per-UE local agent first. Later: do sharing decisions stay local, or do we add a network-side orchestrator across UEs?",
         PURPLE),
        ("Policy action latency",
         "The agent itself must act inside the safety-critical window. A slow decision-maker cannot control a fast perception system.",
         TEAL),
    ]
    cw = Inches(4.18)
    ch = Inches(2.72)
    gx = Inches(0.13)
    gy = Inches(0.12)
    x0 = Inches(0.45)
    y0 = Inches(1.5)
    for i, (head, body, col) in enumerate(items):
        r = i // 3
        c = i % 3
        lx = x0 + (cw + gx) * c
        ty = y0 + (ch + gy) * r
        add_rect(s, lx, ty, cw, ch, fill=WHITE, line=col, line_width=1.5)
        band = add_rect(s, lx, ty, cw, Inches(0.58),
                        fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        set_shape_text(band, head, size=14, bold=True, color=WHITE)
        add_text_box(s, lx + Inches(0.15), ty + Inches(0.7),
                     cw - Inches(0.3), ch - Inches(0.8),
                     body, font_size=12, color=DARK_GRAY)
    return s


# ================================================================
#  SLIDE 18  — Big picture architecture
# ================================================================
def slide_big_picture():
    s = add_blank_slide()
    add_chrome(s,
               "SceneSense — the big picture",
               "UE RL agents + OAI 5G + Edge AI server + Shared Spatial Map + Map-sharing agent + guardrails",
               section=5)

    ue_cols = [ID_BLUE, ORANGE, GREEN]
    ue_lbl  = ["Car 1", "Car 2", "Car 3"]

    for i in range(3):
        y0 = Inches(1.45 + i * 1.55)
        car = add_rect(s, Inches(0.4), y0 + Inches(0.18),
                       Inches(0.9), Inches(0.55),
                       fill=ue_cols[i], line=DARK_GRAY,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        set_shape_text(car, ue_lbl[i], size=10, color=WHITE)
        ag = add_rect(s, Inches(1.4), y0, Inches(2.65), Inches(1.2),
                      fill=WHITE, line=ue_cols[i], line_width=1.5)
        add_text_box(s, Inches(1.4), y0 + Inches(0.02),
                     Inches(2.65), Inches(0.3),
                     "Split-Control RL agent", font_size=11, bold=True,
                     color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(1.5), y0 + Inches(0.33),
                     Inches(2.55), Inches(0.8),
                     "state: scene + model + network\naction: AE/ROI/quant/sched/FEC",
                     font_size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        add_arrow(s, Inches(4.05), y0 + Inches(0.6),
                  Inches(5.0), Inches(3.75),
                  color=ue_cols[i], line_width=1.75)

    # gNB
    gnb = add_rect(s, Inches(5.0), Inches(3.38), Inches(1.5), Inches(0.9),
                   fill=WHITE, line=NAVY, line_width=2)
    set_shape_text(gnb, "gNB\nOAI 5G", size=11, bold=True, color=NAVY)
    upf = add_rect(s, Inches(5.0), Inches(4.4), Inches(1.5), Inches(0.65),
                   fill=WHITE, line=NAVY, line_width=2)
    set_shape_text(upf, "5G Core / UPF", size=10, color=NAVY)

    # Edge AI server
    srv = add_rect(s, Inches(7.0), Inches(2.9), Inches(2.8), Inches(2.5),
                   fill=SOFT_BG, line=NAVY, line_width=2)
    add_text_box(s, Inches(7.0), Inches(2.93), Inches(2.8), Inches(0.38),
                 "Edge AI server", font_size=12, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)
    od = add_rect(s, Inches(7.2), Inches(3.38), Inches(2.4), Inches(0.5),
                  fill=NAVY, line=NAVY)
    set_shape_text(od, "OD head → boxes / pose", size=10, color=WHITE)
    sg = add_rect(s, Inches(7.2), Inches(3.98), Inches(2.4), Inches(0.5),
                  fill=NAVY, line=NAVY)
    set_shape_text(sg, "SEG head → mask", size=10, color=WHITE)
    cf = add_rect(s, Inches(7.2), Inches(4.58), Inches(2.4), Inches(0.65),
                  fill=PURPLE, line=PURPLE)
    set_shape_text(cf, "Confidence /\nuncertainty estimator", size=10,
                   color=WHITE)

    add_arrow(s, Inches(6.5), Inches(3.83),
              Inches(7.0), Inches(3.83), color=NAVY, line_width=2)

    # Shared spatial map
    sm = add_rect(s, Inches(10.1), Inches(1.5), Inches(2.9), Inches(2.2),
                  fill=SOFT_BG, line=ORANGE, line_width=2)
    add_text_box(s, Inches(10.1), Inches(1.53), Inches(2.9), Inches(0.42),
                 "Shared spatial map", font_size=13, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(10.2), Inches(1.97), Inches(2.7), Inches(1.7),
                 "class · pose · velocity\nconfidence · provenance\nfreshness · occlusion state",
                 font_size=10, italic=True, color=DARK_GRAY,
                 align=PP_ALIGN.CENTER)

    # Map-sharing agent
    msa = add_rect(s, Inches(10.1), Inches(3.85), Inches(2.9), Inches(1.4),
                   fill=WHITE, line=ORANGE, line_width=2)
    add_text_box(s, Inches(10.1), Inches(3.88), Inches(2.9), Inches(0.35),
                 "Map-sharing RL agent", font_size=12, bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(10.15), Inches(4.25), Inches(2.75), Inches(1.0),
                 "action: what/when/who/detail\nreward: utility−bytes−latency−stale",
                 font_size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    add_arrow(s, Inches(9.8), Inches(3.35),
              Inches(10.1), Inches(2.4), color=ORANGE, line_width=2)

    # Guardrails ribbon
    grd = add_rect(s, Inches(0.4), Inches(6.1), Inches(12.55), Inches(0.88),
                   fill=RGBColor(0xFD, 0xEC, 0xEA),
                   line=RED, line_width=1.5, shape=MSO_SHAPE.RECTANGLE)
    add_text_box(s, Inches(0.6), Inches(6.14), Inches(12.2), Inches(0.35),
                 "Task-precision guardrails  ·  Must hold on every accepted action",
                 font_size=12, bold=True, color=RED)
    add_text_box(s, Inches(0.6), Inches(6.47), Inches(12.2), Inches(0.45),
                 "AP drop ≤ ε_OD  ·  mIoU drop ≤ ε_SEG  ·  Recall_ped/cyclist ≥ τ  ·  latency ≤ L_max  →  Accept / Clamp / Reject",
                 font_size=11, italic=True, color=DARK_GRAY)
    return s


# ================================================================
#  SLIDE 19  — Breakdown: UE-side RL agent
# ================================================================
def slide_breakdown_ue():
    s = add_blank_slide()
    add_chrome(s,
               "Breakdown · UE-side Split-Control RL Agent",
               "Per-UE policy that picks compression actions from current observed state",
               section=5)

    add_text_box(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.48),
                 "State input  s_t", font_size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(2.08), Inches(6.0), Inches(4.5),
                [
                    ("Scene  x_scene", "density, foreground fraction, occlusion ratio, scale distribution of visible objects."),
                    ("Model  x_model", "head confidence, per-class uncertainty, predicted task risk score."),
                    ("Network  x_net", "RTT, packet loss, throughput estimate, queue delay, scheduling grant sizes."),
                    ("Context  q_t", "freshness, previous action, time since last successfully delivered frame."),
                ], size=14)

    add_text_box(s, Inches(7.0), Inches(1.55), Inches(5.9), Inches(0.48),
                 "Action output  a_t", font_size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(7.0), Inches(2.08), Inches(5.9), Inches(2.8),
                [
                    "AE channels  (128 / 64 / 32)",
                    "ROI threshold  (0.1 / 0.3 / 0.5)",
                    "Quantization  (8 / 6 / 4 bit)",
                    "Frame send / skip",
                    "Redundancy add / drop  (FEC, duplication)",
                ], size=14)

    rb = add_rect(s, Inches(0.5), Inches(5.6), Inches(12.45), Inches(1.52),
                  fill=SOFT_BG, line=ID_BLUE, line_width=2)
    add_text_box(s, Inches(0.7), Inches(5.65), Inches(12.0), Inches(0.4),
                 "Reward  R_t  and constraints", font_size=14, bold=True,
                 color=NAVY)
    add_text_box(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.6),
                 "R_t  =  task_utility (AP, mIoU, foreground IoU, recall)  −  λ_b · bytes  −  λ_l · latency  −  λ_p · loss\nsubject to guardrails: AP_drop ≤ ε_OD, mIoU_drop ≤ ε_SEG, Recall_ped ≥ τ, latency ≤ L_max",
                 font_size=13, color=DARK_GRAY)
    return s


# ================================================================
#  SLIDE 20  — Breakdown: Edge AI server + Spatial Map
# ================================================================
def slide_breakdown_server():
    s = add_blank_slide()
    add_chrome(s,
               "Breakdown · Edge AI server and Shared Spatial Map",
               "Where feature tensors become tagged object evidence and are persisted in the map",
               section=5)

    # Left: server pipeline
    add_text_box(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.48),
                 "Edge AI server", font_size=18, bold=True, color=NAVY)
    stages = [
        ("De-quant / decompress",
         "Reverse the per-UE compression using the agreed settings.", WHITE, DARK_GRAY),
        ("OD head",
         "Faster R-CNN → bounding boxes, classes, confidence scores.", ID_BLUE, WHITE),
        ("SEG head",
         "LR-ASPP → per-pixel semantic mask.", ID_BLUE, WHITE),
        ("Confidence / uncertainty head",
         "Calibrated per-detection and per-region confidence.", PURPLE, WHITE),
        ("Per-UE assembly",
         "Tag outputs: provenance (UE-id), timestamp, pose, sensor type.", ORANGE, WHITE),
    ]
    ty = Inches(2.08)
    for head, body, fill, fg in stages:
        bx = add_rect(s, Inches(0.5), ty, Inches(6.0), Inches(0.85),
                      fill=fill, line=NAVY, line_width=1.0)
        add_text_box(s, Inches(0.7), ty + Inches(0.05),
                     Inches(5.7), Inches(0.4),
                     head, font_size=13, bold=True,
                     color=fg if fill != WHITE else NAVY)
        add_text_box(s, Inches(0.7), ty + Inches(0.43),
                     Inches(5.7), Inches(0.38),
                     body, font_size=11,
                     color=fg if fill != WHITE else DARK_GRAY)
        ty += Inches(0.95)

    # Right: shared spatial map fields
    add_text_box(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(0.48),
                 "Shared spatial map  m_t", font_size=18, bold=True, color=NAVY)
    mp = add_rect(s, Inches(7.0), Inches(2.08), Inches(6.0), Inches(4.9),
                  fill=SOFT_BG, line=ORANGE, line_width=2)
    fields = [
        ("class",         "vehicle, pedestrian, cyclist, unknown"),
        ("pose",          "(x, y, yaw) in shared world frame"),
        ("velocity",      "estimated from tracker or upstream"),
        ("confidence",    "per-object posterior, calibrated"),
        ("provenance",    "which UE / sensor, and its known position"),
        ("freshness",     "time since last update from any UE"),
        ("occlusion",     "visible / partial / full, per-object"),
    ]
    fy = Inches(2.22)
    for k, v in fields:
        chip = add_rect(s, Inches(7.2), fy, Inches(1.85), Inches(0.52),
                        fill=ORANGE, line=ORANGE,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        set_shape_text(chip, k, size=11, color=WHITE)
        add_text_box(s, Inches(9.15), fy + Inches(0.06),
                     Inches(3.7), Inches(0.4),
                     v, font_size=11, color=DARK_GRAY,
                     anchor=MSO_ANCHOR.MIDDLE)
        fy += Inches(0.62)
    add_text_box(s, Inches(7.2), Inches(6.55), Inches(5.6), Inches(0.38),
                 "The map is the interface contract between the perception pipeline and all downstream consumers (vehicles, planners, monitors).",
                 font_size=10, italic=True, color=NAVY)
    return s


# ================================================================
#  SLIDE 21  — Breakdown: Map-sharing agent + guardrails
# ================================================================
def slide_breakdown_mapsharing():
    s = add_blank_slide()
    add_chrome(s,
               "Breakdown · Map-sharing RL agent and task-precision guardrails",
               "Decide WHAT update to push to which vehicle, WHEN, and at what DETAIL — with safety checks on every action",
               section=5)

    # Left: map-sharing agent
    add_text_box(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.48),
                 "Map-sharing RL agent", font_size=18, bold=True, color=NAVY)

    add_text_box(s, Inches(0.5), Inches(2.08), Inches(6.0), Inches(0.38),
                 "State  z_t", font_size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(2.48), Inches(6.0), Inches(1.55),
                ["map risk (collision-relevant objects)",
                 "object freshness scores",
                 "vehicle trajectory and heading",
                 "network load on the downlink"], size=13)

    add_text_box(s, Inches(0.5), Inches(4.1), Inches(6.0), Inches(0.38),
                 "Action  u_t", font_size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(4.5), Inches(6.0), Inches(1.5),
                ["what to share (which objects, which fields)",
                 "when to share (event-triggered vs scheduled)",
                 "who receives (specific UEs, broadcast, RSU)",
                 "detail level (compact summary vs full record)"], size=13)

    rb = add_rect(s, Inches(0.5), Inches(6.05), Inches(6.0), Inches(0.9),
                  fill=SOFT_BG, line=ORANGE, line_width=1.5)
    add_text_box(s, Inches(0.7), Inches(6.1), Inches(5.6), Inches(0.8),
                 "R_t  =  task_utility  −  λ_b · bytes  −  λ_l · latency  −  λ_s · stale_risk",
                 font_size=13, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

    # Right: guardrails pipeline
    add_text_box(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(0.48),
                 "Task-precision guardrails", font_size=18, bold=True,
                 color=RED)
    add_text_box(s, Inches(7.0), Inches(2.08), Inches(6.0), Inches(0.42),
                 "The RL policy is a PROPOSER. The guardrail layer is the GATEKEEPER.",
                 font_size=12, italic=True, color=DARK_GRAY)
    pipe = [
        ("Proposed action from RL agent", ID_BLUE),
        ("Guardrail checks\n(deterministic, cheap, always run)", RED),
        ("Accept  /  Clamp  /  Reject", GREEN),
    ]
    py = Inches(2.6)
    for i, (lbl, col) in enumerate(pipe):
        bx = add_rect(s, Inches(7.0), py, Inches(6.0), Inches(0.85),
                      fill=col, line=col)
        set_shape_text(bx, lbl, size=14, bold=True, color=WHITE)
        if i < len(pipe) - 1:
            add_arrow(s, Inches(10.0), py + Inches(0.85),
                      Inches(10.0), py + Inches(1.05),
                      color=DARK_GRAY)
        py += Inches(1.05)

    add_text_box(s, Inches(7.0), Inches(5.5), Inches(6.0), Inches(0.38),
                 "Constraints (must hold before Accept)", font_size=13,
                 bold=True, color=RED)
    add_bullets(s, Inches(7.0), Inches(5.9), Inches(6.0), Inches(1.5),
                ["AP drop ≤ ε_OD",
                 "mIoU drop ≤ ε_SEG",
                 "Recall_pedestrian / cyclist ≥ τ",
                 "End-to-end latency ≤ L_max"], size=13)
    return s


# ================================================================
#  SLIDE 22  — Evaluation campaigns
# ================================================================
def slide_evaluation():
    s = add_blank_slide()
    add_chrome(s,
               "Evaluation — four CARLA experiment campaigns",
               "Two to validate the controller, two to close the spatial-map / navigation-override loop",
               section=6)

    cps = [
        ("A — Static knobs vs learned controller",
         "Single CARLA ego vehicle, OD + SEG split routes, repeatable network stress profiles.",
         "Task utility retained at lower byte / latency / loss cost than the best static knob policy.",
         ID_BLUE),
        ("B — Guardrail stress test",
         "Crowded, sparse, high-jitter, packet-error, and queueing stress scenarios.",
         "No accepted policy violates AP / mIoU / pedestrian / cyclist recall thresholds.",
         ORANGE),
        ("C — Physical-AI map update",
         "CARLA intersection with occluded objects; split-model outputs feed the shared spatial map.",
         "Map freshness and localization sufficient for AV risk assessment.",
         GREEN),
        ("D — Navigation override demo",
         "AV approaches an occluded course-conflict object revealed through learned map sharing.",
         "Vehicle triggers safe override; collision avoided across the full scenario battery.",
         RED),
    ]
    cw = Inches(6.25)
    ch = Inches(2.65)
    gx = Inches(0.18)
    gy = Inches(0.18)
    x0 = Inches(0.45)
    y0 = Inches(1.5)
    for i, (head, setup, metric, col) in enumerate(cps):
        r = i // 2
        c = i % 2
        lx = x0 + (cw + gx) * c
        ty = y0 + (ch + gy) * r
        add_rect(s, lx, ty, cw, ch, fill=WHITE, line=col, line_width=1.5)
        band = add_rect(s, lx, ty, cw, Inches(0.55),
                        fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        set_shape_text(band, head, size=13, bold=True, color=WHITE)
        add_text_box(s, lx + Inches(0.18), ty + Inches(0.65),
                     cw - Inches(0.36), Inches(0.38),
                     "Setup", font_size=11, bold=True, color=NAVY)
        add_text_box(s, lx + Inches(0.18), ty + Inches(1.0),
                     cw - Inches(0.36), Inches(0.7),
                     setup, font_size=12, color=DARK_GRAY)
        add_text_box(s, lx + Inches(0.18), ty + Inches(1.68),
                     cw - Inches(0.36), Inches(0.38),
                     "Headline metric", font_size=11, bold=True, color=NAVY)
        add_text_box(s, lx + Inches(0.18), ty + Inches(2.05),
                     cw - Inches(0.36), Inches(0.55),
                     metric, font_size=12, color=DARK_GRAY)
    return s


# ================================================================
#  SLIDE 23  — Timeline
# ================================================================
def slide_timeline():
    s = add_blank_slide()
    add_chrome(s,
               "Six-month plan  ·  May 2026 – November 2026",
               "Months 1–3 prove controlled split inference; Months 4–6 close the spatial-map loop",
               section=6)

    months = [
        ("Month 1",
         "Reproduce split-inference baselines in CARLA.\nFreeze RL state / action / reward schema.",
         "Repeatable OD/SEG traces: bytes, latency, loss, AP/mIoU, foreground IoU.",
         ID_BLUE),
        ("Month 2",
         "Implement constrained RL controller over AE / ROI / quant / scheduling / redundancy.",
         "Policy trains/evaluates against static policies on logged metrics.",
         ID_BLUE),
        ("Month 3",
         "Evaluate task-precision guardrails under network stress.",
         "Plots: learned control beats static knobs without guardrail violations.",
         ID_BLUE),
        ("Month 4",
         "Spatial-map ingestion from split-model outputs.\nCARLA ground truth for validation.",
         "Map stores class, pose, velocity, confidence, provenance, freshness, occlusion.",
         ORANGE),
        ("Month 5",
         "Train/evaluate learned map-sharing policies for occluded-object updates.",
         "Policy ranks what/when/how to share under bandwidth and freshness constraints.",
         ORANGE),
        ("Month 6",
         "Intersection collision-avoidance demo.\nPaper + figures + disclosure prep.",
         "End-to-end CARLA demo, paper outline, invention-disclosure candidate.",
         ORANGE),
    ]
    n = len(months)
    track_x = Inches(0.5)
    track_y = Inches(1.62)
    seg_w = Inches(12.35 / n)

    # Track line
    tl = add_rect(s, track_x, track_y + Inches(0.22),
                  Inches(12.35), Inches(0.08),
                  fill=LIGHT_GRAY, line=LIGHT_GRAY,
                  shape=MSO_SHAPE.RECTANGLE)
    # Phase label blocks
    ph1 = add_rect(s, track_x, track_y - Inches(0.4),
                   Inches(12.35 / 2), Inches(0.35),
                   fill=ID_BLUE, line=None, shape=MSO_SHAPE.RECTANGLE)
    set_shape_text(ph1, "Phase 1: Learned split-model control  (Months 1–3)",
                   size=11, bold=True, color=WHITE)
    ph2 = add_rect(s, track_x + Inches(12.35 / 2), track_y - Inches(0.4),
                   Inches(12.35 / 2), Inches(0.35),
                   fill=ORANGE, line=None, shape=MSO_SHAPE.RECTANGLE)
    set_shape_text(ph2, "Phase 2: Spatial-map + navigation override  (Months 4–6)",
                   size=11, bold=True, color=WHITE)

    for i, (label, work, exit_, col) in enumerate(months):
        cx = track_x + seg_w * i + seg_w / 2 - Inches(0.25)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, track_y + Inches(0.05),
                                  Inches(0.5), Inches(0.5))
        dot.fill.solid()
        dot.fill.fore_color.rgb = col
        dot.line.color.rgb = WHITE
        dot.line.width = Pt(2)
        set_shape_text(dot, str(i + 1), size=14, bold=True, color=WHITE)

        card_x = track_x + seg_w * i + Inches(0.08)
        card_w = seg_w - Inches(0.16)
        card_y = Inches(2.42)
        card_h = Inches(4.65)
        add_rect(s, card_x, card_y, card_w, card_h,
                 fill=WHITE, line=col, line_width=1.5)
        band = add_rect(s, card_x, card_y, card_w, Inches(0.48),
                        fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        set_shape_text(band, label, size=12, bold=True, color=WHITE)
        add_text_box(s, card_x + Inches(0.08), card_y + Inches(0.55),
                     card_w - Inches(0.16), Inches(0.3),
                     "Work", font_size=9, bold=True, color=NAVY)
        add_text_box(s, card_x + Inches(0.08), card_y + Inches(0.85),
                     card_w - Inches(0.16), Inches(1.85),
                     work, font_size=9, color=DARK_GRAY)
        add_text_box(s, card_x + Inches(0.08), card_y + Inches(2.65),
                     card_w - Inches(0.16), Inches(0.3),
                     "Exit criterion", font_size=9, bold=True, color=NAVY)
        add_text_box(s, card_x + Inches(0.08), card_y + Inches(2.95),
                     card_w - Inches(0.16), Inches(1.65),
                     exit_, font_size=9, color=DARK_GRAY)
    return s


# ================================================================
#  SLIDE 24  — Risks, open questions, next steps
# ================================================================
def slide_risks():
    s = add_blank_slide()
    add_chrome(s,
               "Risks, open questions, and what I am doing next",
               "Where the project could fail — and how we will know early",
               section=6)

    cols = [
        ("Risks", [
            ("Controller minimises bytes only.",
             "Hard guardrails on AP, mIoU, foreground IoU, pedestrian / cyclist recall prevent this."),
            ("Action space too broad to train.",
             "Start AE / ROI / quant only; add scheduling + FEC after stable baseline."),
            ("Aggregate mIoU hides safety failures.",
             "Track foreground IoU and per-class recall (pedestrian, cyclist, small-object) separately."),
            ("Map scope expands too far.",
             "Months 4–6 strictly: ingestion + sharing + one override scenario battery."),
            ("OAI / PC5 integration consumes intern time.",
             "Trace-driven stress first; OAI live only after controller proof point."),
        ], RED),
        ("Open questions", [
            ("Coordination scope.",
             "Per-UE local agent first; decide centralized vs distributed later based on data."),
            ("Best importance metric.",
             "L2 saliency vs gradient-based vs learned signal — comparison is month 2 work."),
            ("Policy action latency.",
             "Can the agent itself stay within the safety-critical decision window?"),
            ("Multi-UE scaling.",
             "Where does radio resource / fusion server saturate as N grows?"),
        ], ORANGE),
        ("Next 4 weeks", [
            ("Lock the logging schema.",
             "Application + network CSVs aligned by run_group (in flight)."),
            ("Run static-knob baseline curves.",
             "Payload vs AP / mIoU at varying ROI / AE / quantization over OAI multi-UE."),
            ("Draft RL state / action / reward schema.",
             "Concrete feature list, action discretization, reward shape — as a document for review."),
            ("CARLA occlusion scenario harness.",
             "Repeatable curbside-parked-vehicle pedestrian scenario for campaigns B and D."),
        ], GREEN),
    ]
    cw = Inches(4.18)
    h = Inches(5.6)
    gx = Inches(0.13)
    x0 = Inches(0.45)
    y0 = Inches(1.5)
    for i, (head, items, col) in enumerate(cols):
        lx = x0 + (cw + gx) * i
        add_rect(s, lx, y0, cw, h, fill=WHITE, line=col, line_width=1.5)
        band = add_rect(s, lx, y0, cw, Inches(0.55),
                        fill=col, line=col, shape=MSO_SHAPE.RECTANGLE)
        set_shape_text(band, head, size=16, bold=True, color=WHITE)
        iy = y0 + Inches(0.68)
        for hd, body in items:
            add_text_box(s, lx + Inches(0.18), iy,
                         cw - Inches(0.36), Inches(0.38),
                         hd, font_size=12, bold=True, color=NAVY)
            add_text_box(s, lx + Inches(0.18), iy + Inches(0.36),
                         cw - Inches(0.36), Inches(0.85),
                         body, font_size=11, color=DARK_GRAY)
            iy += Inches(1.06)
    return s


# ================================================================
#  Main
# ================================================================
def main():
    slide_title()               # 1
    slide_outline()             # 2
    slide_introduction()        # 3
    slide_motivation()          # 4
    slide_problem_occlusion()   # 5
    slide_coop_concept()        # 6
    slide_use_cases()           # 7
    slide_sharing_options()     # 8
    slide_split_what()          # 9
    slide_split_why()           # 10
    slide_methodology()         # 11
    slide_compression_knobs()   # 12
    slide_measurements()        # 13
    slide_bridge()              # 14
    slide_scenesense_idea()     # 15
    slide_hypothesis()          # 16
    slide_things_to_consider()  # 17
    slide_big_picture()         # 18
    slide_breakdown_ue()        # 19
    slide_breakdown_server()    # 20
    slide_breakdown_mapsharing()# 21
    slide_evaluation()          # 22
    slide_timeline()            # 23
    slide_risks()               # 24

    out = "SceneSense_Agent_Intro_Deck.pptx"
    prs.save(out)
    print(f"Wrote {out} with {len(prs.slides)} slides.")


if __name__ == "__main__":
    main()
