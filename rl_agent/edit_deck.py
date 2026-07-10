#!/usr/bin/env python3
"""Update 'Standup Meeting.pptx' -> '..._updated.pptx':
 (1) slide 4: replace RGB-only LR-ASPP image with the RGB+radar fusion diagram (remove redundant side text)
 (2) new slide after AE: Quantization -> payload  (3) new slide: Radar pps -> accuracy, not payload
 (4) Action Space slide: add the 3 knob plots (entropy/quant/ROI)."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE as ST

R = "rl_agent/plots/_render"
SRC = "rl_agent/plots/Standup Meeting.pptx"
DST = "rl_agent/plots/Standup Meeting_updated.pptx"
prs = Presentation(SRC)
SW = prs.slide_width


def dup_slide(prs, index):
    src = prs.slides[index]
    new = prs.slides.add_slide(src.slide_layout)
    for ph in list(new.shapes):                      # drop layout-added placeholders
        ph._element.getparent().remove(ph._element)
    for shp in src.shapes:                           # copy every source shape
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def set_title(slide, text, needle):
    for shp in slide.shapes:
        if shp.has_text_frame and needle.lower() in shp.text_frame.text.lower():
            p = shp.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = text
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = text
            return True
    return False


def remove_pictures(slide):
    for shp in list(slide.shapes):
        if shp.shape_type == ST.PICTURE:
            shp._element.getparent().remove(shp._element)


# ---------- (1) SLIDE 4: swap in the fusion diagram ----------
s4 = prs.slides[3]
remove_pictures(s4)
for shp in list(s4.shapes):                          # remove the redundant right-side "low-stage" text box
    if shp.has_text_frame and "low-stage feature map" in shp.text_frame.text.lower():
        shp._element.getparent().remove(shp._element)
s4.shapes.add_picture(f"{R}/c_fusion.png", Inches(0.33), Inches(1.55), width=Inches(12.66))

# ---------- (2)+(3) two new slides duplicated from the AE slide (index 5) ----------
q = dup_slide(prs, 5); remove_pictures(q)
set_title(q, "Quantization to compress features", "compress features")
q.shapes.add_picture(f"{R}/c_quant.png", Inches(1.1), Inches(2.2), width=Inches(11.1))

rd = dup_slide(prs, 5); remove_pictures(rd)
set_title(rd, "Radar density enriches input & accuracy — not payload", "compress features")
rd.shapes.add_picture(f"{R}/c_radarpps.png", Inches(1.9), Inches(1.5), width=Inches(9.5))

# reorder: place q at index 6 (slide 7), rd at index 7 (slide 8) -- right after the AE slide
lst = prs.slides._sldIdLst
ids = list(lst)
q_el, rd_el = ids[-2], ids[-1]
lst.remove(q_el); lst.remove(rd_el)
lst.insert(6, q_el); lst.insert(7, rd_el)

# ---------- (4) ACTION SPACE slide: add the 3 knob plots ----------
act = None
for s in prs.slides:
    if any(sh.has_text_frame and sh.text_frame.text.strip().startswith("Action Space") for sh in s.shapes):
        act = s; break
# shrink the big (mostly empty) description box to a top strip
for shp in act.shapes:
    if shp.has_text_frame and shp.text_frame.text.strip().startswith("We characterize"):
        shp.height = Inches(1.1)
plots = ["knob_entropy-1.png", "knob_quant-1.png", "knob_roi-1.png"]
lefts = [0.35, 4.55, 8.75]
for fn, lf in zip(plots, lefts):
    act.shapes.add_picture(f"rl_agent/plots/{fn}", Inches(lf), Inches(2.5), width=Inches(4.05))

prs.save(DST)
print("saved", DST, "slides:", len(prs.slides))
for i, s in enumerate(prs.slides, 1):
    t = next((sh.text_frame.text.strip().split("\n")[0][:42] for sh in s.shapes
              if sh.has_text_frame and sh.text_frame.text.strip()
              and "InterDigital" not in sh.text_frame.text), "")
    print(f"  [{i}] {t}")
