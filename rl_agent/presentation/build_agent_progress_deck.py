#!/usr/bin/env python3
"""Build the 2026-08-19 SceneSense agent advisor-meeting slide deck.

The deck is intentionally rendered as high-resolution slide artwork and then
embedded one image per PowerPoint slide.  This keeps mathematical notation,
diagrams, and layout identical across PowerPoint and LibreOffice while leaving
all generation logic and source paths reproducible in this script.
"""

from __future__ import annotations

import hashlib
import json
import math
import textwrap
from pathlib import Path
from typing import Iterable, Sequence

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Circle, FancyArrowPatch, FancyBboxPatch, Polygon
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "rl_agent" / "presentation" / "agent_progress_20260819"
SLIDES = OUT / "slides"
PPTX = OUT / "SceneSense_Agent_Progress_2026-08-19.pptx"
NOTES = OUT / "PRESENTER_NOTES.md"
MANIFEST = OUT / "SOURCE_MANIFEST.json"
CONTACT_SHEET = OUT / "CONTACT_SHEET.png"

W, H = 16.0, 9.0
DPI = 160

BG = "#07111F"
PANEL = "#0E2034"
PANEL_2 = "#132A42"
INK = "#F7FAFC"
MUTED = "#AFC2D7"
CYAN = "#40D4E6"
TEAL = "#2ED6A1"
AMBER = "#FFB84D"
CORAL = "#FF6B7A"
BLUE = "#6FA8FF"
PURPLE = "#A98BFF"
GRID = "#28415C"
WHITE = "#FFFFFF"
BLACK = "#07111F"

FONT = "Lato"
MONO = "DejaVu Sans Mono"


ASSETS = {
    "worked": Path(
        "/tmp/phase2_geometry_review_midblock_van_positive_20260818_001417/"
        "20260818_001435_896264_combined.png"
    ),
    "curbside": Path(
        "/tmp/phase2_geometry_review_positive_20260815_041815/"
        "20260815_041825_634287_combined.png"
    ),
    "signalized": Path(
        "/tmp/phase2_geometry_review_signalized_corner_positive_20260817_224735/"
        "20260817_224753_327439_combined.png"
    ),
    "cross_traffic": Path(
        "/tmp/phase2_geometry_review_cross_traffic_vehicle_positive_20260818_024451/"
        "20260818_024513_531786_combined.png"
    ),
    "pullout": Path(
        "/tmp/phase2_geometry_review_parked_vehicle_pullout_positive_20260818_031013/"
        "20260818_031036_689670_combined.png"
    ),
    "queue": Path(
        "/tmp/phase2_geometry_review_queue_reveal_vehicle_positive_20260818_033946/"
        "20260818_034002_180110_combined.png"
    ),
    "natural_signal": Path(
        "/tmp/phase2_geometry_review_naturalistic_pair_naturalistic_20260818_043912/"
        "20260818_043928_809709_combined.png"
    ),
    "natural_perimeter": Path(
        "/tmp/phase2_geometry_review_naturalistic_pair_naturalistic_20260818_044034/"
        "20260818_044050_035125_combined.png"
    ),
    "frontier": REPO / "rl_agent" / "plots" / "knob_accuracy_frontier.png",
    "delivery": REPO
    / "channel_condition_sweep"
    / "plots"
    / "fig2_delivery_heatmap.png",
}

PILOT = (
    REPO
    / "data_collection"
    / "experiments"
    / "phase2_paired_causal_v1"
    / "20260817_181354_pilot"
    / "phase2_pilot_positive_001"
)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def check_inputs() -> None:
    missing = [str(path) for path in ASSETS.values() if not path.is_file()]
    for role in ("helper", "recipient"):
        path = PILOT / role / "retained_inputs" / "frame_00156300_inputs.npz"
        if not path.is_file():
            missing.append(str(path))
    if missing:
        raise FileNotFoundError("missing deck inputs:\n" + "\n".join(missing))


def new_slide(
    title: str,
    *,
    subtitle: str | None = None,
    section: str | None = None,
    dark: bool = True,
) -> tuple[plt.Figure, plt.Axes]:
    fig = plt.figure(figsize=(W, H), dpi=DPI, facecolor=BG if dark else "#F6F8FB")
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W)
    ax.set_ylim(H, 0)
    ax.axis("off")
    ax.set_facecolor(BG if dark else "#F6F8FB")
    if section:
        text(ax, 0.62, 0.43, section.upper(), 10.5, CYAN, weight="bold")
    text(ax, 0.62, 0.78, title, 28, INK if dark else BLACK, weight="bold")
    if subtitle:
        text(ax, 0.64, 1.21, subtitle, 13.5, MUTED if dark else "#536579")
    ax.plot([0.62, 15.38], [1.48, 1.48], color=GRID if dark else "#D6DEE8", lw=1)
    return fig, ax


def text(
    ax: plt.Axes,
    x: float,
    y: float,
    value: str,
    size: float,
    color: str = INK,
    *,
    weight: str = "normal",
    ha: str = "left",
    va: str = "top",
    family: str = FONT,
    alpha: float = 1.0,
    linespacing: float = 1.15,
    rotation: float = 0,
) -> None:
    ax.text(
        x,
        y,
        value,
        fontsize=size,
        color=color,
        fontweight=weight,
        ha=ha,
        va=va,
        fontfamily=family,
        alpha=alpha,
        linespacing=linespacing,
        rotation=rotation,
    )


def wrapped(value: str, width: int) -> str:
    return "\n".join(textwrap.wrap(value, width=width, break_long_words=False))


def box(
    ax: plt.Axes,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = PANEL,
    edge: str | None = None,
    lw: float = 1.2,
    radius: float = 0.18,
    alpha: float = 1.0,
) -> FancyBboxPatch:
    patch = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle=f"round,pad=0.02,rounding_size={radius}",
        facecolor=fill,
        edgecolor=edge or fill,
        linewidth=lw,
        alpha=alpha,
    )
    ax.add_patch(patch)
    return patch


def pill(
    ax: plt.Axes,
    x: float,
    y: float,
    label: str,
    *,
    fill: str = PANEL_2,
    color: str = INK,
    edge: str | None = None,
    width: float | None = None,
) -> float:
    w = width or max(1.1, 0.095 * len(label) + 0.5)
    box(ax, x, y, w, 0.38, fill=fill, edge=edge or fill, radius=0.18)
    text(ax, x + w / 2, y + 0.19, label, 9.3, color, weight="bold", ha="center", va="center")
    return w


def arrow(
    ax: plt.Axes,
    start: tuple[float, float],
    end: tuple[float, float],
    *,
    color: str = CYAN,
    lw: float = 2.1,
    style: str = "-|>",
    connection: str = "arc3,rad=0",
    alpha: float = 1.0,
) -> None:
    ax.add_patch(
        FancyArrowPatch(
            start,
            end,
            arrowstyle=style,
            mutation_scale=14,
            linewidth=lw,
            color=color,
            connectionstyle=connection,
            alpha=alpha,
        )
    )


def add_image(
    ax: plt.Axes,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    mode: str = "cover",
    segment: str | None = None,
    edge: str | None = GRID,
    lw: float = 1.0,
    alpha: float = 1.0,
) -> None:
    image = Image.open(path).convert("RGB")
    if segment == "left":
        image = image.crop((0, 0, image.width // 2, image.height))
    elif segment == "right":
        image = image.crop((image.width // 2, 0, image.width, image.height))
    elif segment is not None:
        raise ValueError(f"unsupported image segment: {segment}")
    src_ratio = image.width / image.height
    dst_ratio = w / h
    if mode == "cover":
        if src_ratio > dst_ratio:
            new_width = int(image.height * dst_ratio)
            left = (image.width - new_width) // 2
            image = image.crop((left, 0, left + new_width, image.height))
        else:
            new_height = int(image.width / dst_ratio)
            top = (image.height - new_height) // 2
            image = image.crop((0, top, image.width, top + new_height))
    elif mode == "contain":
        if src_ratio > dst_ratio:
            draw_w = w
            draw_h = w / src_ratio
            draw_x = x
            draw_y = y + (h - draw_h) / 2
        else:
            draw_h = h
            draw_w = h * src_ratio
            draw_x = x + (w - draw_w) / 2
            draw_y = y
        ax.imshow(
            image,
            extent=(draw_x, draw_x + draw_w, draw_y + draw_h, draw_y),
            alpha=alpha,
            zorder=1,
        )
        if edge:
            box(ax, x, y, w, h, fill="none", edge=edge, lw=lw, radius=0.08, alpha=1.0)
        return
    else:
        raise ValueError(f"unsupported image placement mode: {mode}")
    ax.imshow(image, extent=(x, x + w, y + h, y), alpha=alpha, zorder=1)
    if edge:
        box(ax, x, y, w, h, fill="none", edge=edge, lw=lw, radius=0.08, alpha=1.0)


def footer(ax: plt.Axes, number: int, source: str | None = None) -> None:
    if source:
        text(ax, 0.64, 8.52, source, 8.1, "#7890A8")
    text(
        ax,
        15.35,
        8.52,
        f"{number:02d}  •  SCENESENSE AGENT  •  19 AUG 2026",
        8.1,
        "#7890A8",
        ha="right",
    )


def card(
    ax: plt.Axes,
    x: float,
    y: float,
    w: float,
    h: float,
    title_value: str,
    body: str,
    *,
    accent: str = CYAN,
    body_size: float = 12.5,
    icon: str | None = None,
) -> None:
    box(ax, x, y, w, h, fill=PANEL, edge=GRID, radius=0.16)
    ax.add_patch(FancyBboxPatch((x, y), 0.08, h, boxstyle="square,pad=0", facecolor=accent, edgecolor=accent))
    if icon:
        text(ax, x + 0.35, y + 0.30, icon, 19, accent, weight="bold", va="center")
        title_x = x + 0.77
    else:
        title_x = x + 0.32
    text(ax, title_x, y + 0.20, title_value, 13.2, INK, weight="bold")
    text(ax, x + 0.32, y + 0.78, body, body_size, MUTED, linespacing=1.22)


def equation_box(
    ax: plt.Axes,
    x: float,
    y: float,
    w: float,
    h: float,
    equation: str,
    *,
    label: str | None = None,
    accent: str = CYAN,
    size: float = 19,
) -> None:
    box(ax, x, y, w, h, fill="#0A1828", edge=accent, lw=1.2, radius=0.16)
    if label:
        pill(ax, x + 0.22, y + 0.18, label, fill=accent, color=BLACK)
    text(ax, x + w / 2, y + h * 0.60, equation, size, INK, ha="center", va="center", family="DejaVu Sans")


def save_slide(fig: plt.Figure, number: int, slug: str) -> Path:
    path = SLIDES / f"{number:02d}_{slug}.png"
    fig.savefig(path, dpi=DPI, facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def render_pilot_rgb_assets() -> dict[str, Path]:
    import numpy as np

    paths: dict[str, Path] = {}
    for role in ("helper", "recipient"):
        src = PILOT / role / "retained_inputs" / "frame_00156300_inputs.npz"
        with np.load(src) as arrays:
            rgb = arrays["frame_bgr"][..., ::-1]
        dst = OUT / f"pilot_frame_156300_{role}.png"
        Image.fromarray(rgb).save(dst)
        paths[role] = dst
    return paths


def slide_01() -> plt.Figure:
    fig = plt.figure(figsize=(W, H), dpi=DPI, facecolor=BG)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W)
    ax.set_ylim(H, 0)
    ax.axis("off")
    add_image(ax, ASSETS["worked"], 0, 0, W, H, mode="cover", edge=None)
    ax.add_patch(FancyBboxPatch((0, 0), W, H, boxstyle="square,pad=0", facecolor=BG, edgecolor=BG, alpha=0.72))
    ax.add_patch(FancyBboxPatch((0, 0), 9.9, H, boxstyle="square,pad=0", facecolor=BG, edgecolor=BG, alpha=0.74))
    pill(ax, 0.85, 0.72, "ADVISOR PROGRESS REVIEW", fill=CYAN, color=BLACK, width=2.55)
    text(ax, 0.85, 1.55, "SceneSense Agent", 34, INK, weight="bold")
    text(ax, 0.85, 2.30, "Causal, network-aware\ncooperative perception", 30, INK, weight="bold", linespacing=1.08)
    text(
        ax,
        0.88,
        4.10,
        "From state and action design → constrained reward →\npaired helper–recipient environment → RL decision gate",
        15.2,
        MUTED,
        linespacing=1.25,
    )
    box(ax, 0.85, 5.45, 6.45, 1.34, fill="#0B1C2F", edge=GRID, radius=0.14, alpha=0.94)
    text(ax, 1.18, 5.74, "Design principle", 11.2, CYAN, weight="bold")
    text(ax, 1.18, 6.15, "The network determines feasible fusion;\nthe observed hazard determines the deadline.", 16.5, INK, weight="bold")
    text(ax, 0.88, 8.20, "Abiodun  •  InterDigital / Northeastern University collaboration", 11.2, MUTED)
    return fig


def slide_02() -> plt.Figure:
    fig, ax = new_slide(
        "One scene defines the control problem",
        subtitle="The helper can expose a hazard before the recipient can see around the parked van.",
        section="Motivation",
    )
    add_image(ax, ASSETS["worked"], 0.62, 1.78, 14.76, 4.15, mode="contain", edge=GRID)
    # Approximate annotation positions in the accepted visual-review frame.
    ax.add_patch(Circle((5.20, 4.15), 0.33, fill=False, edgecolor=AMBER, linewidth=3))
    arrow(ax, (5.20, 3.77), (4.55, 3.25), color=AMBER, lw=2)
    pill(ax, 3.18, 2.82, "HELPER SEES PEDESTRIAN", fill=AMBER, color=BLACK, width=2.75)
    ax.add_patch(Circle((11.05, 3.66), 0.48, fill=False, edgecolor=CORAL, linewidth=3))
    arrow(ax, (11.05, 3.18), (11.83, 2.78), color=CORAL, lw=2)
    pill(ax, 11.72, 2.47, "RECIPIENT OCCLUDED", fill=CORAL, color=BLACK, width=2.45)
    card(ax, 0.62, 6.28, 4.52, 1.72, "Perception asymmetry", "The two vehicles do not possess the same evidence at the same time.", accent=AMBER, body_size=11.8)
    card(ax, 5.39, 6.28, 4.52, 1.72, "Resource asymmetry", "SPLIT consumes radio; LOCAL consumes UE compute; SKIP consumes freshness.", accent=CYAN, body_size=11.8)
    card(ax, 10.16, 6.28, 5.22, 1.72, "Decision consequence", "The action changes queue state, map uncertainty, warning time, and the next decision.", accent=TEAL, body_size=11.8)
    footer(ax, 2, "Source: visually accepted Epic midblock-van pedestrian geometry")
    return fig


def slide_03() -> plt.Figure:
    fig, ax = new_slide(
        "The agent is a causal two-stage controller",
        subtitle="Placement is decided before inference; publication is decided only after a result exists.",
        section="State diagram",
    )
    stages = [
        (0.72, "CAUSAL STATE", "past map • lagged network\nprior tracks • compute headroom", CYAN),
        (3.75, "PLACEMENT", "SPLIT_FEATURE\nLOCAL_INFER\nSKIP_INFERENCE", BLUE),
        (6.78, "INFERENCE", "head + edge tail\nor full local model\nor no new result", PURPLE),
        (9.81, "PUBLICATION", "PUBLISH_ALL\nHAZARD_SUBSET\nSKIP_PUBLICATION", AMBER),
        (12.84, "RECIPIENT", "transport • install\npropagate uncertainty\nwarning", TEAL),
    ]
    for x, name, body, accent in stages:
        box(ax, x, 2.23, 2.48, 3.24, fill=PANEL, edge=accent, lw=1.7, radius=0.16)
        pill(ax, x + 0.20, 2.48, name, fill=accent, color=BLACK, width=2.08)
        text(ax, x + 1.24, 3.35, body, 13.2, INK, ha="center", va="center", linespacing=1.36)
    for idx in range(len(stages) - 1):
        arrow(ax, (stages[idx][0] + 2.48, 3.85), (stages[idx + 1][0], 3.85), color=CYAN, lw=2.0)
    box(ax, 1.05, 6.20, 13.90, 1.16, fill="#26141C", edge=CORAL, lw=1.4, radius=0.14)
    text(ax, 1.34, 6.46, "FORBIDDEN FEEDBACK", 11.5, CORAL, weight="bold")
    text(ax, 3.25, 6.46, "same-frame detections • confidence • current track identity • map quality • CARLA truth", 13.6, INK, weight="bold")
    text(ax, 1.34, 6.94, "These outputs may influence the next decision only after their availability timestamp.", 11.8, MUTED)
    box(ax, 3.40, 7.73, 9.20, 0.50, fill="#0B2B2A", edge=TEAL, radius=0.18)
    text(ax, 8.0, 7.98, r"Causal invariant:   $t_{available}(x) \leq t_{decision}(x)$", 15.0, TEAL, weight="bold", ha="center", va="center", family="DejaVu Sans")
    footer(ax, 3, "Source: rl_agent/state_diagram.md and Phase-2 paired causal contract")
    return fig


def slide_04() -> plt.Figure:
    fig, ax = new_slide(
        "State design: what is known before the action?",
        subtitle="Every observation field is timestamped and audited at its consuming decision.",
        section="State",
    )
    equation_box(
        ax,
        0.72,
        1.83,
        14.56,
        1.18,
        r"$s_t^-=[\hat C_{t-\ell},\,\sigma_C,\,q_t,\,b_t^{flight},\,M_{t^-},P_{t^-},\,x_t^H,\,\tilde x_{t-\delta}^R,\,h_t^{local},\,a_{t-1},y_{t-1}]$",
        label="PRE-ACTION OBSERVATION",
        size=18,
    )
    allowed = [
        ("Network", "lagged capacity + confidence\nMCS / BLER / BSR • prior delivery", CYAN),
        ("Map", "installed tracks • AoI • covariance\nsource + capture/install provenance", TEAL),
        ("Motion", "helper-local pose • newest causally\nreceived recipient pose and age", AMBER),
        ("Runtime", "scheduler credit • in-flight work\nLOCAL latency/headroom • prior action", PURPLE),
    ]
    for idx, (name, body, accent) in enumerate(allowed):
        x = 0.72 + idx * 3.67
        card(ax, x, 3.42, 3.36, 2.18, name, body, accent=accent, body_size=11.5)
    box(ax, 0.72, 6.00, 7.08, 1.55, fill="#0B2B2A", edge=TEAL, radius=0.15)
    text(ax, 1.05, 6.28, "ALLOWED", 11, TEAL, weight="bold")
    text(ax, 1.05, 6.72, "prior completed causal tracks\nwith source/observed/available timestamps", 13.5, INK, weight="bold")
    box(ax, 8.20, 6.00, 7.08, 1.55, fill="#28131C", edge=CORAL, radius=0.15)
    text(ax, 8.53, 6.28, "FORBIDDEN", 11, CORAL, weight="bold")
    text(ax, 8.53, 6.72, "current-action outputs • shadow outputs\nGT actor IDs • future hazard labels", 13.5, INK, weight="bold")
    footer(ax, 4, "Design correction: no post-action information may select the action that produced it")
    return fig


def slide_05() -> plt.Figure:
    fig, ax = new_slide(
        "Action space: discrete semantics, measured operating points",
        subtitle="Continuous state and outcomes do not justify inventing unmeasured continuous actions.",
        section="Actions",
    )
    placement = [
        ("SPLIT_FEATURE", "head at helper\nfeature over uplink\nedge-tail inference", BLUE),
        ("LOCAL_INFER", "full model at helper\ncompact object record\nlate/object fusion", TEAL),
        ("SKIP_INFERENCE", "no new inference\nprior map propagates\nuncertainty grows", CORAL),
    ]
    for idx, (name, body, accent) in enumerate(placement):
        x = 0.72 + idx * 3.35
        box(ax, x, 2.05, 3.04, 2.02, fill=PANEL, edge=accent, lw=1.6, radius=0.14)
        pill(ax, x + 0.20, 2.28, name, fill=accent, color=BLACK, width=2.64)
        text(ax, x + 1.52, 3.24, body, 12.4, INK, ha="center", va="center", linespacing=1.28)
    text(ax, 0.72, 4.45, "AFTER A RESULT EXISTS", 10.5, AMBER, weight="bold")
    publications = [
        ("PUBLISH_ALL", AMBER),
        ("PUBLISH_HAZARD_SUBSET", AMBER),
        ("SKIP_PUBLICATION", CORAL),
    ]
    for idx, (name, accent) in enumerate(publications):
        x = 0.72 + idx * 3.35
        box(ax, x, 4.80, 3.04, 0.80, fill="#18243A", edge=accent, radius=0.13)
        text(ax, x + 1.52, 5.20, name, 11.4, accent, weight="bold", ha="center", va="center")
    add_image(ax, ASSETS["frontier"], 11.10, 2.05, 4.18, 3.55, mode="contain", edge=GRID)
    text(ax, 13.19, 5.77, "Measured frontier", 11.5, CYAN, weight="bold", ha="center")
    rows = [
        ("Discrete now", "placement • profile • publication", CYAN),
        ("Bounded grid now", "target FPS / update interval", AMBER),
        ("Continuous later", "FPS only after held-out interpolation", TEAL),
        ("Continuous outcomes", "latency • AoI • covariance • clearance", MUTED),
    ]
    for idx, (label, body, accent) in enumerate(rows):
        y = 6.20 + idx * 0.47
        text(ax, 0.82, y, label, 11.3, accent, weight="bold")
        text(ax, 3.02, y, body, 11.3, INK)
    box(ax, 8.55, 6.15, 6.73, 1.72, fill="#111F31", edge=GRID, radius=0.14)
    text(ax, 8.90, 6.45, "Algorithm implication", 11.0, PURPLE, weight="bold")
    text(ax, 8.90, 6.90, "Start with exact/rule/MPC over a masked finite catalog.\nIf learning is needed: DQN, discrete SAC, or masked PPO—not continuous SAC.", 12.8, INK, weight="bold", linespacing=1.25)
    footer(ax, 5, "Measured profiles are categorical; payload/accuracy interpolation is not assumed")
    return fig


def slide_06() -> plt.Figure:
    fig, ax = new_slide(
        "Worked loop: one midblock pedestrian decision",
        subtitle="Illustrative controller step using the reviewed scenario—not a measured policy result.",
        section="Scenario loop",
    )
    add_image(
        ax,
        ASSETS["worked"],
        0.62,
        1.80,
        7.20,
        4.72,
        mode="cover",
        segment="left",
        edge=GRID,
    )
    box(ax, 0.62, 1.80, 7.20, 0.34, fill="#111111", edge="#111111", radius=0.01)
    text(
        ax,
        0.82,
        1.97,
        "HELPER  |  controlled positive occlusion  |  Epic / 1280×720 / FOV 120°",
        10.0,
        INK,
        weight="bold",
        va="center",
    )
    ax.add_patch(Circle((4.98, 4.10), 0.31, fill=False, edgecolor=AMBER, linewidth=3))
    pill(ax, 0.90, 6.75, "HELPER EVIDENCE", fill=AMBER, color=BLACK, width=2.10)
    text(ax, 3.18, 6.79, "pedestrian emerges beside van", 11.5, INK, va="center")

    steps = [
        ("1  STATE", "Mild channel estimate\nrecipient map aging\nLOCAL headroom available", CYAN),
        ("2  ACTION", "SPLIT profile + FPS\nOR LOCAL + publication\nOR SKIP", BLUE),
        ("3  OUTCOME", "delivery / latency / bytes\nnew causal objects\nwarning or no warning", AMBER),
        ("4  NEXT STATE", "queue and estimator update\nmap age/covariance update\nprevious outcome recorded", TEAL),
    ]
    for idx, (name, body, accent) in enumerate(steps):
        y = 1.87 + idx * 1.48
        card(ax, 8.25, y, 7.03, 1.18, name, body, accent=accent, body_size=11.4)
        if idx < len(steps) - 1:
            arrow(ax, (11.77, y + 1.18), (11.77, y + 1.47), color=accent, lw=1.8)
    arrow(ax, (14.97, 7.52), (8.46, 7.52), color=TEAL, lw=1.8, connection="arc3,rad=0")
    arrow(ax, (8.46, 7.52), (8.46, 2.10), color=TEAL, lw=1.8, connection="arc3,rad=0")
    text(ax, 11.71, 7.70, "repeat at the next decision boundary", 10.6, TEAL, weight="bold", ha="center")
    footer(ax, 6, "Scenario: parked-van midblock pedestrian • helper/recipient synchronized Epic views")
    return fig


def slide_07() -> plt.Figure:
    fig, ax = new_slide(
        "The action changes network, map, and hazard state",
        subtitle="This temporal coupling—not a large state vector—is the reason learning may eventually help.",
        section="Transition model",
    )
    equation_box(ax, 0.72, 1.90, 4.62, 1.60, r"$q_{t+1}=\left[q_t+B(a_t)-S_t\right]_+$", label="NETWORK QUEUE", accent=BLUE, size=19)
    equation_box(ax, 5.69, 1.90, 4.62, 1.60, r"$z_{j,t+\Delta}=F(\Delta)z_{j,t}$", label="OBJECT STATE", accent=TEAL, size=19)
    equation_box(ax, 10.66, 1.90, 4.62, 1.60, r"$P_{j,t+\Delta}=FP_{j,t}F^\top+Q(\Delta)$", label="UNCERTAINTY", accent=AMBER, size=18)
    branches = [
        ("DELIVERED + INSTALLED", "sequence accepted\nAoI resets to capture age\ncovariance starts from new evidence", TEAL),
        ("DROPPED / EXPIRED", "queue service consumed\nno map installation\nold state keeps propagating", CORAL),
        ("SKIP", "no new work created\ncorrect if no service demand\nrisky if hazard debt accumulates", AMBER),
    ]
    for idx, (name, body, accent) in enumerate(branches):
        x = 0.72 + idx * 4.97
        card(ax, x, 4.10, 4.62, 2.03, name, body, accent=accent, body_size=11.8)
    equation_box(
        ax,
        2.00,
        6.62,
        12.00,
        1.12,
        r"$t_{install}(a_t)\;\leq\;t_{conflict}-\left(\tau_{pipeline,p95}+\tau_{reaction}+\tau_{brake}+\tau_{margin}\right)$",
        label="ACTIONABLE DEADLINE",
        accent=CORAL,
        size=17.5,
    )
    text(ax, 8.0, 8.05, "Deadline terms are frozen only when the common warning-to-braking adapter exists.", 10.8, MUTED, ha="center")
    footer(ax, 7, "Phase-2 uses state + covariance propagation; it does not reuse the frozen-object speed × AoI shortcut")
    return fig


def slide_08() -> plt.Figure:
    fig, ax = new_slide(
        "Reward v5: task value inside a constrained safe set",
        subtitle="Causality, feasibility, and safety are masks—not quantities the reward may buy its way around.",
        section="Reward",
    )
    equation_box(
        ax,
        0.72,
        1.83,
        14.56,
        1.24,
        r"$U_{task}=0.35\,\frac{mIoU_{seg}}{mIoU_{ref}}+0.40\,\frac{Recall_{ped}}{Recall_{ped,ref}}+0.25\,\frac{Recall_{veh}}{Recall_{veh,ref}}$",
        label="PERCEPTION UTILITY",
        accent=TEAL,
        size=19.5,
    )
    # Weight bars.
    bars = [
        ("Pedestrian recall", 0.40, CORAL, "highest supported safety-critical class weight"),
        ("Segmentation", 0.35, CYAN, "protect scene understanding from ROI collapse"),
        ("Vehicle recall", 0.25, AMBER, "explicit rather than hidden in object recall"),
    ]
    for idx, (label, value, color, note) in enumerate(bars):
        y = 3.52 + idx * 0.86
        text(ax, 0.78, y, label, 12.4, INK, weight="bold")
        box(ax, 3.25, y + 0.02, 4.15, 0.30, fill="#18304A", edge="#18304A", radius=0.12)
        box(ax, 3.25, y + 0.02, 4.15 * value / 0.40, 0.30, fill=color, edge=color, radius=0.12)
        text(ax, 7.64, y - 0.01, f"{int(value*100)}%", 13.2, color, weight="bold")
        text(ax, 8.38, y - 0.01, note, 11.2, MUTED)
    equation_box(
        ax,
        0.72,
        6.32,
        14.56,
        1.27,
        r"$R_t=w_UU_{task}-\lambda_r C_{PRB}-\lambda_c C_{local}-\lambda_q D_{hazard}-\lambda_s\mathbf{1}[a_t\neq a_{t-1}]-w_E M_{uncertainty}$",
        label="CANDIDATE PHASE-2 INNER SCORE",
        accent=CYAN,
        size=17.5,
    )
    text(ax, 0.88, 7.88, "No explicit ROI cost", 10.8, AMBER, weight="bold")
    text(ax, 2.68, 7.88, "ROI damage is already measured in task utility—another term would double count it.", 10.8, MUTED)
    text(ax, 9.17, 7.88, "Weights remain unfrozen", 10.8, CORAL, weight="bold")
    text(ax, 11.37, 7.88, "until causal LOCAL/OAI tables exist.", 10.8, MUTED)
    footer(ax, 8, "Source: REWARD_FORMULATION.md v5 and Phase-2 constraint catalog")
    return fig


def slide_09() -> plt.Figure:
    fig, ax = new_slide(
        "Reward terms shape behavior—without replacing constraints",
        subtitle="The desired controller is useful, efficient, stable, and honest about graceful degradation.",
        section="Reward effects",
    )
    effects = [
        ("Task utility ↑", "Preserve segmentation and class recall; avoid profiles whose compression destroys useful evidence.", TEAL),
        ("PRB-time cost ↑", "Prefer compact measured profiles or LOCAL when the same task value can be delivered more cheaply.", BLUE),
        ("Hazard debt ↑", "Publish before the recipient deadline; SKIP becomes costly only when service is actually owed.", CORAL),
        ("Switching cost ↑", "Avoid oscillating SPLIT ↔ LOCAL on noisy estimates; retain hysteresis without hiding deadline misses.", PURPLE),
        ("Expected uncertainty ↑", "Prefer margin inside the admitted set; tail uncertainty remains a shield/service constraint.", AMBER),
    ]
    for idx, (name, body, accent) in enumerate(effects):
        y = 1.88 + idx * 1.08
        box(ax, 0.72, y, 9.25, 0.82, fill=PANEL, edge=GRID, radius=0.13)
        pill(ax, 0.94, y + 0.20, name, fill=accent, color=BLACK, width=2.08)
        text(ax, 3.28, y + 0.17, wrapped(body, 72), 11.8, INK, linespacing=1.16)
    box(ax, 10.35, 1.88, 4.93, 2.25, fill="#28131C", edge=CORAL, radius=0.15)
    text(ax, 10.72, 2.19, "NO GLOBAL SKIP PENALTY", 12.0, CORAL, weight="bold")
    text(ax, 10.72, 2.70, "Correct abstention in empty/fresh scenes must remain free.\n\nCharge the consequence: missed warning, deadline debt, or growing map uncertainty.", 12.3, INK, weight="bold", linespacing=1.22)
    equation_box(
        ax,
        10.35,
        4.48,
        4.93,
        2.08,
        r"$C_{stop}=\left(\frac{(d_{min}-d)_+}{d_{min}}\right)^2+\lambda_e\left(\frac{(d-d_{comfort})_+}{d_s}\right)^2$",
        label="FUTURE STOP COST",
        accent=AMBER,
        size=14.5,
    )
    text(ax, 10.56, 6.82, "Only after every arm uses the same warning-to-braking adapter.", 10.5, MUTED)
    box(ax, 10.35, 7.26, 4.93, 0.56, fill="#0B2B2A", edge=TEAL, radius=0.14)
    text(ax, 12.82, 7.54, "Collision + minimum clearance remain hard constraints", 10.6, TEAL, weight="bold", ha="center", va="center")
    footer(ax, 9, "Stopping too early and stopping too close are both undesirable; clearance alone is gameable")
    return fig


def slide_10() -> plt.Figure:
    fig, ax = new_slide(
        "Constraints are lexicographic, not one weighted soup",
        subtitle="A scalar reward ranks only the actions that survive the higher-priority contracts.",
        section="Constraint rank",
    )
    levels = [
        ("0", "CAUSAL + STRUCTURAL", "S0–S2 • no leakage • measured actions • recipient/order isolation", CORAL, 5.2),
        ("1", "PHYSICAL SAFETY", "P1–P2 • collision + minimum surface clearance after common actuation", AMBER, 6.7),
        ("2", "DEADLINE / SERVICE", "N2 + M2 • useful observed-hazard evidence before recipient deadline", PURPLE, 8.2),
        ("3", "NETWORK + COMPUTE FEASIBILITY", "N1/N3 + K1/K2 + M1 • capacity, queue, LOCAL headroom, valid contribution", BLUE, 10.1),
        ("4", "TASK UTILITY", "M3 • segmentation + pedestrian recall + vehicle recall", TEAL, 12.1),
        ("5", "EFFICIENCY + COMFORT", "PRB-time • bytes • compute • switching • stop placement • jerk • progress", CYAN, 14.1),
    ]
    y0 = 1.93
    center = 8.0
    for idx, (rank, name, body, color, width) in enumerate(levels):
        x = center - width / 2
        h = 0.88
        y = y0 + idx * 1.02
        box(ax, x, y, width, h, fill=PANEL, edge=color, lw=1.7, radius=0.14)
        pill(ax, x + 0.18, y + 0.25, f"RANK {rank}", fill=color, color=BLACK, width=1.15)
        text(ax, x + 1.55, y + 0.18, name, 12.2, INK, weight="bold")
        text(ax, x + 1.55, y + 0.52, body, 9.9, MUTED)
    box(ax, 2.38, 8.00, 11.24, 0.34, fill="#26141C", edge=CORAL, radius=0.16)
    text(ax, 8.0, 8.17, "No reward weight can authorize causal leakage, an unsupported action, or a hard safety violation.", 10.9, INK, weight="bold", ha="center", va="center")
    footer(ax, 10, "Operational IDs S/N/K/M/P/O are kept separate from paper contributions C1–C4")
    return fig


def slide_11() -> plt.Figure:
    fig, ax = new_slide(
        "The constraint domains are coupled",
        subtitle="Network and physical context meet through map installation time, uncertainty, and the hazard deadline.",
        section="Relationships",
    )
    nodes = [
        (0.72, 2.15, 3.20, 1.20, "PHYSICAL ENVIRONMENT", "object motion • occlusion\nclosing speed • recipient path", AMBER),
        (0.72, 5.26, 3.20, 1.20, "NETWORK + COMPUTE", "capacity estimate • queue\nLOCAL p95 • in-flight work", BLUE),
        (5.02, 2.15, 3.20, 1.20, "HAZARD DEADLINE", "required response margin\nservice debt", CORAL),
        (5.02, 5.26, 3.20, 1.20, "ACTION ADMISSION", "measured catalog ∩ radio\n∩ compute ∩ service", CYAN),
        (9.32, 2.15, 3.20, 1.20, "MAP STATE", "install time • AoI\nstate + covariance", TEAL),
        (9.32, 5.26, 3.20, 1.20, "REWARD RANKING", "task value − costs\n− debt − switching", PURPLE),
        (13.20, 3.70, 2.08, 1.28, "ACTION", "SPLIT\nLOCAL\nSKIP", WHITE),
    ]
    for x, y, w, h, name, body, accent in nodes:
        box(ax, x, y, w, h, fill=PANEL if accent != WHITE else "#E9F0F8", edge=accent, lw=1.6, radius=0.14)
        text(ax, x + 0.22, y + 0.20, name, 10.7, accent if accent != WHITE else BLACK, weight="bold")
        text(ax, x + 0.22, y + 0.61, body, 10.8, INK if accent != WHITE else BLACK, weight="bold", linespacing=1.16)
    arrow(ax, (3.92, 2.75), (5.02, 2.75), color=AMBER)
    arrow(ax, (3.92, 5.86), (5.02, 5.86), color=BLUE)
    arrow(ax, (8.22, 2.75), (9.32, 2.75), color=CORAL)
    arrow(ax, (8.22, 5.86), (9.32, 5.86), color=CYAN)
    arrow(ax, (10.92, 3.35), (10.92, 5.26), color=TEAL)
    arrow(ax, (12.52, 5.86), (13.20, 4.58), color=PURPLE)
    arrow(ax, (13.20, 4.14), (12.52, 2.75), color=CYAN, connection="arc3,rad=-0.18")
    arrow(ax, (14.24, 4.98), (2.32, 7.36), color=TEAL, connection="arc3,rad=0.12", alpha=0.8)
    text(ax, 7.40, 7.30, "outcome becomes lagged state at the next decision", 11.2, TEAL, weight="bold", ha="center")
    box(ax, 2.14, 7.78, 11.72, 0.46, fill="#0B2B2A", edge=TEAL, radius=0.18)
    text(ax, 8.0, 8.01, "Network decides feasible fusion level  •  observed hazard decides the deadline  •  reward ranks the survivors", 11.3, INK, weight="bold", ha="center", va="center")
    footer(ax, 11, "Graceful degradation is explicit when no action satisfies every service target")
    return fig


def slide_12() -> plt.Figure:
    fig, ax = new_slide(
        "Physical and environmental constraints",
        subtitle="The agent must be evaluated in traffic that is legal, realistic, diverse, and attributable.",
        section="Physical world",
    )
    add_image(ax, ASSETS["signalized"], 0.62, 1.82, 7.58, 4.25, mode="cover", edge=GRID)
    add_image(ax, ASSETS["queue"], 8.55, 1.82, 6.83, 4.25, mode="cover", edge=GRID)
    pill(ax, 0.88, 5.60, "PEDESTRIAN HAZARD", fill=AMBER, color=BLACK, width=2.05)
    pill(ax, 8.82, 5.60, "VEHICLE REVEAL", fill=BLUE, color=BLACK, width=1.84)
    constraints = [
        ("P1", "Collision", "report now; hard/terminal after warning actuation", CORAL),
        ("P2", "Surface clearance", "oriented ego/hazard boxes; advisor freezes d_min", AMBER),
        ("P3–P5", "Stopping + comfort", "too-close / too-early • deceleration • jerk • route progress", PURPLE),
        ("O1", "Scene validity", "legal lanes • realistic pedestrian speed • no gridlock • zero actor leaks", TEAL),
    ]
    for idx, (cid, name, body, accent) in enumerate(constraints):
        x = 0.72 + (idx % 2) * 7.36
        y = 6.42 + (idx // 2) * 0.78
        box(ax, x, y, 6.95, 0.62, fill=PANEL, edge=GRID, radius=0.12)
        pill(ax, x + 0.16, y + 0.12, cid, fill=accent, color=BLACK, width=0.82)
        text(ax, x + 1.17, y + 0.12, name, 11.4, INK, weight="bold")
        text(ax, x + 2.63, y + 0.12, body, 10.2, MUTED)
    footer(ax, 12, "Primary CARLA capture: Town10HD_Opt • Epic • native 10 Hz • matched positive/benign futures")
    return fig


def slide_13() -> plt.Figure:
    fig, ax = new_slide(
        "Network and compute constraints decide what is feasible",
        subtitle="The measured uplink has a payload-dependent delivery cliff; LOCAL is constrained by its own p95 latency.",
        section="Network + compute",
    )
    add_image(ax, ASSETS["delivery"], 0.62, 1.80, 8.45, 5.85, mode="contain", edge=GRID)
    equation_box(
        ax,
        9.45,
        1.85,
        5.83,
        1.30,
        r"$r_{offer}(a)=B_{frame}(a)\,f(a)\leq \kappa\,C_{LCB}(s_t)$",
        label="N1 RADIO ADMISSION",
        accent=BLUE,
        size=17,
    )
    equation_box(
        ax,
        9.45,
        3.47,
        5.83,
        1.30,
        r"$T_{local,p95}(p)\leq\Delta t\quad\wedge\quad f\leq f_{sust}(p)$",
        label="K1 LOCAL ADMISSION",
        accent=TEAL,
        size=17,
    )
    cards = [
        ("Large feature", "high task ceiling • queue/latency risk", CORAL),
        ("Compact feature", "measured quality/cost compromise", AMBER),
        ("LOCAL compact", "radio-light • compute-heavy • calibration pending", TEAL),
    ]
    for idx, (name, body, accent) in enumerate(cards):
        y = 5.12 + idx * 0.80
        box(ax, 9.45, y, 5.83, 0.60, fill=PANEL, edge=GRID, radius=0.12)
        text(ax, 9.75, y + 0.12, name, 11.0, accent, weight="bold")
        text(ax, 11.55, y + 0.12, body, 10.4, INK)
    footer(ax, 13, "Measured OAI RFsim channel sweep • delivery, queue latency, and capacity are action-dependent")
    return fig


def slide_14() -> plt.Figure:
    fig, ax = new_slide(
        "Building the agent environment",
        subtitle="An offline/replay Gym-style environment preserves causal timing and keeps training off the CARLA data path.",
        section="Environment",
    )
    api = [
        ("reset(group_id, arm)", "load one immutable trajectory\nclear isolated map/queue/controller state", CYAN),
        ("observe()", "emit pre-action allowlisted state\nwith provenance and availability", BLUE),
        ("step(a_place)", "advance selected inference/transport model\nproduce result or no-result", PURPLE),
        ("step(a_publish)", "install contribution or retain old map\ncompute reward + next state", AMBER),
        ("close()", "flush structured logs\nrelease replay handles deterministically", TEAL),
    ]
    for idx, (name, body, accent) in enumerate(api):
        x = 0.72 + idx * 3.06
        box(ax, x, 2.05, 2.72, 2.42, fill=PANEL, edge=accent, lw=1.5, radius=0.15)
        text(ax, x + 1.36, 2.38, name, 11.1, accent, weight="bold", ha="center")
        text(ax, x + 1.36, 3.18, body, 10.8, INK, weight="bold", ha="center", va="center", linespacing=1.25)
        if idx < len(api) - 1:
            arrow(ax, (x + 2.72, 3.25), (x + 3.06, 3.25), color=accent, lw=1.7)
    box(ax, 0.72, 5.05, 7.02, 2.16, fill="#0A1828", edge=GRID, radius=0.15)
    text(ax, 1.05, 5.35, "CLOCKS + DATA", 11.2, CYAN, weight="bold")
    text(ax, 1.05, 5.82, "• CARLA sensors and detection: native 10 Hz\n• optional surrogate policy clock: 20 Hz, interpolated only\n• paired arms replay identical captured evidence\n• truth stream is joined only by the evaluator", 12.1, INK, linespacing=1.35)
    box(ax, 8.08, 5.05, 7.20, 2.16, fill="#0A1828", edge=GRID, radius=0.15)
    text(ax, 8.42, 5.35, "STRUCTURED OUTPUT", 11.2, TEAL, weight="bold")
    text(ax, 8.42, 5.82, "episode_id • trajectory_group • frame/time\nstate provenance • masks • both actions • outcome\nreward terms • deadline debt • bytes/PRB/compute\nwarning/physical evaluation in separate namespaces", 12.1, INK, linespacing=1.35)
    box(ax, 2.80, 7.58, 10.40, 0.52, fill="#26141C", edge=CORAL, radius=0.18)
    text(ax, 8.0, 7.84, "Training never blocks CARLA or the real-time sensing path—use replay buffers and offline/async updates.", 11.7, INK, weight="bold", ha="center", va="center")
    footer(ax, 14, "Environment contract: explicit observation/action spaces, deterministic reset, structured episode logs")
    return fig


def thumbnail(ax: plt.Axes, path: Path, x: float, y: float, w: float, h: float, label: str, accent: str) -> None:
    add_image(ax, path, x, y, w, h, mode="cover", edge=accent, lw=1.2)
    box(ax, x, y + h - 0.46, w, 0.46, fill=BG, edge=BG, radius=0.0, alpha=0.84)
    text(ax, x + 0.16, y + h - 0.23, label, 9.2, INK, weight="bold", va="center")


def slide_15() -> plt.Figure:
    fig, ax = new_slide(
        "Suite A creates decision opportunities",
        subtitle="Six visually accepted hazard families span class, occlusion geometry, speed, deadline, and traffic density.",
        section="Scenarios",
    )
    items = [
        (ASSETS["curbside"], "1  Curbside bus • pedestrian", AMBER),
        (ASSETS["signalized"], "2  Signalized corner • pedestrian", AMBER),
        (ASSETS["worked"], "3  Midblock van • pedestrian", AMBER),
        (ASSETS["cross_traffic"], "4  Occluded cross-traffic • vehicle", BLUE),
        (ASSETS["pullout"], "5  Parked pull-out • vehicle", BLUE),
        (ASSETS["queue"], "6  Queue reveal • vehicle", BLUE),
    ]
    for idx, (path, label, accent) in enumerate(items):
        row, col = divmod(idx, 3)
        thumbnail(ax, path, 0.72 + col * 5.02, 1.88 + row * 2.80, 4.65, 2.42, label, accent)
    box(ax, 2.52, 7.65, 10.96, 0.48, fill="#0B2B2A", edge=TEAL, radius=0.16)
    text(ax, 8.0, 7.89, "Each positive has a matched benign twin; traffic/weather/seeds stay paired and group-locked.", 11.4, INK, weight="bold", ha="center", va="center")
    footer(ax, 15, "Suite A = designed decision opportunities • cyclists excluded until the perception contract supports them")
    return fig


def slide_16() -> plt.Figure:
    fig, ax = new_slide(
        "Suite B keeps the result honest",
        subtitle="Naturalistic operation is reported with the same endpoints so the curated suite cannot flatter the controller invisibly.",
        section="Evaluation design",
    )
    add_image(ax, ASSETS["natural_signal"], 0.72, 1.88, 7.04, 3.60, mode="cover", edge=TEAL)
    add_image(ax, ASSETS["natural_perimeter"], 8.24, 1.88, 7.04, 3.60, mode="cover", edge=TEAL)
    pill(ax, 1.00, 5.05, "SIGNALIZED DEMO ROUTE", fill=TEAL, color=BLACK, width=2.30)
    pill(ax, 8.52, 5.05, "SAFE-PERIMETER ROUTE", fill=TEAL, color=BLACK, width=2.30)
    cards = [
        ("Designed Suite A", "helper-visible hazards + matched benign negatives\nregime-bounded cooperation claim", AMBER),
        ("Naturalistic Suite B", "hazards at natural prevalence\nno-hazard runs remain in denominator", TEAL),
        ("Grouped inference", "trajectory/scenario clusters—not frames\npositive + benign never split", CYAN),
    ]
    for idx, (name, body, accent) in enumerate(cards):
        x = 0.72 + idx * 5.02
        card(ax, x, 6.02, 4.65, 1.47, name, body, accent=accent, body_size=11.2)
    # 20/20/60 split visual.
    x0, y0, width = 4.10, 7.83, 7.80
    box(ax, x0, y0, width, 0.34, fill="#132A42", edge="#132A42", radius=0.12)
    box(ax, x0, y0, width * 0.20, 0.34, fill=CYAN, edge=CYAN, radius=0.12)
    box(ax, x0 + width * 0.20, y0, width * 0.20, 0.34, fill=AMBER, edge=AMBER, radius=0.0)
    box(ax, x0 + width * 0.40, y0, width * 0.60, 0.34, fill=TEAL, edge=TEAL, radius=0.12)
    text(ax, x0 - 0.20, y0 + 0.17, "split", 9.6, MUTED, ha="right", va="center")
    text(ax, x0 + width * 0.10, y0 + 0.17, "20% CAL", 8.7, BLACK, weight="bold", ha="center", va="center")
    text(ax, x0 + width * 0.30, y0 + 0.17, "20% VAL", 8.7, BLACK, weight="bold", ha="center", va="center")
    text(ax, x0 + width * 0.70, y0 + 0.17, "60% UNTOUCHED TEST", 8.7, BLACK, weight="bold", ha="center", va="center")
    footer(ax, 16, "Pilot excluded • primary C2 endpoint: recipient-specific actionable warning success + warning lead")
    return fig


def slide_17(pilot_assets: dict[str, Path]) -> plt.Figure:
    fig, ax = new_slide(
        "What is already banked",
        subtitle="The pilot proves the causal capture and evaluation chain—not cooperation performance.",
        section="Progress",
    )
    add_image(ax, pilot_assets["helper"], 0.72, 1.84, 4.42, 2.58, mode="cover", edge=CYAN)
    add_image(ax, pilot_assets["recipient"], 5.35, 1.84, 4.42, 2.58, mode="cover", edge=TEAL)
    pill(ax, 0.94, 3.93, "HELPER • FRAME 156300", fill=CYAN, color=BLACK, width=2.32)
    pill(ax, 5.57, 3.93, "RECIPIENT • SAME FRAME", fill=TEAL, color=BLACK, width=2.45)
    box(ax, 10.10, 1.84, 5.18, 2.58, fill=PANEL, edge=GRID, radius=0.14)
    checks = [
        "aligned RGB + radar retained",
        "unfiltered detections + causal tracks",
        "placement and publication audit",
        "recipient map install + warning",
        "separate truth + future adjudication",
    ]
    text(ax, 10.45, 2.12, "ACCEPTED PILOT CHAIN", 11.2, TEAL, weight="bold")
    for idx, item in enumerate(checks):
        text(ax, 10.47, 2.60 + idx * 0.34, "•", 14.0, TEAL, weight="bold")
        text(ax, 10.78, 2.60 + idx * 0.34, item, 10.7, INK)
    milestones = [
        ("Measured inputs", "profile payload/quality\nchannel/queue surface\nexact sensor contract", CYAN),
        ("Control design", "causal state/action\nreward v5\nranked constraints", BLUE),
        ("Map + evaluation", "v2 contribution schema\nwarning engine\nfuture-hazard adjudicator", PURPLE),
        ("Scenario design", "6 designed families\n2 naturalistic routes\nEpic visual acceptance", AMBER),
        ("Pilot", "helper→map→warning\nraw recoverability\nstructural gates PASS", TEAL),
    ]
    for idx, (name, body, accent) in enumerate(milestones):
        x = 0.72 + idx * 3.03
        card(ax, x, 5.04, 2.70, 2.23, name, body, accent=accent, body_size=10.7)
    box(ax, 2.90, 7.67, 10.20, 0.48, fill="#1E1B35", edge=PURPLE, radius=0.17)
    text(ax, 8.0, 7.91, "Current scientific gate: calibration sufficiency → frozen operating point → causal baseline ladder", 11.5, INK, weight="bold", ha="center", va="center")
    footer(ax, 17, "Pilot batch: 20260817_181354_pilot • authoritative evaluation/verification namespaces retained")
    return fig


def slide_18() -> plt.Figure:
    fig, ax = new_slide(
        "Next steps: answer whether learning is necessary",
        subtitle="The staged plan stops at every scientific gate; no long collection or RL training is automatic.",
        section="Roadmap",
    )
    stages = [
        ("1", "CALIBRATION AUDIT", "prove raw replay sufficiency\nrun bounded 96-setting replay", CYAN),
        ("2", "FREEZE + POWER", "choose one warning setting\nverify nuisance / miss / power gates", AMBER),
        ("3", "VALIDATION → TEST", "grouped Suite A + B\nlocal C2, then identical bytes over OAI", TEAL),
        ("4", "LOCAL TABLE", "payload vs object count\nlocal p50/p95 + sustainable FPS\ncompact-record OAI latency", BLUE),
        ("5", "CONTROLLER LADDER", "exact → rule → λ-RDO/AoI → MPC\nRL only after residual sequential gap", PURPLE),
        ("6", "PHYSICAL OVERRIDE", "fixed warning-to-braking adapter\nclearance • stop band • jerk • collision", CORAL),
    ]
    for idx, (num, name, body, accent) in enumerate(stages):
        row, col = divmod(idx, 3)
        x = 0.72 + col * 5.02
        y = 1.92 + row * 2.35
        box(ax, x, y, 4.65, 1.88, fill=PANEL, edge=accent, lw=1.5, radius=0.15)
        ax.add_patch(Circle((x + 0.48, y + 0.48), 0.26, facecolor=accent, edgecolor=accent))
        text(ax, x + 0.48, y + 0.48, num, 11.5, BLACK, weight="bold", ha="center", va="center")
        text(ax, x + 0.86, y + 0.25, name, 11.4, accent, weight="bold")
        text(ax, x + 0.32, y + 0.87, body, 11.6, INK, weight="bold", linespacing=1.24)
        if col < 2:
            arrow(ax, (x + 4.65, y + 0.94), (x + 5.02, y + 0.94), color=accent, lw=1.7)
    box(ax, 0.72, 6.80, 14.56, 1.02, fill="#0A1828", edge=GRID, radius=0.15)
    text(ax, 1.02, 7.06, "ADVISOR DECISIONS TO FREEZE", 10.8, AMBER, weight="bold")
    text(ax, 1.02, 7.47, "minimum clearance + comfort stop band  •  reaction/braking assumptions  •  LOCAL hardware target  •  staged calibration authorization", 12.2, INK, weight="bold")
    box(ax, 3.17, 8.01, 9.66, 0.35, fill="#0B2B2A", edge=TEAL, radius=0.15)
    text(ax, 8.0, 8.185, "Success is a defensible controller and feasibility envelope—even if the simplest controller wins.", 10.9, INK, weight="bold", ha="center", va="center")
    footer(ax, 18, "RL gate: held-out causal service/reward gap after exact, rule, and MPC baselines")
    return fig


def presenter_notes() -> str:
    return """# Presenter notes — SceneSense Agent progress review

## Slide 1 — SceneSense Agent
Lead with the refined objective: the project is not committed to RL. We are
building a causal cooperative-perception controller and will use measured
evidence to decide whether learning is necessary.

## Slide 2 — One scene defines the problem
Point to the pedestrian in the helper view and the van/occlusion in the
recipient view. The agent trades three scarce resources: uplink, local compute,
and map freshness.

## Slide 3 — Causal two-stage controller
Emphasize that placement and publication happen at different times. The red
band is the key validity safeguard: the current inference output cannot select
the inference action that produced it.

## Slide 4 — State
Walk through network, prior map, motion, and runtime fields. Each field carries
source and availability timestamps. CARLA truth is evaluation-only.

## Slide 5 — Action space
Profiles and semantic actions are discrete because they are measured. FPS may
be made continuous only after interpolation is validated. This is why
continuous SAC is not the current choice.

## Slide 6 — Worked scenario loop
This is an illustrative step, not a result. Describe how the same scene may
lead to SPLIT, LOCAL, or SKIP depending on lagged channel, map state, and local
headroom—and how that outcome changes the next state.

## Slide 7 — Transition model
Explain the three coupled dynamics: queue service, object state, and covariance.
The actionable deadline becomes physically meaningful only after reaction and
braking assumptions are frozen.

## Slide 8 — Reward v5
Pedestrian recall is highest, segmentation remains substantial, and vehicle
recall stays explicit. The masks happen before the inner objective. Reward
weights are not yet tuned from the tiny pilot.

## Slide 9 — Reward effects
Explain why there is no global SKIP penalty. Correct abstention is useful;
unserved hazard debt is the failure. Stopping distance enters only after a
common warning-actuation adapter makes it attributable.

## Slide 10 — Constraint ranking
This is a constrained/lexicographic problem. A large reward cannot compensate
for a causal leak or unsupported action.

## Slide 11 — Constraint relationships
Physical context sets the deadline. Network/compute determine feasibility.
Transport outcomes determine map age and uncertainty. Reward ranks the safe,
feasible survivors.

## Slide 12 — Physical constraints
Distinguish current report-only physical outcomes from later hard constraints.
Mention legal lanes, realistic traffic, matched futures, and actor cleanup as
experiment-validity requirements.

## Slide 13 — Network/compute constraints
Use the measured heatmap to show the payload-dependent cliff. LOCAL is not a
free fallback: it needs measured local latency and sustainable FPS.

## Slide 14 — Environment
Training is offline/replay and Gym-style, decoupled from CARLA. Paired arms see
the same immutable evidence; truth is attached only by evaluation.

## Slide 15 — Designed scenarios
The six families intentionally create decision opportunities across pedestrians
and vehicles. Every positive has a benign twin.

## Slide 16 — Naturalistic suite
Designed cases can flatter a controller. Suite B is the honest denominator,
reported with the same metrics and grouped confidence intervals.

## Slide 17 — Banked progress
The pilot proves the complete causal artifact chain and warning computability.
Do not quote its lead as performance evidence.

## Slide 18 — Next steps
Ask the advisor to freeze the physical response parameters and LOCAL hardware
target, and to approve only the next bounded calibration stage. End with the
simplest-controller-that-works principle.
"""


def build_pptx(slide_paths: Sequence[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    note_sections = presenter_notes().split("\n## Slide ")[1:]
    note_bodies = []
    for section in note_sections:
        lines = section.splitlines()
        note_bodies.append("\n".join(lines[1:]).strip())
    sources_by_slide = [
        ["rl_agent/ADVISOR_PROGRESS_BRIEF_2026-08-19.md"],
        ["phase2_map_sharing/PHASE2_SUITE_AB_DESIGN.md", "SOURCE_MANIFEST.json: worked"],
        ["rl_agent/state_diagram.md", "phase2_map_sharing/PHASE2_PAIRED_CAUSAL_CORPUS_SPEC.md"],
        ["rl_agent/state_diagram.md", "phase2_map_sharing/PHASE2_CONSTRAINT_CATALOG.md"],
        ["rl_agent/plots/knob_accuracy_frontier.png", "rl_agent/ADVISOR_PROGRESS_BRIEF_2026-08-19.md"],
        ["SOURCE_MANIFEST.json: worked", "phase2_map_sharing/PHASE2_SUITE_AB_DESIGN.md"],
        ["phase2_map_sharing/PHASE2_CONSTRAINT_CATALOG.md"],
        ["rl_agent/REWARD_FORMULATION.md", "phase2_map_sharing/PHASE2_CONSTRAINT_CATALOG.md"],
        ["phase2_map_sharing/PHASE2_CONSTRAINT_CATALOG.md"],
        ["phase2_map_sharing/PHASE2_CONSTRAINT_CATALOG.md"],
        ["phase2_map_sharing/PHASE2_CONSTRAINT_CATALOG.md", "rl_agent/state_diagram.md"],
        ["phase2_map_sharing/PHASE2_CONSTRAINT_CATALOG.md", "SOURCE_MANIFEST.json: CARLA geometry assets"],
        ["channel_condition_sweep/CHANNEL_SWEEP_RESULTS.md", "channel_condition_sweep/plots/fig2_delivery_heatmap.png"],
        ["phase2_map_sharing/PHASE2_PAIRED_CAUSAL_CORPUS_SPEC.md"],
        ["phase2_map_sharing/PHASE2_SUITE_AB_DESIGN.md", "SOURCE_MANIFEST.json: Suite A assets"],
        ["phase2_map_sharing/PHASE2_SUITE_AB_DESIGN.md", "SOURCE_MANIFEST.json: Suite B assets"],
        ["data_collection/experiments/phase2_paired_causal_v1/20260817_181354_pilot"],
        ["rl_agent/ADVISOR_PROGRESS_BRIEF_2026-08-19.md", "phase2_map_sharing/PHASE2_PAIRED_CAUSAL_CORPUS_SPEC.md"],
    ]
    if len(note_bodies) != len(slide_paths) or len(sources_by_slide) != len(slide_paths):
        raise RuntimeError("speaker-note/source count does not match slide count")
    for index, path in enumerate(slide_paths):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        source_lines = "\n".join(f"- {source}" for source in sources_by_slide[index])
        slide.notes_slide.notes_text_frame.text = (
            f"{note_bodies[index]}\n\nReproducibility sources:\n{source_lines}\n\n"
            "Full hashes: rl_agent/presentation/agent_progress_20260819/SOURCE_MANIFEST.json"
        )
    # Core document metadata.
    props = prs.core_properties
    props.title = "SceneSense Agent — Causal Network-Aware Cooperative Perception"
    props.subject = "Advisor progress review: state, action, reward, constraints, environment, and next gates"
    props.author = "Abiodun"
    props.comments = "Generated from versioned SceneSense agent design artifacts; slide artwork is reproducible."
    prs.save(PPTX)


def build_manifest(slide_paths: Sequence[Path], pilot_assets: dict[str, Path]) -> None:
    manifest = {
        "schema": "scenesense.agent_advisor_deck_manifest.v1",
        "deck": str(PPTX.relative_to(REPO)),
        "slide_count": len(slide_paths),
        "slides": [
            {"ordinal": index + 1, "path": str(path.relative_to(REPO)), "sha256": sha256(path)}
            for index, path in enumerate(slide_paths)
        ],
        "source_assets": [
            {
                "name": name,
                "path": str(path) if path.is_absolute() and not str(path).startswith(str(REPO)) else str(path.relative_to(REPO)),
                "sha256": sha256(path),
            }
            for name, path in {**ASSETS, **{f"pilot_{k}": v for k, v in pilot_assets.items()}}.items()
        ],
        "claim_sources": [
            "rl_agent/state_diagram.md",
            "rl_agent/REWARD_FORMULATION.md",
            "phase2_map_sharing/PHASE2_CONSTRAINT_CATALOG.md",
            "phase2_map_sharing/PHASE2_PAIRED_CAUSAL_CORPUS_SPEC.md",
            "phase2_map_sharing/WARNING_EVALUATION_DESIGN_FREEZE.md",
            "phase2_map_sharing/PHASE2_SUITE_AB_DESIGN.md",
            "rl_agent/ADVISOR_PROGRESS_BRIEF_2026-08-19.md",
        ],
        "claim_boundary": (
            "The accepted two-trajectory pilot is structural/computability evidence only. "
            "The deck reports no historical noncausal controller result as a deployable policy result."
        ),
    }
    MANIFEST.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")


def build_contact_sheet(slide_paths: Sequence[Path]) -> None:
    """Create a lightweight visual-QA sheet without changing deck contents."""
    thumb_w, thumb_h = 480, 270
    margin, label_h, columns = 24, 34, 3
    rows = math.ceil(len(slide_paths) / columns)
    canvas = Image.new(
        "RGB",
        (
            columns * thumb_w + (columns + 1) * margin,
            rows * (thumb_h + label_h) + (rows + 1) * margin,
        ),
        color=(7, 17, 31),
    )
    from PIL import ImageDraw, ImageFont

    draw = ImageDraw.Draw(canvas)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 20)
    except OSError:
        font = ImageFont.load_default()
    for index, path in enumerate(slide_paths):
        row, column = divmod(index, columns)
        x = margin + column * (thumb_w + margin)
        y = margin + row * (thumb_h + label_h + margin)
        with Image.open(path) as slide_image:
            thumbnail = slide_image.convert("RGB").resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        canvas.paste(thumbnail, (x, y))
        draw.text((x, y + thumb_h + 5), f"{index + 1:02d}  {path.stem}", fill=(175, 194, 215), font=font)
    canvas.save(CONTACT_SHEET)


def build() -> None:
    check_inputs()
    OUT.mkdir(parents=True, exist_ok=True)
    SLIDES.mkdir(parents=True, exist_ok=True)
    pilot_assets = render_pilot_rgb_assets()
    builders = [
        ("title", slide_01),
        ("problem_scene", slide_02),
        ("causal_state_diagram", slide_03),
        ("state_contract", slide_04),
        ("action_space", slide_05),
        ("worked_loop", slide_06),
        ("transition_model", slide_07),
        ("reward_v5", slide_08),
        ("reward_effects", slide_09),
        ("constraint_rank", slide_10),
        ("constraint_relationships", slide_11),
        ("physical_constraints", slide_12),
        ("network_compute", slide_13),
        ("agent_environment", slide_14),
        ("suite_a", slide_15),
        ("suite_b", slide_16),
        ("banked_progress", lambda: slide_17(pilot_assets)),
        ("next_steps", slide_18),
    ]
    slide_paths = [save_slide(builder(), idx, slug) for idx, (slug, builder) in enumerate(builders, start=1)]
    build_contact_sheet(slide_paths)
    build_pptx(slide_paths)
    NOTES.write_text(presenter_notes(), encoding="utf-8")
    build_manifest(slide_paths, pilot_assets)
    print(
        json.dumps(
            {
                "verdict": "PASS",
                "pptx": str(PPTX),
                "slide_count": len(slide_paths),
                "presenter_notes": str(NOTES),
                "manifest": str(MANIFEST),
                "contact_sheet": str(CONTACT_SHEET),
            },
            indent=2,
        )
    )


if __name__ == "__main__":
    build()
