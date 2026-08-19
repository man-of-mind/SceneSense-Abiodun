#!/usr/bin/env python3
"""Rebuild the 2026-08-19 SceneSense agent deck on the InterDigital template.

Differences from ``build_agent_progress_deck.py`` (v1):

* The deck is built **on top of** ``interdigital_template.pptx`` so the brand
  master, theme palette, Century Gothic type, logo, and footer are inherited
  rather than imitated.
* Titles, subtitles, cards, tables, bars, and flow diagrams are **native
  PowerPoint shapes and text** instead of one flattened bitmap per slide, so
  every slide stays editable and stays crisp at any zoom.
* Mathematics is rendered as **LaTeX-style Computer Modern** (matplotlib
  mathtext, ``fontset="cm"``) into transparent PNGs and placed inline.  Word
  identifiers are wrapped in ``\\mathrm{}`` so ``mIoU_seg`` reads as an upright
  roman label rather than a product of italic letters.

Outputs land in ``rl_agent/presentation/agent_progress_20260819``.
"""

from __future__ import annotations

import hashlib
import json
import math
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "rl_agent" / "presentation" / "agent_progress_20260819"
TEMPLATE = OUT / "interdigital_template.pptx"
BUILD = OUT / "build_v2"
PPTX = OUT / "SceneSense_Agent_Progress_2026-08-19_ID.pptx"
NOTES_MD = OUT / "PRESENTER_NOTES.md"
MANIFEST = OUT / "SOURCE_MANIFEST_v2.json"

DECK_DATE = "19 August 2026"
COPYRIGHT = "©2026 InterDigital, Inc. All Rights Reserved."

# ---------------------------------------------------------------- brand palette
NAVY = "010644"  # theme accent2 - primary ink
BLUE = "00A8EC"  # theme accent1 - primary accent
DEEP = "0A5C93"  # theme accent5
TEAL = "00C9C8"  # theme accent4
ORANGE = "FFAA59"  # theme accent3
GREY = "49515C"  # theme dk2 - secondary text
LINE = "D5D9DF"  # hairline rules and card borders
LT = "E6E6E7"  # theme lt2 - inert fills
RED = "D8455F"  # functional only: forbidden / hazard / danger
WHITE = "FFFFFF"

FONT = "Century Gothic"

# ------------------------------------------------------------------ page metrics
PAGE_W = 13.3333
PAGE_H = 7.5
ML = 0.92
CW = 11.50
MR = ML + CW
EYEBROW_T = 0.30
TITLE_T = 0.60
TITLE_H = 0.62
SUB_T = 1.22
RULE_Y = 1.66
BODY_T = 1.84
BODY_B = 6.76
FOOT_T = 6.95

DPI_EQ = 460


def mix(hex_a: str, hex_b: str, t: float) -> str:
    """Blend two hex colours; ``t`` is the weight of ``hex_b``."""

    a = tuple(int(hex_a[i : i + 2], 16) for i in (0, 2, 4))
    b = tuple(int(hex_b[i : i + 2], 16) for i in (0, 2, 4))
    return "".join(f"{round(a[i] + (b[i] - a[i]) * t):02X}" for i in range(3))


def tint(accent: str, t: float = 0.90) -> str:
    """Pale wash of an accent, used for panel fills on a white page."""

    return mix(accent, WHITE, t)


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


# ------------------------------------------------------------------- equations
#
# Every identifier that is a *word* is wrapped in \mathrm{} so it renders
# upright, the way it would in a paper.  Single-letter quantities stay italic.
EQ: dict[str, str] = {
    "causal_invariant": r"$t_{\mathrm{available}}(x)\ \leq\ t_{\mathrm{decision}}(x)$",
    "state": (
        r"$s_t^{-}=\left[\,\hat{C}_{t-\ell},\ \sigma_C,\ q_t,\ b_t^{\mathrm{flight}},\ "
        r"M_{t^-},\ P_{t^-},\ x_t^{H},\ \tilde{x}^{R}_{t-\delta},\ "
        r"h_t^{\mathrm{local}},\ a_{t-1},\ y_{t-1}\,\right]$"
    ),
    "queue": r"$q_{t+1}=\left[\,q_t+B(a_t)-S_t\,\right]_{+}$",
    "object": r"$z_{j,\,t+\Delta}=F(\Delta)\,z_{j,t}$",
    "covariance": r"$P_{j,\,t+\Delta}=F\,P_{j,t}\,F^{\top}+Q(\Delta)$",
    "deadline": (
        r"$t_{\mathrm{install}}(a_t)\ \leq\ t_{\mathrm{conflict}}-\left("
        r"\tau_{\mathrm{pipe},95}+\tau_{\mathrm{react}}+\tau_{\mathrm{brake}}"
        r"+\tau_{\mathrm{margin}}\right)$"
    ),
    "utility": (
        r"$U_{\mathrm{task}}=0.35\,\frac{\mathrm{mIoU}_{\mathrm{seg}}}{\mathrm{mIoU}_{\mathrm{ref}}}"
        r"+0.40\,\frac{\mathrm{Recall}_{\mathrm{ped}}}{\mathrm{Recall}_{\mathrm{ped,ref}}}"
        r"+0.25\,\frac{\mathrm{Recall}_{\mathrm{veh}}}{\mathrm{Recall}_{\mathrm{veh,ref}}}$"
    ),
    "reward": (
        r"$R_t=w_U\,U_{\mathrm{task}}-\lambda_r C_{\mathrm{PRB}}"
        r"-\lambda_c C_{\mathrm{local}}-\lambda_q D_{\mathrm{hazard}}"
        r"-\lambda_s\,\mathbf{1}\!\left[a_t\neq a_{t-1}\right]-w_E\,M_{\mathrm{unc}}$"
    ),
    "stop_cost": (
        r"$C_{\mathrm{stop}}=\left(\frac{(d_{\min}-d)_{+}}{d_{\min}}\right)^{\!2}"
        r"+\lambda_e\left(\frac{(d-d_{\mathrm{comf}})_{+}}{d_s}\right)^{\!2}$"
    ),
    "radio_admission": (
        r"$r_{\mathrm{offer}}(a)=B_{\mathrm{frame}}(a)\,f(a)"
        r"\ \leq\ \kappa\,C_{\mathrm{LCB}}(s_t)$"
    ),
    "local_admission": (
        r"$T^{\,95}_{\mathrm{local}}(p)\ \leq\ \Delta t"
        r"\quad\wedge\quad f\ \leq\ f_{\mathrm{sust}}(p)$"
    ),
}

# ----------------------------------------------------------------- source assets
ASSETS = {
    "worked": Path(
        "/tmp/phase2_geometry_review_midblock_van_positive_20260818_001417/"
        "20260818_001435_896264_combined.png"
    ),
    "curbside": Path(
        "/tmp/phase2_geometry_review_positive_20260815_041815/"
        "20260815_041825_634287_combined.png"
    ),
    "signalized": Path(
        "/tmp/phase2_geometry_review_signalized_corner_positive_20260817_224735/"
        "20260817_224753_327439_combined.png"
    ),
    "cross_traffic": Path(
        "/tmp/phase2_geometry_review_cross_traffic_vehicle_positive_20260818_024451/"
        "20260818_024513_531786_combined.png"
    ),
    "pullout": Path(
        "/tmp/phase2_geometry_review_parked_vehicle_pullout_positive_20260818_031013/"
        "20260818_031036_689670_combined.png"
    ),
    "queue": Path(
        "/tmp/phase2_geometry_review_queue_reveal_vehicle_positive_20260818_033946/"
        "20260818_034002_180110_combined.png"
    ),
    "natural_signal": Path(
        "/tmp/phase2_geometry_review_naturalistic_pair_naturalistic_20260818_043912/"
        "20260818_043928_809709_combined.png"
    ),
    "natural_perimeter": Path(
        "/tmp/phase2_geometry_review_naturalistic_pair_naturalistic_20260818_044034/"
        "20260818_044050_035125_combined.png"
    ),
    "frontier": REPO / "rl_agent" / "plots" / "knob_accuracy_frontier.png",
    "delivery": REPO
    / "channel_condition_sweep"
    / "plots"
    / "fig2_delivery_heatmap.png",
    "pilot_helper": OUT / "pilot_frame_156300_helper.png",
    "pilot_recipient": OUT / "pilot_frame_156300_recipient.png",
}


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1 << 20), b""):
            digest.update(block)
    return digest.hexdigest()


def check_inputs() -> None:
    missing = [str(p) for p in ASSETS.values() if not p.is_file()]
    if not TEMPLATE.is_file():
        missing.append(str(TEMPLATE))
    if missing:
        raise FileNotFoundError("missing deck inputs:\n  " + "\n  ".join(missing))


# ============================================================ rendering helpers
class EquationRenderer:
    """Render mathtext to tight transparent PNGs, cached on disk."""

    def __init__(self, directory: Path) -> None:
        self.directory = directory
        self.directory.mkdir(parents=True, exist_ok=True)
        plt.rcParams["mathtext.fontset"] = "cm"
        plt.rcParams["mathtext.default"] = "it"
        self._cache: dict[tuple[str, str], tuple[Path, float]] = {}

    def render(self, key: str, expression: str, colour: str) -> tuple[Path, float]:
        """Return (png path, width/height aspect) for one equation."""

        cache_key = (key, colour)
        if cache_key in self._cache:
            return self._cache[cache_key]
        path = self.directory / f"eq_{key}_{colour}.png"
        figure = plt.figure(figsize=(0.1, 0.1))
        figure.text(0, 0, expression, fontsize=24, color=f"#{colour}")
        try:
            figure.savefig(
                path,
                dpi=DPI_EQ,
                transparent=True,
                bbox_inches="tight",
                pad_inches=0.02,
            )
        except Exception as exc:  # pragma: no cover - surfaces bad mathtext loudly
            raise RuntimeError(f"could not render equation {key!r}: {exc}") from exc
        finally:
            plt.close(figure)
        with Image.open(path) as image:
            aspect = image.width / image.height
        self._cache[cache_key] = (path, aspect)
        return path, aspect


# The geometry-review composites carry a dark instrumentation caption strip
# across the top 7-8% of the frame.  Trimming it keeps the slide's own labels
# the only caption the audience reads.
REVIEW_BANNER_FRACTION = 0.085


def cover_crop(
    path: Path,
    aspect: float,
    cache: Path,
    *,
    segment: str | None = None,
    crop_top: float = 0.0,
) -> Path:
    """Centre-crop an image to ``aspect`` so it fills a box without distortion."""

    cache.mkdir(parents=True, exist_ok=True)
    tag = f"{path.stem}_{segment or 'full'}_{aspect:.4f}_{crop_top:.3f}.png"
    destination = cache / tag
    if destination.is_file():
        return destination
    with Image.open(path) as source:
        image = source.convert("RGB")
        if crop_top > 0:
            image = image.crop((0, int(round(image.height * crop_top)), image.width, image.height))
        if segment == "left":
            image = image.crop((0, 0, image.width // 2, image.height))
        elif segment == "right":
            image = image.crop((image.width // 2, 0, image.width, image.height))
        source_aspect = image.width / image.height
        if source_aspect > aspect:
            width = int(round(image.height * aspect))
            left = (image.width - width) // 2
            image = image.crop((left, 0, left + width, image.height))
        else:
            height = int(round(image.width / aspect))
            top = (image.height - height) // 2
            image = image.crop((0, top, image.width, top + height))
        image.save(destination)
    return destination


def pad_to_aspect(path: Path, aspect: float, cache: Path, *, background: str = WHITE) -> Path:
    """Letterbox an image onto a flat panel so charts are never cropped."""

    cache.mkdir(parents=True, exist_ok=True)
    destination = cache / f"{path.stem}_pad_{aspect:.4f}_{background}.png"
    if destination.is_file():
        return destination
    with Image.open(path) as source:
        image = source.convert("RGB")
        source_aspect = image.width / image.height
        if source_aspect > aspect:
            width, height = image.width, int(round(image.width / aspect))
        else:
            width, height = int(round(image.height * aspect)), image.height
        canvas = Image.new("RGB", (width, height), tuple(int(background[i : i + 2], 16) for i in (0, 2, 4)))
        canvas.paste(image, ((width - image.width) // 2, (height - image.height) // 2))
        canvas.save(destination)
    return destination


# ================================================================= deck builder
@dataclass
class Slide:
    number: int
    title: str
    source: str


class Deck:
    LAYOUT_TITLE = 0
    LAYOUT_CONTENT = 6  # "Title Only" - brand title band, everything else free

    def __init__(self) -> None:
        self.prs = Presentation(str(TEMPLATE))
        self._drop_template_slides()
        self.equations = EquationRenderer(BUILD / "equations")
        self.image_cache = BUILD / "images"
        self.index: list[Slide] = []

    # ------------------------------------------------------------------ plumbing
    def _drop_template_slides(self) -> None:
        id_list = self.prs.slides._sldIdLst
        for element in list(id_list):
            self.prs.part.drop_rel(element.rId)
            id_list.remove(element)

    def save(self) -> None:
        self.prs.save(str(PPTX))

    # -------------------------------------------------------------------- slides
    def title_slide(self, title: str, subtitle: str, notes: str) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_TITLE])
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        head = placeholders[0]
        head.text_frame.text = title
        for paragraph in head.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.name = FONT
                run.font.color.rgb = rgb(WHITE)
        sub = placeholders[1]
        frame = sub.text_frame
        frame.word_wrap = True
        lines = subtitle.split("\n")
        for position, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if position == 0 else frame.add_paragraph()
            run = paragraph.add_run()
            run.text = line
            run.font.name = FONT
            run.font.size = Pt(15 if position == 0 else 12)
            run.font.bold = position == 0
            run.font.color.rgb = rgb(BLUE if position == 0 else mix(LT, NAVY, 0.20))
            paragraph.space_after = Pt(6)
        self._fill_footer(slide, None)
        self._set_notes(slide, notes)
        self.index.append(Slide(1, title, "title"))

    def content_slide(
        self,
        number: int,
        eyebrow: str,
        title: str,
        subtitle: str,
        *,
        source: str,
        notes: str,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_CONTENT])
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        head = placeholders[0]
        head.left, head.top = Inches(ML), Inches(TITLE_T)
        head.width, head.height = Inches(CW), Inches(TITLE_H)
        head.text_frame.word_wrap = True
        head.text_frame.text = title
        for paragraph in head.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(self._title_size(title))
                run.font.bold = True
                run.font.name = FONT
                run.font.color.rgb = rgb(NAVY)

        # Eyebrow: accent tab + uppercase section label.
        self.rect(slide, ML, EYEBROW_T + 0.045, 0.13, 0.13, fill=BLUE)
        self.text(
            slide,
            ML + 0.24,
            EYEBROW_T - 0.02,
            2.9,
            0.26,
            eyebrow.upper(),
            size=8.5,
            bold=True,
            colour=BLUE,
            spacing=90,
        )
        self.text(
            slide,
            MR - 3.2,
            EYEBROW_T - 0.02,
            3.2,
            0.26,
            f"SCENESENSE AGENT  •  {number:02d}",
            size=8.5,
            bold=True,
            colour=mix(GREY, WHITE, 0.35),
            align=PP_ALIGN.RIGHT,
        )
        if subtitle:
            self.text(
                slide,
                ML,
                SUB_T,
                CW,
                0.34,
                subtitle,
                size=11,
                colour=GREY,
            )
        self.hline(slide, ML, RULE_Y, CW, colour=LINE)
        if source:
            self.text(
                slide,
                ML,
                BODY_B + 0.06,
                CW - 1.2,
                0.24,
                source,
                size=7.5,
                colour=mix(GREY, WHITE, 0.3),
            )
        self._fill_footer(slide, number)
        self._set_notes(slide, notes)
        self.index.append(Slide(number, title, source))
        return slide

    @staticmethod
    def _title_size(title: str) -> float:
        """Keep every content title on a single line inside the brand band."""

        length = len(title)
        if length <= 46:
            return 24.0
        if length <= 56:
            return 21.0
        return 19.0

    def _fill_footer(self, slide, number: int | None) -> None:
        for placeholder in slide.placeholders:
            idx = placeholder.placeholder_format.idx
            frame = placeholder.text_frame
            if str(placeholder.placeholder_format.type).startswith("FOOTER"):
                frame.text = COPYRIGHT
            elif str(placeholder.placeholder_format.type).startswith("SLIDE_NUMBER"):
                frame.text = "" if number is None else f"{number}"
            else:
                continue
            for paragraph in frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(8)
                    run.font.color.rgb = rgb(mix(GREY, WHITE, 0.3))

    @staticmethod
    def _set_notes(slide, notes: str) -> None:
        slide.notes_slide.notes_text_frame.text = notes.strip()

    # ------------------------------------------------------------------ graphics
    @staticmethod
    def _plain(shape) -> None:
        shape.shadow.inherit = False
        if shape.has_text_frame:
            shape.text_frame.text = ""

    def rect(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        fill: str | None = None,
        edge: str | None = None,
        lw: float = 0.75,
        rounded: bool = False,
        radius: float = 0.10,
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        if rounded:
            shape.adjustments[0] = min(0.5, radius / min(w, h))
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)
        if edge is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = rgb(edge)
            shape.line.width = Pt(lw)
        self._plain(shape)
        return shape

    def ellipse(self, slide, x: float, y: float, w: float, h: float, *, fill=None, edge=None, lw=2.0):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)
        if edge is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = rgb(edge)
            shape.line.width = Pt(lw)
        self._plain(shape)
        return shape

    def chevron(self, slide, x: float, y: float, w: float, h: float, *, fill: str, edge: str | None = None):
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if edge is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = rgb(edge)
            shape.line.width = Pt(0.75)
        self._plain(shape)
        return shape

    def hline(self, slide, x: float, y: float, w: float, *, colour: str = LINE, lw: float = 0.75):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y)
        )
        line.line.color.rgb = rgb(colour)
        line.line.width = Pt(lw)
        return line

    def arrow(
        self,
        slide,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        *,
        colour: str = BLUE,
        lw: float = 1.5,
    ):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = rgb(colour)
        line.line.width = Pt(lw)
        properties = line.line._get_or_add_ln()
        tail = properties.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        properties.append(tail)
        return line

    def text(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        *,
        size: float = 10,
        bold: bool = False,
        italic: bool = False,
        colour: str = NAVY,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
        line_spacing: float = 1.12,
        space_after: float = 2.0,
        spacing: int | None = None,
        font: str = FONT,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = anchor
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        for position, line in enumerate(str(value).split("\n")):
            paragraph = frame.paragraphs[0] if position == 0 else frame.add_paragraph()
            paragraph.alignment = align
            paragraph.line_spacing = line_spacing
            paragraph.space_after = Pt(space_after)
            run = paragraph.add_run()
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = rgb(colour)
            if spacing is not None:
                run.font._rPr.set("spc", str(spacing))
        return box

    def pill(
        self,
        slide,
        x: float,
        y: float,
        label: str,
        *,
        fill: str = BLUE,
        colour: str = WHITE,
        w: float | None = None,
        h: float = 0.24,
        size: float = 8,
    ) -> float:
        width = w if w is not None else max(0.74, 0.082 * len(label) + 0.30)
        shape = self.rect(slide, x, y, width, h, fill=fill, rounded=True, radius=h / 2)
        frame = shape.text_frame
        frame.word_wrap = False
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = label
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = rgb(colour)
        run.font._rPr.set("spc", "60")
        return width

    def card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        heading: str,
        body: str,
        *,
        accent: str = BLUE,
        heading_size: float = 10.5,
        body_size: float = 8.8,
        washed: bool = False,
    ):
        self.rect(
            slide,
            x,
            y,
            w,
            h,
            fill=tint(accent, 0.94) if washed else WHITE,
            edge=mix(accent, WHITE, 0.55),
            rounded=True,
            radius=0.09,
        )
        self.rect(slide, x, y + 0.10, 0.075, h - 0.20, fill=accent)
        self.text(slide, x + 0.24, y + 0.12, w - 0.42, 0.26, heading, size=heading_size, bold=True, colour=NAVY)
        if body:
            self.text(
                slide,
                x + 0.24,
                y + 0.12 + 0.235 * (heading_size / 10.5),
                w - 0.42,
                h - 0.42,
                body,
                size=body_size,
                colour=GREY,
                line_spacing=1.22,
            )

    def panel(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        accent: str = BLUE,
        label: str | None = None,
        wash: float = 0.93,
    ):
        self.rect(
            slide,
            x,
            y,
            w,
            h,
            fill=tint(accent, wash),
            edge=mix(accent, WHITE, 0.45),
            rounded=True,
            radius=0.09,
        )
        if label:
            self.pill(slide, x + 0.16, y + 0.14, label, fill=accent)

    def picture(
        self,
        slide,
        path: Path,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        mode: str = "cover",
        segment: str | None = None,
        edge: str | None = LINE,
        crop_top: float = 0.0,
    ):
        if mode == "cover":
            prepared = cover_crop(path, w / h, self.image_cache, segment=segment, crop_top=crop_top)
            box = (x, y, w, h)
        elif mode == "contain":
            prepared = pad_to_aspect(path, w / h, self.image_cache)
            box = (x, y, w, h)
        elif mode == "fit":
            # Preserve the figure's own aspect: scale to fit and centre, so a
            # chart is never cropped and never padded with dead white bands.
            prepared = path
            with Image.open(path) as image:
                aspect = image.width / image.height
            width = min(w, h * aspect)
            height = width / aspect
            box = (x + (w - width) / 2, y + (h - height) / 2, width, height)
        else:
            raise ValueError(f"unsupported picture mode: {mode}")
        slide.shapes.add_picture(
            str(prepared), Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3])
        )
        if edge:
            self.rect(slide, box[0], box[1], box[2], box[3], fill=None, edge=edge, lw=0.75)
        return box

    def equation(
        self,
        slide,
        key: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        colour: str = NAVY,
        max_height: float | None = None,
    ) -> None:
        """Place one LaTeX-rendered equation centred in the box (x, y, w, h)."""

        path, aspect = self.equations.render(key, EQ[key], colour)
        limit_h = max_height if max_height is not None else h
        width = min(w, limit_h * aspect)
        height = width / aspect
        slide.shapes.add_picture(
            str(path),
            Inches(x + (w - width) / 2),
            Inches(y + (h - height) / 2),
            Inches(width),
            Inches(height),
        )

    def equation_panel(
        self,
        slide,
        key: str,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        label: str,
        accent: str = BLUE,
        eq_colour: str | None = None,
        pad: float = 0.30,
    ) -> None:
        self.panel(slide, x, y, w, h, accent=accent, label=label)
        top = y + 0.46
        self.equation(
            slide,
            key,
            x + pad,
            top,
            w - 2 * pad,
            h - (top - y) - 0.18,
            colour=eq_colour or NAVY,
            max_height=max(0.30, h - (top - y) - 0.30),
        )

    def bar(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        fraction: float,
        *,
        colour: str,
    ) -> None:
        self.rect(slide, x, y, w, h, fill=LT, rounded=True, radius=h / 2)
        if fraction > 0:
            self.rect(slide, x, y, max(h, w * fraction), h, fill=colour, rounded=True, radius=h / 2)

    def kv_row(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        key: str,
        value: str,
        *,
        accent: str = BLUE,
        key_w: float = 2.1,
        size: float = 8.8,
    ) -> None:
        self.text(slide, x, y, key_w, 0.24, key, size=size, bold=True, colour=accent)
        self.text(slide, x + key_w, y, w - key_w, 0.24, value, size=size, colour=GREY)


# =================================================================== the content
NOTES = {
    1: """Lead with the refined objective: the project is not committed to RL. We are building a
causal cooperative-perception controller and will use measured evidence to decide whether
learning is necessary.""",
    2: """Point to the pedestrian in the helper view and the van/occlusion in the recipient view.
The agent trades three scarce resources: uplink, local compute, and map freshness.""",
    3: """Emphasize that placement and publication happen at different times. The red band is the
key validity safeguard: the current inference output cannot select the inference action that
produced it.""",
    4: """Walk through network, prior map, motion, and runtime fields. Each field carries source and
availability timestamps. CARLA truth is evaluation-only.""",
    5: """Profiles and semantic actions are discrete because they are measured. FPS may be made
continuous only after interpolation is validated. This is why continuous SAC is not the current
choice.""",
    6: """This is an illustrative step, not a result. Describe how the same scene may lead to SPLIT,
LOCAL, or SKIP depending on lagged channel, map state, and local headroom - and how that outcome
changes the next state.""",
    7: """Explain the three coupled dynamics: queue service, object state, and covariance. The
actionable deadline becomes physically meaningful only after reaction and braking assumptions are
frozen.""",
    8: """Pedestrian recall is highest, segmentation remains substantial, and vehicle recall stays
explicit. The masks happen before the inner objective. Reward weights are not yet tuned from the
tiny pilot.""",
    9: """Explain why there is no global SKIP penalty. Correct abstention is useful; unserved hazard
debt is the failure. Stopping distance enters only after a common warning-actuation adapter makes
it attributable.""",
    10: """This is a constrained/lexicographic problem. A large reward cannot compensate for a causal
leak or unsupported action.""",
    11: """Physical context sets the deadline. Network/compute determine feasibility. Transport
outcomes determine map age and uncertainty. Reward ranks the safe, feasible survivors.""",
    12: """Distinguish current report-only physical outcomes from later hard constraints. Mention
legal lanes, realistic traffic, matched futures, and actor cleanup as experiment-validity
requirements.""",
    13: """Use the measured heatmap to show the payload-dependent cliff. LOCAL is not a free
fallback: it needs measured local latency and sustainable FPS.""",
    14: """Training is offline/replay and Gym-style, decoupled from CARLA. Paired arms see the same
immutable evidence; truth is attached only by evaluation.""",
    15: """The six families intentionally create decision opportunities across pedestrians and
vehicles. Every positive has a benign twin.""",
    16: """Designed cases can flatter a controller. Suite B is the honest denominator, reported with
the same metrics and grouped confidence intervals.""",
    17: """The pilot proves the complete causal artifact chain and warning computability. Do not
quote its lead as performance evidence.""",
    18: """Ask the advisor to freeze the physical response parameters and LOCAL hardware target, and
to approve only the next bounded calibration stage. End with the simplest-controller-that-works
principle.""",
}


def slide_01(deck: Deck) -> None:
    deck.title_slide(
        "SceneSense Agent",
        "Causal, network-aware cooperative perception\n"
        "State and action design → constrained reward → paired helper–recipient "
        "environment → RL decision gate\n"
        f"Advisor progress review  •  {DECK_DATE}  •  "
        "Abiodun Ganiyu  •  InterDigital / Northeastern University",
        NOTES[1],
    )
    # The InterDigital title layout is a dark navy field carrying the logo, so
    # title-slide type has to be light rather than the navy used on body slides.
    slide = deck.prs.slides[0]
    deck.rect(
        slide, 1.67, 6.16, 10.0, 0.62,
        fill=mix(NAVY, BLUE, 0.16), edge=BLUE, rounded=True, radius=0.09,
    )
    deck.text(
        slide, 1.87, 6.30, 9.6, 0.38,
        "The network determines feasible fusion;  the observed hazard determines the deadline.",
        size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER,
    )


def slide_02(deck: Deck) -> None:
    slide = deck.content_slide(
        2,
        "Motivation",
        "One scene defines the control problem",
        "The helper can expose a hazard before the recipient can see around the parked van.",
        source="Source: visually accepted Epic midblock-van pedestrian geometry (Town10HD_Opt)",
        notes=NOTES[2],
    )
    deck.picture(slide, ASSETS["worked"], ML, BODY_T, CW, 3.02, mode="cover", crop_top=REVIEW_BANNER_FRACTION)
    deck.pill(slide, ML + 0.14, BODY_T + 2.62, "HELPER  •  SEES PEDESTRIAN", fill=ORANGE, colour=NAVY, w=2.35)
    deck.pill(slide, ML + 5.90, BODY_T + 2.62, "RECIPIENT  •  OCCLUDED", fill=RED, w=2.05)
    deck.ellipse(slide, ML + 4.05, BODY_T + 1.28, 0.46, 0.46, edge=ORANGE, lw=2.25)
    deck.ellipse(slide, ML + 8.62, BODY_T + 1.16, 0.60, 0.60, edge=RED, lw=2.25)
    cards = [
        (
            "Perception asymmetry",
            "The two vehicles do not possess the same evidence at the same time.",
            ORANGE,
        ),
        (
            "Resource asymmetry",
            "SPLIT consumes radio; LOCAL consumes UE compute; SKIP consumes freshness.",
            BLUE,
        ),
        (
            "Decision consequence",
            "The action changes queue state, map uncertainty, warning time, and the next decision.",
            TEAL,
        ),
    ]
    for index, (heading, body, accent) in enumerate(cards):
        deck.card(slide, ML + index * 3.90, 5.20, 3.70, 1.32, heading, body, accent=accent)


def slide_03(deck: Deck) -> None:
    slide = deck.content_slide(
        3,
        "State diagram",
        "The agent is a causal two-stage controller",
        "Placement is decided before inference; publication is decided only after a result exists.",
        source="Source: rl_agent/state_diagram.md and the Phase-2 paired causal contract",
        notes=NOTES[3],
    )
    stages = [
        ("CAUSAL STATE", "past map • lagged network\nprior tracks • compute headroom", BLUE),
        ("PLACEMENT", "SPLIT_FEATURE\nLOCAL_INFER\nSKIP_INFERENCE", DEEP),
        ("INFERENCE", "head + edge tail\nfull local model\nor no new result", NAVY),
        ("PUBLICATION", "PUBLISH_ALL\nHAZARD_SUBSET\nSKIP_PUBLICATION", ORANGE),
        ("RECIPIENT", "transport • install\npropagate uncertainty\nwarning", TEAL),
    ]
    width, gap = 2.10, 0.25
    top, height = BODY_T + 0.14, 1.30
    for index, (name, body, accent) in enumerate(stages):
        x = ML + index * (width + gap)
        deck.rect(slide, x, top, width, height, fill=tint(accent, 0.94), edge=mix(accent, WHITE, 0.4), rounded=True, radius=0.09)
        deck.rect(slide, x, top, width, 0.30, fill=accent, rounded=True, radius=0.09)
        deck.rect(slide, x, top + 0.16, width, 0.14, fill=accent)
        deck.text(
            slide, x, top + 0.055, width, 0.24, name,
            size=8.5, bold=True, colour=WHITE if accent != ORANGE else NAVY,
            align=PP_ALIGN.CENTER, spacing=50,
        )
        deck.text(
            slide, x + 0.12, top + 0.34, width - 0.24, height - 0.44, body,
            size=8.6, colour=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.30,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        if index < len(stages) - 1:
            deck.arrow(slide, x + width + 0.03, top + height / 2, x + width + gap - 0.03, top + height / 2, colour=BLUE, lw=1.4)

    deck.panel(slide, ML, 3.62, CW, 0.94, accent=RED, wash=0.94)
    deck.pill(slide, ML + 0.18, 3.75, "FORBIDDEN FEEDBACK", fill=RED, w=1.78)
    deck.text(
        slide, ML + 2.12, 3.74, CW - 2.35, 0.30,
        "same-frame detections  •  confidence  •  current track identity  •  map quality  •  CARLA truth",
        size=10.5, bold=True, colour=NAVY,
    )
    deck.text(
        slide, ML + 0.20, 4.10, CW - 0.40, 0.34,
        "These outputs may influence the next decision only after their availability timestamp.",
        size=9, colour=GREY,
    )
    deck.panel(slide, ML + 2.75, 4.86, 6.00, 0.90, accent=TEAL, wash=0.93)
    deck.text(slide, ML + 2.75, 4.98, 6.00, 0.22, "CAUSAL INVARIANT", size=8, bold=True, colour=TEAL, align=PP_ALIGN.CENTER, spacing=70)
    deck.equation(slide, "causal_invariant", ML + 3.00, 5.24, 5.50, 0.38, colour=NAVY, max_height=0.34)
    deck.text(
        slide, ML, 6.06, CW, 0.30,
        "Two decisions, two timestamps: the placement action is chosen from the past; the publication action "
        "is chosen only once a result physically exists.",
        size=9, colour=GREY, align=PP_ALIGN.CENTER,
    )


def slide_04(deck: Deck) -> None:
    slide = deck.content_slide(
        4,
        "State",
        "State design: what is known before the action?",
        "Every observation field is timestamped and audited at its consuming decision.",
        source="Design correction: no post-action information may select the action that produced it",
        notes=NOTES[4],
    )
    deck.equation_panel(
        slide, "state", ML, BODY_T, CW, 1.12,
        label="PRE-ACTION OBSERVATION", accent=BLUE,
    )
    fields = [
        ("Network", "lagged capacity + confidence\nMCS / BLER / BSR • prior delivery", BLUE),
        ("Map", "installed tracks • AoI • covariance\nsource + capture/install provenance", TEAL),
        ("Motion", "helper-local pose • newest causally\nreceived recipient pose and age", ORANGE),
        ("Runtime", "scheduler credit • in-flight work\nLOCAL latency/headroom • prior action", DEEP),
    ]
    width = (CW - 3 * 0.24) / 4
    for index, (heading, body, accent) in enumerate(fields):
        deck.card(slide, ML + index * (width + 0.24), 3.26, width, 1.40, heading, body, accent=accent)

    deck.panel(slide, ML, 4.98, 5.60, 1.30, accent=TEAL, wash=0.93)
    deck.pill(slide, ML + 0.18, 5.12, "ALLOWED", fill=TEAL, w=1.00)
    deck.text(
        slide, ML + 0.20, 5.50, 5.20, 0.66,
        "Prior completed causal tracks, each carrying source, observed, and available timestamps.",
        size=10, bold=True, colour=NAVY, line_spacing=1.24,
    )
    deck.panel(slide, ML + 5.90, 4.98, 5.60, 1.30, accent=RED, wash=0.94)
    deck.pill(slide, ML + 6.08, 5.12, "FORBIDDEN", fill=RED, w=1.15)
    deck.text(
        slide, ML + 6.10, 5.50, 5.20, 0.66,
        "Current-action outputs • shadow outputs • ground-truth actor IDs • future hazard labels.",
        size=10, bold=True, colour=NAVY, line_spacing=1.24,
    )


def slide_05(deck: Deck) -> None:
    slide = deck.content_slide(
        5,
        "Actions",
        "Action space: discrete semantics, measured operating points",
        "Continuous state and outcomes do not justify inventing unmeasured continuous actions.",
        source="Measured profiles are categorical; payload/accuracy interpolation is not assumed",
        notes=NOTES[5],
    )
    placement = [
        ("SPLIT_FEATURE", "head at helper\nfeature over uplink\nedge-tail inference", DEEP),
        ("LOCAL_INFER", "full model at helper\ncompact object record\nlate/object fusion", TEAL),
        ("SKIP_INFERENCE", "no new inference\nprior map propagates\nuncertainty grows", RED),
    ]
    width = 2.24
    for index, (name, body, accent) in enumerate(placement):
        x = ML + index * (width + 0.20)
        deck.rect(slide, x, BODY_T, width, 1.44, fill=tint(accent, 0.94), edge=mix(accent, WHITE, 0.45), rounded=True, radius=0.09)
        deck.rect(slide, x, BODY_T, width, 0.28, fill=accent, rounded=True, radius=0.09)
        deck.rect(slide, x, BODY_T + 0.14, width, 0.14, fill=accent)
        deck.text(slide, x, BODY_T + 0.05, width, 0.22, name, size=8.5, bold=True, colour=WHITE, align=PP_ALIGN.CENTER, spacing=40)
        deck.text(slide, x + 0.12, BODY_T + 0.38, width - 0.24, 1.00, body, size=8.6, colour=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.28)

    deck.text(slide, ML, 3.42, 4.0, 0.24, "AFTER A RESULT EXISTS", size=8, bold=True, colour=ORANGE, spacing=70)
    for index, name in enumerate(("PUBLISH_ALL", "PUBLISH_HAZARD_SUBSET", "SKIP_PUBLICATION")):
        accent = RED if name.startswith("SKIP") else ORANGE
        x = ML + index * (width + 0.20)
        deck.rect(slide, x, 3.70, width, 0.42, fill=tint(accent, 0.90), edge=mix(accent, WHITE, 0.35), rounded=True, radius=0.08)
        deck.text(slide, x, 3.80, width, 0.24, name, size=8.5, bold=True, colour=NAVY if accent == ORANGE else RED, align=PP_ALIGN.CENTER)

    rows = [
        ("Discrete now", "placement • profile • publication", BLUE),
        ("Bounded grid now", "target FPS / update interval", ORANGE),
        ("Continuous later", "FPS only after held-out interpolation is validated", TEAL),
        ("Continuous outcomes", "latency • AoI • covariance • clearance", GREY),
    ]
    for index, (key, value, accent) in enumerate(rows):
        deck.kv_row(slide, ML, 4.42 + index * 0.27, 7.10, key, value, accent=accent, key_w=1.95, size=9)

    deck.panel(slide, ML, 5.42, 7.10, 0.86, accent=ORANGE, wash=0.95)
    deck.text(
        slide, ML + 0.20, 5.55, 6.70, 0.62,
        "A profile is a measured operating point, not a knob value: payload, accuracy, and latency are "
        "known only at the points that were actually captured.",
        size=9, bold=True, colour=NAVY, line_spacing=1.24,
    )

    deck.picture(slide, ASSETS["frontier"], ML + 7.42, BODY_T, 4.08, 1.86, mode="fit")
    deck.text(slide, ML + 7.42, 3.80, 4.08, 0.22, "MEASURED ACCURACY / PAYLOAD FRONTIER", size=7.5, bold=True, colour=BLUE, align=PP_ALIGN.CENTER, spacing=40)

    deck.panel(slide, ML + 7.42, 4.20, 4.08, 2.08, accent=DEEP, wash=0.94)
    deck.pill(slide, ML + 7.58, 4.34, "ALGORITHM IMPLICATION", fill=DEEP, w=1.95)
    deck.text(
        slide, ML + 7.60, 4.72, 3.74, 1.42,
        "Start with exact / rule / MPC over a masked finite catalog.\n"
        "If learning is needed: DQN, discrete SAC, or masked PPO — not continuous SAC.\n"
        "Masks are applied before scoring, so an unsupported profile is never scored at all.",
        size=9, bold=True, colour=NAVY, line_spacing=1.26,
    )


def slide_06(deck: Deck) -> None:
    slide = deck.content_slide(
        6,
        "Scenario loop",
        "Worked loop: one midblock pedestrian decision",
        "Illustrative controller step using the reviewed scenario — not a measured policy result.",
        source="Scenario: parked-van midblock pedestrian • helper/recipient synchronized Epic views",
        notes=NOTES[6],
    )
    deck.picture(slide, ASSETS["worked"], ML, BODY_T, 5.35, 3.55, mode="cover", segment="left", crop_top=REVIEW_BANNER_FRACTION)
    deck.rect(slide, ML, BODY_T, 5.35, 0.28, fill=NAVY)
    deck.text(
        slide, ML + 0.12, BODY_T + 0.055, 5.10, 0.22,
        "HELPER  |  controlled positive occlusion  |  Epic • 1280×720 • FOV 120°",
        size=7.5, bold=True, colour=WHITE,
    )
    deck.ellipse(slide, ML + 2.97, BODY_T + 1.48, 0.46, 0.52, edge=ORANGE, lw=2.25)
    deck.pill(slide, ML + 0.14, BODY_T + 3.10, "PEDESTRIAN EMERGES BESIDE VAN", fill=ORANGE, colour=NAVY, w=2.55)

    steps = [
        ("1  STATE", "mild channel estimate • recipient map aging • LOCAL headroom available", BLUE),
        ("2  ACTION", "SPLIT profile + FPS, or LOCAL + publication, or SKIP", DEEP),
        ("3  OUTCOME", "delivery / latency / bytes • new causal objects • warning or no warning", ORANGE),
        ("4  NEXT STATE", "queue and estimator update • map age / covariance update • outcome recorded", TEAL),
    ]
    x = ML + 5.65
    width = CW - 5.65
    for index, (name, body, accent) in enumerate(steps):
        y = BODY_T + index * 0.92
        deck.card(slide, x, y, width, 0.74, name, body, accent=accent, heading_size=9.5, body_size=8.6)
        if index < len(steps) - 1:
            deck.arrow(slide, x + 0.42, y + 0.76, x + 0.42, y + 0.90, colour=accent, lw=1.3)
    deck.text(
        slide, x, BODY_T + 3.72, width, 0.26,
        "↻  the outcome becomes lagged state at the next decision boundary",
        size=9, bold=True, colour=TEAL,
    )
    deck.text(
        slide, ML, 5.62, 5.35, 0.60,
        "The same scene can justify SPLIT, LOCAL, or SKIP. Which one is correct depends on state "
        "the agent already holds — never on the result it has not produced yet.",
        size=9, colour=GREY, line_spacing=1.24,
    )


def slide_07(deck: Deck) -> None:
    slide = deck.content_slide(
        7,
        "Transition model",
        "The action changes network, map, and hazard state",
        "This temporal coupling — not a large state vector — is the reason learning may eventually help.",
        source="Phase-2 propagates state + covariance; it does not reuse the frozen-object speed × AoI shortcut",
        notes=NOTES[7],
    )
    equations = [
        ("queue", "NETWORK QUEUE", DEEP),
        ("object", "OBJECT STATE", TEAL),
        ("covariance", "UNCERTAINTY", ORANGE),
    ]
    width = (CW - 2 * 0.26) / 3
    for index, (key, label, accent) in enumerate(equations):
        deck.equation_panel(
            slide, key, ML + index * (width + 0.26), BODY_T, width, 1.14,
            label=label, accent=accent, pad=0.24,
        )
    branches = [
        ("DELIVERED + INSTALLED", "sequence accepted\nAoI resets to capture age\ncovariance starts from new evidence", TEAL),
        ("DROPPED / EXPIRED", "queue service consumed\nno map installation\nold state keeps propagating", RED),
        ("SKIP", "no new work created\ncorrect if no service demand\nrisky if hazard debt accumulates", ORANGE),
    ]
    for index, (heading, body, accent) in enumerate(branches):
        deck.card(slide, ML + index * (width + 0.26), 3.24, width, 1.42, heading, body, accent=accent, heading_size=9.5)

    deck.equation_panel(
        slide, "deadline", ML + 0.55, 4.90, CW - 1.10, 1.10,
        label="ACTIONABLE DEADLINE", accent=RED,
    )
    deck.text(
        slide, ML, 6.14, CW, 0.28,
        "The deadline terms are frozen only once the common warning-to-braking adapter exists.",
        size=9, colour=GREY, align=PP_ALIGN.CENTER,
    )


def slide_08(deck: Deck) -> None:
    slide = deck.content_slide(
        8,
        "Reward",
        "Reward v5: task value inside a constrained safe set",
        "Causality, feasibility, and safety are masks — not quantities the reward may buy its way around.",
        source="Source: REWARD_FORMULATION.md v5 and the Phase-2 constraint catalog",
        notes=NOTES[8],
    )
    deck.equation_panel(
        slide, "utility", ML, BODY_T, CW, 1.32,
        label="PERCEPTION UTILITY", accent=TEAL,
    )
    bars = [
        ("Pedestrian recall", 0.40, RED, "highest supported safety-critical class weight"),
        ("Segmentation", 0.35, BLUE, "protects scene understanding from ROI collapse"),
        ("Vehicle recall", 0.25, ORANGE, "explicit rather than hidden inside object recall"),
    ]
    for index, (label, value, colour, note) in enumerate(bars):
        y = 3.42 + index * 0.44
        deck.text(slide, ML, y, 1.95, 0.24, label, size=9.5, bold=True, colour=NAVY)
        deck.bar(slide, ML + 2.00, y + 0.035, 3.30, 0.19, value / 0.40, colour=colour)
        deck.text(slide, ML + 5.42, y, 0.55, 0.24, f"{int(value * 100)}%", size=9.5, bold=True, colour=colour)
        deck.text(slide, ML + 6.10, y, CW - 6.10, 0.24, note, size=9, colour=GREY)

    deck.equation_panel(
        slide, "reward", ML, 4.86, CW, 1.16,
        label="CANDIDATE PHASE-2 INNER SCORE", accent=BLUE,
    )
    deck.text(slide, ML, 6.16, 1.95, 0.24, "No explicit ROI cost", size=9, bold=True, colour=ORANGE)
    deck.text(
        slide, ML + 1.98, 6.16, 4.20, 0.24,
        "ROI damage already shows up in task utility.",
        size=9, colour=GREY,
    )
    deck.text(slide, ML + 6.35, 6.16, 2.05, 0.24, "Weights unfrozen", size=9, bold=True, colour=RED)
    deck.text(
        slide, ML + 8.10, 6.16, CW - 8.10, 0.24,
        "until causal LOCAL / OAI tables exist.",
        size=9, colour=GREY,
    )


def slide_09(deck: Deck) -> None:
    slide = deck.content_slide(
        9,
        "Reward effects",
        "Reward terms shape behavior — without replacing constraints",
        "The desired controller is useful, efficient, stable, and honest about graceful degradation.",
        source="Stopping too early and stopping too close are both undesirable; clearance alone is gameable",
        notes=NOTES[9],
    )
    effects = [
        ("Task utility ↑", "Preserve segmentation and class recall; avoid profiles whose compression destroys useful evidence.", TEAL),
        ("PRB-time cost ↑", "Prefer compact measured profiles or LOCAL when the same task value can be delivered more cheaply.", DEEP),
        ("Hazard debt ↑", "Publish before the recipient deadline; SKIP becomes costly only when service is actually owed.", RED),
        ("Switching cost ↑", "Avoid oscillating SPLIT ↔ LOCAL on noisy estimates; keep hysteresis without hiding deadline misses.", NAVY),
        ("Expected uncertainty ↑", "Prefer margin inside the admitted set; tail uncertainty remains a shield/service constraint.", ORANGE),
    ]
    left_w = 6.80
    for index, (name, body, accent) in enumerate(effects):
        y = BODY_T + index * 0.78
        deck.rect(slide, ML, y, left_w, 0.66, fill=WHITE, edge=LINE, rounded=True, radius=0.08)
        deck.rect(slide, ML, y + 0.09, 0.065, 0.48, fill=accent)
        deck.text(slide, ML + 0.20, y + 0.09, 1.95, 0.22, name, size=9, bold=True, colour=accent)
        deck.text(slide, ML + 2.10, y + 0.08, left_w - 2.30, 0.52, body, size=8.4, colour=GREY, line_spacing=1.18)

    right_x = ML + 7.10
    right_w = CW - 7.10
    deck.panel(slide, right_x, BODY_T, right_w, 1.62, accent=RED, wash=0.94)
    deck.pill(slide, right_x + 0.16, BODY_T + 0.14, "NO GLOBAL SKIP PENALTY", fill=RED, w=1.95)
    deck.text(
        slide, right_x + 0.18, BODY_T + 0.50, right_w - 0.36, 1.00,
        "Correct abstention in empty or fresh scenes must stay free.\n"
        "Charge the consequence instead: a missed warning, deadline debt, or growing map uncertainty.",
        size=9, bold=True, colour=NAVY, line_spacing=1.24,
    )
    deck.equation_panel(
        slide, "stop_cost", right_x, 3.62, right_w, 1.42,
        label="FUTURE STOP COST", accent=ORANGE, pad=0.22,
    )
    deck.text(
        slide, right_x + 0.02, 5.14, right_w, 0.28,
        "Enabled only once every arm uses the same warning-to-braking adapter.",
        size=8.4, colour=GREY,
    )
    deck.panel(slide, right_x, 5.50, right_w, 0.60, accent=TEAL, wash=0.92)
    deck.text(
        slide, right_x + 0.16, 5.63, right_w - 0.32, 0.36,
        "Collision and minimum clearance remain hard constraints.",
        size=9, bold=True, colour=DEEP, align=PP_ALIGN.CENTER,
    )


def slide_10(deck: Deck) -> None:
    slide = deck.content_slide(
        10,
        "Constraint rank",
        "Constraints are lexicographic, not one weighted soup",
        "A scalar reward ranks only the actions that survive the higher-priority contracts.",
        source="Operational IDs S/N/K/M/P/O stay separate from paper contributions C1–C4",
        notes=NOTES[10],
    )
    levels = [
        ("0", "CAUSAL + STRUCTURAL", "S0–S2 • no leakage • measured actions • recipient isolation", RED, 5.40),
        ("1", "PHYSICAL SAFETY", "P1–P2 • collision + minimum surface clearance", ORANGE, 6.62),
        ("2", "DEADLINE / SERVICE", "N2 + M2 • hazard evidence before the recipient deadline", NAVY, 7.84),
        ("3", "NETWORK + COMPUTE FEASIBILITY", "N1/N3 + K1/K2 + M1 • capacity, queue, LOCAL headroom, valid contribution", DEEP, 8.74),
        ("4", "TASK UTILITY", "M3 • segmentation + pedestrian recall + vehicle recall", TEAL, 10.12),
        ("5", "EFFICIENCY + COMFORT", "PRB-time • bytes • compute • switching • stop placement • jerk • progress", BLUE, 11.50),
    ]
    centre = ML + CW / 2
    for index, (rank, name, body, accent, width) in enumerate(levels):
        y = BODY_T + index * 0.70
        x = centre - width / 2
        deck.rect(slide, x, y, width, 0.60, fill=tint(accent, 0.94), edge=mix(accent, WHITE, 0.45), rounded=True, radius=0.08)
        deck.rect(slide, x, y, 0.90, 0.60, fill=accent, rounded=True, radius=0.08)
        deck.rect(slide, x + 0.45, y, 0.45, 0.60, fill=accent)
        deck.text(slide, x, y + 0.19, 0.90, 0.24, f"RANK {rank}", size=8, bold=True, colour=WHITE if accent != ORANGE else NAVY, align=PP_ALIGN.CENTER)
        deck.text(slide, x + 1.05, y + 0.08, width - 1.20, 0.24, name, size=9.5, bold=True, colour=NAVY)
        deck.text(slide, x + 1.05, y + 0.33, width - 1.20, 0.24, body, size=8.2, colour=GREY)
    deck.arrow(slide, ML - 0.30, BODY_T + 0.10, ML - 0.30, BODY_T + 4.00, colour=mix(GREY, WHITE, 0.4), lw=1.2)
    deck.text(
        slide, 0.10, BODY_T + 1.55, 0.46, 0.90,
        "higher\npriority\nfirst",
        size=8, bold=True, colour=mix(GREY, WHITE, 0.25), align=PP_ALIGN.RIGHT,
    )
    deck.panel(slide, ML + 0.60, 6.14, CW - 1.20, 0.52, accent=RED, wash=0.94)
    deck.text(
        slide, ML + 0.70, 6.26, CW - 1.40, 0.30,
        "No reward weight can authorize causal leakage, an unsupported action, or a hard safety violation.",
        size=9.5, bold=True, colour=NAVY, align=PP_ALIGN.CENTER,
    )


def slide_11(deck: Deck) -> None:
    slide = deck.content_slide(
        11,
        "Relationships",
        "The constraint domains are coupled",
        "Network and physical context meet through map installation time, uncertainty, and the hazard deadline.",
        source="Graceful degradation is explicit when no action satisfies every service target",
        notes=NOTES[11],
    )
    node_w, node_h = 2.42, 1.00
    columns = [ML, ML + 3.05, ML + 6.10]
    rows = [BODY_T + 0.10, BODY_T + 1.90]
    nodes = [
        (columns[0], rows[0], "PHYSICAL ENVIRONMENT", "object motion • occlusion\nclosing speed • recipient path", ORANGE),
        (columns[0], rows[1], "NETWORK + COMPUTE", "capacity estimate • queue\nLOCAL p95 • in-flight work", DEEP),
        (columns[1], rows[0], "HAZARD DEADLINE", "required response margin\nservice debt", RED),
        (columns[1], rows[1], "ACTION ADMISSION", "measured catalog ∩ radio\n∩ compute ∩ service", BLUE),
        (columns[2], rows[0], "MAP STATE", "install time • AoI\nstate + covariance", TEAL),
        (columns[2], rows[1], "REWARD RANKING", "task value − costs\n− debt − switching", NAVY),
    ]
    for x, y, name, body, accent in nodes:
        deck.rect(slide, x, y, node_w, node_h, fill=tint(accent, 0.94), edge=mix(accent, WHITE, 0.42), rounded=True, radius=0.08)
        deck.rect(slide, x, y + 0.08, 0.065, node_h - 0.16, fill=accent)
        deck.text(slide, x + 0.20, y + 0.12, node_w - 0.36, 0.22, name, size=8.5, bold=True, colour=accent, spacing=30)
        deck.text(slide, x + 0.20, y + 0.40, node_w - 0.36, 0.50, body, size=8.4, colour=NAVY, line_spacing=1.22)

    action_x = ML + 9.30
    action_y = BODY_T + 0.90
    deck.rect(slide, action_x, action_y, 2.20, 1.20, fill=NAVY, rounded=True, radius=0.10)
    deck.text(slide, action_x, action_y + 0.16, 2.20, 0.24, "ACTION", size=9, bold=True, colour=WHITE, align=PP_ALIGN.CENTER, spacing=60)
    deck.text(
        slide, action_x, action_y + 0.46, 2.20, 0.62,
        "SPLIT  •  LOCAL  •  SKIP",
        size=9, bold=True, colour=BLUE, align=PP_ALIGN.CENTER,
    )

    for row_index, y in enumerate(rows):
        mid = y + node_h / 2
        colour = (ORANGE, DEEP)[row_index]
        deck.arrow(slide, columns[0] + node_w + 0.04, mid, columns[1] - 0.04, mid, colour=colour, lw=1.4)
        colour = (RED, BLUE)[row_index]
        deck.arrow(slide, columns[1] + node_w + 0.04, mid, columns[2] - 0.04, mid, colour=colour, lw=1.4)
    deck.arrow(slide, columns[2] + node_w / 2, rows[0] + node_h + 0.04, columns[2] + node_w / 2, rows[1] - 0.04, colour=TEAL, lw=1.4)
    deck.arrow(slide, columns[2] + node_w + 0.04, rows[1] + node_h / 2, action_x - 0.04, action_y + 0.85, colour=NAVY, lw=1.4)
    deck.arrow(slide, columns[2] + node_w + 0.04, rows[0] + node_h / 2, action_x - 0.04, action_y + 0.35, colour=TEAL, lw=1.4)

    deck.arrow(slide, action_x + 1.10, 5.24, ML + 0.02, 5.24, colour=TEAL, lw=1.4)
    deck.text(
        slide, ML, 4.94, CW, 0.26,
        "the outcome becomes lagged state at the next decision",
        size=8.5, bold=True, colour=TEAL, align=PP_ALIGN.CENTER,
    )
    deck.panel(slide, ML + 0.40, 5.56, CW - 0.80, 0.58, accent=TEAL, wash=0.93)
    deck.text(
        slide, ML + 0.50, 5.69, CW - 1.00, 0.32,
        "Network decides feasible fusion  •  the observed hazard decides the deadline  •  reward ranks the survivors",
        size=9.5, bold=True, colour=DEEP, align=PP_ALIGN.CENTER,
    )


def slide_12(deck: Deck) -> None:
    slide = deck.content_slide(
        12,
        "Physical world",
        "Physical and environmental constraints",
        "The agent must be evaluated in traffic that is legal, realistic, diverse, and attributable.",
        source="Primary CARLA capture: Town10HD_Opt • Epic • native 10 Hz • matched positive/benign futures",
        notes=NOTES[12],
    )
    deck.picture(slide, ASSETS["signalized"], ML, BODY_T, 5.62, 2.32, mode="cover", crop_top=REVIEW_BANNER_FRACTION)
    deck.picture(slide, ASSETS["queue"], ML + 5.88, BODY_T, 5.62, 2.32, mode="cover", crop_top=REVIEW_BANNER_FRACTION)
    deck.pill(slide, ML + 0.14, BODY_T + 1.94, "PEDESTRIAN HAZARD", fill=ORANGE, colour=NAVY, w=1.70)
    deck.pill(slide, ML + 6.02, BODY_T + 1.94, "VEHICLE REVEAL", fill=DEEP, w=1.45)

    constraints = [
        ("P1", "Collision", "report now; hard once warning actuation exists", RED),
        ("P2", "Surface clearance", "oriented ego/hazard boxes; advisor freezes the minimum", ORANGE),
        ("P3–P5", "Stopping + comfort", "too-close / too-early • deceleration • jerk", NAVY),
        ("O1", "Scene validity", "legal lanes • no gridlock • zero actor leaks", TEAL),
    ]
    width = (CW - 0.26) / 2
    for index, (identifier, name, body, accent) in enumerate(constraints):
        x = ML + (index % 2) * (width + 0.26)
        y = 4.44 + (index // 2) * 0.66
        deck.rect(slide, x, y, width, 0.54, fill=WHITE, edge=LINE, rounded=True, radius=0.08)
        deck.pill(slide, x + 0.12, y + 0.145, identifier, fill=accent, colour=WHITE if accent != ORANGE else NAVY, w=0.66)
        deck.text(slide, x + 0.86, y + 0.16, 1.55, 0.24, name, size=9, bold=True, colour=NAVY)
        deck.text(slide, x + 2.46, y + 0.17, width - 2.58, 0.24, body, size=8.2, colour=GREY)
    deck.text(
        slide, ML, 5.90, CW, 0.30,
        "Everything above is report-only today. Each becomes a hard constraint only after every arm shares "
        "one warning-to-braking adapter, so a violation is attributable to the controller.",
        size=9, colour=GREY, line_spacing=1.22,
    )


def slide_13(deck: Deck) -> None:
    slide = deck.content_slide(
        13,
        "Network + compute",
        "Network and compute constraints decide what is feasible",
        "The measured uplink has a payload-dependent delivery cliff; LOCAL is bounded by its own p95 latency.",
        source="Measured OAI RFsim channel sweep • delivery, queue latency, and capacity are action-dependent",
        notes=NOTES[13],
    )
    deck.rect(slide, ML, BODY_T, 6.55, 4.12, fill=WHITE, edge=LINE, rounded=True, radius=0.08)
    deck.picture(slide, ASSETS["delivery"], ML + 0.10, BODY_T + 0.10, 6.35, 3.92, mode="fit", edge=None)
    right_x = ML + 6.85
    right_w = CW - 6.85
    deck.equation_panel(
        slide, "radio_admission", right_x, BODY_T, right_w, 1.16,
        label="N1  RADIO ADMISSION", accent=DEEP, pad=0.20,
    )
    deck.equation_panel(
        slide, "local_admission", right_x, 3.16, right_w, 1.16,
        label="K1  LOCAL ADMISSION", accent=TEAL, pad=0.20,
    )
    cards = [
        ("Large feature", "high task ceiling • queue and latency risk", RED),
        ("Compact feature", "measured quality / cost compromise", ORANGE),
        ("LOCAL compact", "radio-light • compute-heavy • calibration pending", TEAL),
    ]
    for index, (name, body, accent) in enumerate(cards):
        y = 4.54 + index * 0.62
        deck.rect(slide, right_x, y, right_w, 0.52, fill=WHITE, edge=LINE, rounded=True, radius=0.08)
        deck.rect(slide, right_x, y + 0.08, 0.065, 0.36, fill=accent)
        deck.text(slide, right_x + 0.18, y + 0.06, right_w - 0.34, 0.22, name, size=9, bold=True, colour=accent)
        deck.text(slide, right_x + 0.18, y + 0.27, right_w - 0.34, 0.22, body, size=8.2, colour=GREY)


def slide_14(deck: Deck) -> None:
    slide = deck.content_slide(
        14,
        "Environment",
        "Building the agent environment",
        "An offline/replay Gym-style environment preserves causal timing and keeps training off the CARLA data path.",
        source="Environment contract: explicit observation/action spaces, deterministic reset, structured episode logs",
        notes=NOTES[14],
    )
    api = [
        ("reset(group_id, arm)", "load one immutable trajectory\nclear isolated map / queue / controller state", BLUE),
        ("observe()", "emit pre-action allowlisted state\nwith provenance and availability", DEEP),
        ("step(a_place)", "advance selected inference / transport model\nproduce a result or no-result", NAVY),
        ("step(a_publish)", "install contribution or retain the old map\ncompute reward + next state", ORANGE),
        ("close()", "flush structured logs\nrelease replay handles deterministically", TEAL),
    ]
    width = (CW - 4 * 0.22) / 5
    for index, (name, body, accent) in enumerate(api):
        x = ML + index * (width + 0.22)
        deck.rect(slide, x, BODY_T, width, 1.52, fill=tint(accent, 0.94), edge=mix(accent, WHITE, 0.42), rounded=True, radius=0.09)
        deck.text(slide, x + 0.10, BODY_T + 0.14, width - 0.20, 0.36, name, size=8.5, bold=True, colour=accent, align=PP_ALIGN.CENTER, font="Consolas")
        deck.text(slide, x + 0.12, BODY_T + 0.58, width - 0.24, 0.84, body, size=8.2, colour=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.24)
        if index < len(api) - 1:
            deck.arrow(slide, x + width + 0.02, BODY_T + 0.76, x + width + 0.20, BODY_T + 0.76, colour=accent, lw=1.3)

    left_w = (CW - 0.30) / 2
    deck.panel(slide, ML, 3.66, left_w, 1.86, accent=BLUE, wash=0.95)
    deck.pill(slide, ML + 0.16, 3.80, "CLOCKS + DATA", fill=BLUE, w=1.35)
    deck.text(
        slide, ML + 0.20, 4.18, left_w - 0.40, 1.24,
        "•  CARLA sensors and detection: native 10 Hz\n"
        "•  optional surrogate policy clock: 20 Hz, interpolated only\n"
        "•  paired arms replay identical captured evidence\n"
        "•  the truth stream is joined only by the evaluator",
        size=9, colour=NAVY, line_spacing=1.34,
    )
    right_x = ML + left_w + 0.30
    deck.panel(slide, right_x, 3.66, left_w, 1.86, accent=TEAL, wash=0.95)
    deck.pill(slide, right_x + 0.16, 3.80, "STRUCTURED OUTPUT", fill=TEAL, w=1.70)
    deck.text(
        slide, right_x + 0.20, 4.18, left_w - 0.40, 1.24,
        "•  episode_id • trajectory_group • frame / time\n"
        "•  state provenance • masks • both actions • outcome\n"
        "•  reward terms • deadline debt • bytes / PRB / compute\n"
        "•  warning + physical evaluation in separate namespaces",
        size=9, colour=NAVY, line_spacing=1.34,
    )
    deck.panel(slide, ML + 1.00, 5.72, CW - 2.00, 0.56, accent=RED, wash=0.94)
    deck.text(
        slide, ML + 1.10, 5.85, CW - 2.20, 0.32,
        "Training never blocks CARLA or the real-time sensing path — use replay buffers and offline/async updates.",
        size=9.5, bold=True, colour=NAVY, align=PP_ALIGN.CENTER,
    )


def slide_15(deck: Deck) -> None:
    slide = deck.content_slide(
        15,
        "Scenarios",
        "Suite A creates decision opportunities",
        "Six visually accepted hazard families span class, occlusion geometry, speed, deadline, and traffic density.",
        source="Suite A = designed decision opportunities • cyclists excluded until the perception contract supports them",
        notes=NOTES[15],
    )
    items = [
        (ASSETS["curbside"], "1", "Curbside bus • pedestrian", ORANGE),
        (ASSETS["signalized"], "2", "Signalized corner • pedestrian", ORANGE),
        (ASSETS["worked"], "3", "Midblock van • pedestrian", ORANGE),
        (ASSETS["cross_traffic"], "4", "Occluded cross-traffic • vehicle", DEEP),
        (ASSETS["pullout"], "5", "Parked pull-out • vehicle", DEEP),
        (ASSETS["queue"], "6", "Queue reveal • vehicle", DEEP),
    ]
    width = (CW - 2 * 0.24) / 3
    height = 1.62
    for index, (path, number, label, accent) in enumerate(items):
        row, column = divmod(index, 3)
        x = ML + column * (width + 0.24)
        y = BODY_T + row * (height + 0.52)
        deck.picture(slide, path, x, y, width, height, mode="cover", edge=mix(accent, WHITE, 0.4), crop_top=REVIEW_BANNER_FRACTION)
        deck.rect(slide, x, y + height, width, 0.30, fill=tint(accent, 0.92))
        deck.text(slide, x + 0.12, y + height + 0.055, 0.22, 0.22, number, size=8.5, bold=True, colour=accent)
        deck.text(slide, x + 0.36, y + height + 0.055, width - 0.48, 0.22, label, size=8.5, bold=True, colour=NAVY)
    deck.panel(slide, ML + 0.90, 6.10, CW - 1.80, 0.54, accent=TEAL, wash=0.93)
    deck.text(
        slide, ML + 1.00, 6.23, CW - 2.00, 0.30,
        "Each positive has a matched benign twin; traffic, weather, and seeds stay paired and group-locked.",
        size=9.5, bold=True, colour=DEEP, align=PP_ALIGN.CENTER,
    )


def slide_16(deck: Deck) -> None:
    slide = deck.content_slide(
        16,
        "Evaluation design",
        "Suite B keeps the result honest",
        "Naturalistic operation is reported with the same endpoints, so the curated suite cannot flatter the controller invisibly.",
        source="Pilot excluded • primary C2 endpoint: recipient-specific actionable warning success + warning lead",
        notes=NOTES[16],
    )
    width = (CW - 0.30) / 2
    deck.picture(slide, ASSETS["natural_signal"], ML, BODY_T, width, 2.10, mode="cover", edge=mix(TEAL, WHITE, 0.35), crop_top=REVIEW_BANNER_FRACTION)
    deck.picture(slide, ASSETS["natural_perimeter"], ML + width + 0.30, BODY_T, width, 2.10, mode="cover", edge=mix(TEAL, WHITE, 0.35), crop_top=REVIEW_BANNER_FRACTION)
    deck.pill(slide, ML + 0.14, BODY_T + 1.72, "SIGNALIZED DEMO ROUTE", fill=TEAL, w=1.95)
    deck.pill(slide, ML + width + 0.44, BODY_T + 1.72, "SAFE-PERIMETER ROUTE", fill=TEAL, w=1.95)

    cards = [
        ("Designed Suite A", "helper-visible hazards + matched benign negatives\nregime-bounded cooperation claim", ORANGE),
        ("Naturalistic Suite B", "hazards at natural prevalence\nno-hazard runs remain in the denominator", TEAL),
        ("Grouped inference", "trajectory / scenario clusters — not frames\npositive + benign are never split", BLUE),
    ]
    card_w = (CW - 2 * 0.24) / 3
    for index, (heading, body, accent) in enumerate(cards):
        deck.card(slide, ML + index * (card_w + 0.24), 4.20, card_w, 1.10, heading, body, accent=accent, heading_size=10)

    x0, y0, bar_w, bar_h = ML + 1.80, 5.66, CW - 3.60, 0.32
    deck.text(slide, ML + 0.30, y0 + 0.05, 1.40, 0.24, "SPLIT", size=8, bold=True, colour=GREY, align=PP_ALIGN.RIGHT, spacing=60)
    segments = [(0.20, BLUE, "20%  CALIBRATION"), (0.20, ORANGE, "20%  VALIDATION"), (0.60, TEAL, "60%  UNTOUCHED TEST")]
    offset = 0.0
    for fraction, colour, label in segments:
        deck.rect(slide, x0 + bar_w * offset, y0, bar_w * fraction, bar_h, fill=colour)
        deck.text(
            slide, x0 + bar_w * offset, y0 + 0.08, bar_w * fraction, 0.22, label,
            size=7.5, bold=True, colour=NAVY if colour is ORANGE else WHITE, align=PP_ALIGN.CENTER,
        )
        offset += fraction
    deck.text(
        slide, ML, 6.16, CW, 0.26,
        "Hashed, group-locked assignment across 210 independent groups and 330 world trajectories.",
        size=8.5, colour=GREY, align=PP_ALIGN.CENTER,
    )


def slide_17(deck: Deck) -> None:
    slide = deck.content_slide(
        17,
        "Progress",
        "What is already banked",
        "The pilot proves the causal capture and evaluation chain — not cooperation performance.",
        source="Pilot batch: 20260817_181354_pilot • authoritative evaluation/verification namespaces retained",
        notes=NOTES[17],
    )
    img_w = 3.36
    deck.picture(slide, ASSETS["pilot_helper"], ML, BODY_T, img_w, 1.72, mode="cover", edge=mix(BLUE, WHITE, 0.35))
    deck.picture(slide, ASSETS["pilot_recipient"], ML + img_w + 0.22, BODY_T, img_w, 1.72, mode="cover", edge=mix(TEAL, WHITE, 0.35))
    deck.pill(slide, ML + 0.12, BODY_T + 1.36, "HELPER • FRAME 156300", fill=BLUE, w=1.85)
    deck.pill(slide, ML + img_w + 0.34, BODY_T + 1.36, "RECIPIENT • SAME FRAME", fill=TEAL, w=1.95)

    panel_x = ML + 2 * (img_w + 0.22)
    panel_w = CW - 2 * (img_w + 0.22)
    deck.panel(slide, panel_x, BODY_T, panel_w, 1.72, accent=TEAL, wash=0.95)
    deck.pill(slide, panel_x + 0.16, BODY_T + 0.12, "ACCEPTED PILOT CHAIN", fill=TEAL, w=1.80)
    checks = [
        "aligned RGB + radar retained",
        "unfiltered detections + causal tracks",
        "placement and publication audit",
        "recipient map install + warning",
        "separate truth + future adjudication",
    ]
    deck.text(
        slide, panel_x + 0.20, BODY_T + 0.48, panel_w - 0.40, 1.16,
        "\n".join(f"•  {item}" for item in checks),
        size=8.6, colour=NAVY, line_spacing=1.30,
    )

    milestones = [
        ("Measured inputs", "profile payload / quality\nchannel + queue surface\nexact sensor contract", BLUE),
        ("Control design", "causal state and action\nreward v5\nranked constraints", DEEP),
        ("Map + evaluation", "v2 contribution schema\nwarning engine\nfuture-hazard adjudicator", NAVY),
        ("Scenario design", "6 designed families\n2 naturalistic routes\nEpic visual acceptance", ORANGE),
        ("Pilot", "helper → map → warning\nraw recoverability\nstructural gates PASS", TEAL),
    ]
    card_w = (CW - 4 * 0.22) / 5
    for index, (heading, body, accent) in enumerate(milestones):
        deck.card(slide, ML + index * (card_w + 0.22), 4.02, card_w, 1.20, heading, body, accent=accent, heading_size=9.5, body_size=8.2)
    deck.panel(slide, ML + 0.70, 5.56, CW - 1.40, 0.56, accent=NAVY, wash=0.95)
    deck.text(
        slide, ML + 0.80, 5.69, CW - 1.60, 0.32,
        "Current scientific gate: calibration sufficiency → frozen operating point → causal baseline ladder",
        size=9.5, bold=True, colour=NAVY, align=PP_ALIGN.CENTER,
    )


def slide_18(deck: Deck) -> None:
    slide = deck.content_slide(
        18,
        "Roadmap",
        "Next steps: answer whether learning is necessary",
        "The staged plan stops at every scientific gate; no long collection or RL training is automatic.",
        source="RL gate: a held-out causal service/reward gap that survives the exact, rule, and MPC baselines",
        notes=NOTES[18],
    )
    stages = [
        ("1", "CALIBRATION AUDIT", "prove raw replay sufficiency\nrun the bounded 96-setting replay", BLUE),
        ("2", "FREEZE + POWER", "choose one warning setting\nverify nuisance / miss / power gates", ORANGE),
        ("3", "VALIDATION → TEST", "grouped Suite A + B\nlocal C2, then identical bytes over OAI", TEAL),
        ("4", "LOCAL TABLE", "payload vs object count\nlocal p50/p95 + sustainable FPS", DEEP),
        ("5", "CONTROLLER LADDER", "exact → rule → λ-RDO/AoI → MPC\nRL only after a residual gap", NAVY),
        ("6", "PHYSICAL OVERRIDE", "fixed warning-to-braking adapter\nclearance • stop band • jerk", RED),
    ]
    width = (CW - 2 * 0.24) / 3
    for index, (number, name, body, accent) in enumerate(stages):
        row, column = divmod(index, 3)
        x = ML + column * (width + 0.24)
        y = BODY_T + row * 1.36
        deck.rect(slide, x, y, width, 1.16, fill=WHITE, edge=mix(accent, WHITE, 0.5), rounded=True, radius=0.09)
        deck.rect(slide, x, y + 0.10, 0.075, 0.96, fill=accent)
        deck.ellipse(slide, x + 0.22, y + 0.14, 0.30, 0.30, fill=accent, edge=None)
        deck.text(slide, x + 0.22, y + 0.195, 0.30, 0.22, number, size=8.5, bold=True, colour=WHITE if accent != ORANGE else NAVY, align=PP_ALIGN.CENTER)
        deck.text(slide, x + 0.62, y + 0.17, width - 0.78, 0.24, name, size=9, bold=True, colour=accent, spacing=25)
        deck.text(slide, x + 0.24, y + 0.54, width - 0.42, 0.54, body, size=8.4, colour=GREY, line_spacing=1.22)
        if column < 2:
            deck.arrow(slide, x + width + 0.03, y + 0.58, x + width + 0.21, y + 0.58, colour=accent, lw=1.3)

    deck.panel(slide, ML, 4.66, CW, 0.90, accent=ORANGE, wash=0.94)
    deck.pill(slide, ML + 0.16, 4.80, "ADVISOR DECISIONS TO FREEZE", fill=ORANGE, colour=NAVY, w=2.30)
    deck.text(
        slide, ML + 0.20, 5.16, CW - 0.40, 0.30,
        "minimum clearance + comfort stop band   •   reaction and braking assumptions   •   "
        "LOCAL hardware target   •   staged calibration authorization",
        size=9.5, bold=True, colour=NAVY,
    )
    deck.panel(slide, ML + 1.20, 5.78, CW - 2.40, 0.54, accent=TEAL, wash=0.92)
    deck.text(
        slide, ML + 1.30, 5.91, CW - 2.60, 0.30,
        "Success is a defensible controller and feasibility envelope — even if the simplest controller wins.",
        size=9.5, bold=True, colour=DEEP, align=PP_ALIGN.CENTER,
    )


BUILDERS = (
    slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
    slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
    slide_13, slide_14, slide_15, slide_16, slide_17, slide_18,
)


def write_manifest(deck: Deck) -> None:
    payload = {
        "schema": "scenesense.agent_progress_deck.v2",
        "deck": str(PPTX.relative_to(REPO)),
        "template": str(TEMPLATE.relative_to(REPO)),
        "template_sha256": sha256(TEMPLATE),
        "slide_count": len(deck.index),
        "page_size_in": [PAGE_W, PAGE_H],
        "math_rendering": {
            "engine": "matplotlib mathtext",
            "fontset": "cm",
            "dpi": DPI_EQ,
            "equations": {key: value for key, value in EQ.items()},
        },
        "brand": {
            "font": FONT,
            "navy": NAVY, "blue": BLUE, "deep": DEEP,
            "teal": TEAL, "orange": ORANGE, "grey": GREY, "functional_red": RED,
        },
        "slides": [
            {"number": item.number, "title": item.title, "source": item.source}
            for item in deck.index
        ],
        "source_assets": {
            name: {
                "path": str(path),
                "sha256": sha256(path),
                "bytes": path.stat().st_size,
            }
            for name, path in sorted(ASSETS.items())
        },
    }
    MANIFEST.write_text(json.dumps(payload, indent=2, sort_keys=True) + "\n", encoding="utf-8")


def export_previews() -> None:
    """Render PNG/PDF previews with LibreOffice when it is available."""

    soffice = shutil.which("soffice")
    if soffice is None:
        print("libreoffice not found; skipping PNG/PDF preview export")
        return
    preview = BUILD / "preview"
    preview.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(preview), str(PPTX)],
        check=True,
        capture_output=True,
        timeout=600,
    )
    produced = preview / (PPTX.stem + ".pdf")
    if produced.is_file():
        shutil.copy2(produced, OUT / (PPTX.stem + ".pdf"))
        print(f"wrote {OUT / (PPTX.stem + '.pdf')}")


def build() -> None:
    check_inputs()
    BUILD.mkdir(parents=True, exist_ok=True)
    deck = Deck()
    for builder in BUILDERS:
        builder(deck)
    if len(deck.prs.slides._sldIdLst) != len(BUILDERS):
        raise RuntimeError("slide count does not match the builder list")
    deck.save()
    write_manifest(deck)
    print(f"wrote {PPTX}  ({len(BUILDERS)} slides)")
    print(f"wrote {MANIFEST}")
    export_previews()


if __name__ == "__main__":
    build()
