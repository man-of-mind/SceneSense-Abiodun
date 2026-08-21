#!/usr/bin/env python3
"""Build the editable supervisor deck for the UE split-only baseline.

The deck is created on top of the InterDigital PowerPoint template. All
experiment content is native PowerPoint text, tables, shapes, connectors, and
charts; no slide is flattened into an image.
"""

from __future__ import annotations

import csv
import hashlib
import json
import zipfile
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


REPO = Path(__file__).resolve().parents[2]
OUT = REPO / "rl_agent" / "presentation" / "ue_split_supervisor_20260820"
TEMPLATE = (
    REPO
    / "rl_agent"
    / "presentation"
    / "agent_progress_20260819"
    / "interdigital_template.pptx"
)
PLAN = REPO / "rl_agent" / "UE_SPLIT_ONLY_SUPERVISOR_DELIVERABLE.md"
COMBINATIONS = REPO / "rl_agent" / "UE_SPLIT_ONLY_SUPERVISOR_COMBINATIONS_V1.csv"
CANDIDATES = (
    REPO
    / "rl_agent"
    / "experiments"
    / "ue_split_catalog_proposal_v1"
    / "20260820_042414_candidate"
    / "ue_split_candidate_catalog.csv"
)
DECODER_REPORT = (
    REPO
    / "rl_agent"
    / "experiments"
    / "model_precision_decoder_audit_v1"
    / "20260819_210004"
    / "REPORT.md"
)

PPTX = OUT / "SceneSense_UE_Split_Baseline_Supervisor_2026-08-20.pptx"
NOTES = OUT / "PRESENTER_NOTES.md"
MANIFEST = OUT / "SOURCE_MANIFEST.json"

DECK_DATE = "20 August 2026"
COPYRIGHT = "©2026 InterDigital, Inc. All Rights Reserved."
FONT = "Century Gothic"

# InterDigital template palette.
NAVY = "010644"
BLUE = "00A8EC"
DEEP = "0A5C93"
TEAL = "00C9C8"
ORANGE = "FFAA59"
GREY = "49515C"
LINE = "D5D9DF"
LIGHT = "E6E6E7"
RED = "D8455F"
WHITE = "FFFFFF"

PAGE_W = 13.3333
PAGE_H = 7.5
ML = 0.92
CONTENT_W = 11.50
RIGHT = ML + CONTENT_W
BODY_TOP = 1.84
SOURCE_TOP = 6.67


PROFILE_ORDER = ("rescue", "compact", "balanced", "quality")
PROFILE_IDS = {
    "rescue": "ae32__u4__q0.9__zstd3__ckpt10cebbeede4d",
    "compact": "ae32__u4__q0.5__zstd3__ckpt10cebbeede4d",
    "balanced": "ae64__u4__q0.5__zstd3__ckptc6a2362c7c2d",
    "quality": "ae128__u4__q0.3__zstd3__ckpt601984b96d85",
}
PROFILE_LABELS = {
    "rescue": "AE32 / u4 / q0.9",
    "compact": "AE32 / u4 / q0.5",
    "balanced": "AE64 / u4 / q0.5",
    "quality": "AE128 / u4 / q0.3",
}
PROFILE_ROLES = {
    "rescue": "Degraded rescue",
    "compact": "Compact normal",
    "balanced": "Balanced normal",
    "quality": "Quality normal",
}
PROFILE_ACCENTS = {
    "rescue": ORANGE,
    "compact": BLUE,
    "balanced": TEAL,
    "quality": DEEP,
}


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1 << 20), b""):
            digest.update(block)
    return digest.hexdigest()


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def mix(hex_a: str, hex_b: str, t: float) -> str:
    a = tuple(int(hex_a[i : i + 2], 16) for i in (0, 2, 4))
    b = tuple(int(hex_b[i : i + 2], 16) for i in (0, 2, 4))
    return "".join(f"{round(a[i] + (b[i] - a[i]) * t):02X}" for i in range(3))


def tint(accent: str, amount: float = 0.92) -> str:
    return mix(accent, WHITE, amount)


def read_rows(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as stream:
        return list(csv.DictReader(stream))


def check_inputs() -> None:
    required = (TEMPLATE, PLAN, COMBINATIONS, CANDIDATES, DECODER_REPORT)
    missing = [str(path) for path in required if not path.is_file()]
    if missing:
        raise FileNotFoundError("missing deck input(s):\n  " + "\n  ".join(missing))


def load_evidence() -> tuple[dict[str, dict[str, float | str]], list[dict[str, str]]]:
    combinations = read_rows(COMBINATIONS)
    if len(combinations) != 16:
        raise ValueError(f"expected 16 logical combination rows, found {len(combinations)}")
    if len({row["cell_id"] for row in combinations}) != 16:
        raise ValueError("combination cell_id values are not unique")
    if {row["network_regime"] for row in combinations} != {"clear", "mild", "mid", "poor"}:
        raise ValueError("network regimes do not match the locked four-regime contract")
    if any(row["measurement_authorized"].lower() != "false" for row in combinations):
        raise ValueError("the supervisor planning sheet must not authorize a measurement")
    if sum(row["new_measurement_plan"].startswith("INITIAL_DIRECT") for row in combinations) != 2:
        raise ValueError("the planning sheet must contain exactly two initial direct cells")

    catalog_by_id = {row["profile_id"]: row for row in read_rows(CANDIDATES)}
    profiles: dict[str, dict[str, float | str]] = {}
    for key in PROFILE_ORDER:
        profile_id = PROFILE_IDS[key]
        if profile_id not in catalog_by_id:
            raise ValueError(f"missing shortlisted profile in candidate evidence: {profile_id}")
        row = catalog_by_id[profile_id]
        profile_cells = [item for item in combinations if item["profile_id"] == profile_id]
        if len(profile_cells) != 4:
            raise ValueError(f"profile {profile_id} does not have exactly four regime rows")
        reference = profile_cells[0]
        profiles[key] = {
            "profile_id": profile_id,
            "label": PROFILE_LABELS[key],
            "role": PROFILE_ROLES[key],
            "p95_kib": float(reference["p95_payload_kib"]),
            "offer_mbps": float(reference["estimated_udp_ip_offer_mbps_at_10hz"]),
            "vehicle_recall": float(row["recall_vehicle"]),
            "pedestrian_recall": float(row["recall_pedestrian"]),
            "miou": float(row["miou"]),
            "vehicle_xy": float(row["xy_mae_vehicle_m"]),
            "pedestrian_xy": float(row["xy_mae_pedestrian_m"]),
        }
    return profiles, combinations


@dataclass(frozen=True)
class SlideRecord:
    number: int
    title: str
    evidence_status: str
    source: str


class Deck:
    LAYOUT_TITLE = 0
    LAYOUT_CONTENT = 6

    def __init__(self) -> None:
        self.prs = Presentation(str(TEMPLATE))
        self._drop_template_slides()
        self.records: list[SlideRecord] = []

    def _drop_template_slides(self) -> None:
        slide_ids = self.prs.slides._sldIdLst
        for element in list(slide_ids):
            self.prs.part.drop_rel(element.rId)
            slide_ids.remove(element)

    @staticmethod
    def _title_size(value: str) -> float:
        if len(value) <= 48:
            return 24.0
        if len(value) <= 62:
            return 21.0
        return 19.0

    @staticmethod
    def _set_notes(slide, value: str) -> None:
        slide.notes_slide.notes_text_frame.text = value.strip()

    def _fill_footer(self, slide, number: int | None) -> None:
        for placeholder in slide.placeholders:
            placeholder_type = str(placeholder.placeholder_format.type)
            if placeholder_type.startswith("FOOTER"):
                value = COPYRIGHT
            elif placeholder_type.startswith("SLIDE_NUMBER"):
                value = "" if number is None else str(number)
            else:
                continue
            placeholder.text_frame.text = value
            for paragraph in placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(8)
                    run.font.color.rgb = rgb(mix(GREY, WHITE, 0.30))

    def title_slide(self, title: str, subtitle: str, notes: str) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_TITLE])
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        head = placeholders[0]
        head.text_frame.text = title
        for paragraph in head.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT
                run.font.size = Pt(39)
                run.font.bold = True
                run.font.color.rgb = rgb(WHITE)
        sub = placeholders[1]
        sub.text_frame.clear()
        for index, line in enumerate(subtitle.split("\n")):
            paragraph = sub.text_frame.paragraphs[0] if index == 0 else sub.text_frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(7)
            run = paragraph.add_run()
            run.text = line
            run.font.name = FONT
            run.font.size = Pt(15 if index == 0 else 11.5)
            run.font.bold = index == 0
            run.font.color.rgb = rgb(BLUE if index == 0 else mix(LIGHT, NAVY, 0.18))
        self.rect(slide, 1.67, 5.82, 8.65, 0.54, fill=DEEP, edge=BLUE, rounded=True)
        self.text(
            slide,
            1.92,
            5.97,
            8.15,
            0.24,
            "Measure first  →  derive the simplest defensible rule  →  justify learning only if needed",
            size=10,
            bold=True,
            colour=WHITE,
            align=PP_ALIGN.CENTER,
        )
        self._fill_footer(slide, None)
        self._set_notes(slide, notes)
        self.records.append(SlideRecord(1, title, "PLAN", "Supervisor plan"))

    def content_slide(
        self,
        number: int,
        eyebrow: str,
        title: str,
        subtitle: str,
        *,
        source: str,
        evidence_status: str,
        notes: str,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_CONTENT])
        placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        head = placeholders[0]
        head.left = Inches(ML)
        head.top = Inches(0.60)
        head.width = Inches(CONTENT_W)
        head.height = Inches(0.62)
        head.text_frame.word_wrap = True
        head.text_frame.text = title
        for paragraph in head.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = FONT
                run.font.size = Pt(self._title_size(title))
                run.font.bold = True
                run.font.color.rgb = rgb(NAVY)

        self.rect(slide, ML, 0.345, 0.13, 0.13, fill=BLUE)
        self.text(slide, ML + 0.24, 0.28, 3.2, 0.25, eyebrow.upper(), size=8.5, bold=True, colour=BLUE)
        self.text(
            slide,
            RIGHT - 3.6,
            0.28,
            3.6,
            0.25,
            f"UE SPLIT BASELINE  •  {number:02d}",
            size=8.5,
            bold=True,
            colour=mix(GREY, WHITE, 0.35),
            align=PP_ALIGN.RIGHT,
        )
        if subtitle:
            self.text(slide, ML, 1.22, CONTENT_W, 0.30, subtitle, size=10.5, colour=GREY)
        self.hline(slide, ML, 1.66, CONTENT_W, colour=LINE)
        self.text(slide, ML, SOURCE_TOP, CONTENT_W - 1.3, 0.20, source, size=7.0, colour=mix(GREY, WHITE, 0.25))
        self.pill(
            slide,
            RIGHT - 1.25,
            SOURCE_TOP - 0.02,
            evidence_status,
            w=1.25,
            h=0.20,
            size=6.8,
            fill={
                "PLAN": DEEP,
                "OFFLINE": TEAL,
                "HISTORICAL": GREY,
                "PROJECTED": ORANGE,
                "PLANNED": BLUE,
            }.get(evidence_status, DEEP),
            colour=NAVY if evidence_status in {"OFFLINE", "PROJECTED"} else WHITE,
        )
        self._fill_footer(slide, number)
        self._set_notes(slide, notes)
        self.records.append(SlideRecord(number, title, evidence_status, source))
        return slide

    @staticmethod
    def _plain(shape) -> None:
        shape.shadow.inherit = False
        if shape.has_text_frame:
            shape.text_frame.text = ""

    def rect(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        fill: str | None = None,
        edge: str | None = None,
        lw: float = 0.75,
        rounded: bool = False,
        radius: float = 0.08,
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        if rounded:
            shape.adjustments[0] = min(0.5, radius / min(w, h))
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)
        if edge is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = rgb(edge)
            shape.line.width = Pt(lw)
        self._plain(shape)
        return shape

    def panel(self, slide, x: float, y: float, w: float, h: float, *, accent: str, wash: float = 0.94):
        return self.rect(
            slide,
            x,
            y,
            w,
            h,
            fill=tint(accent, wash),
            edge=mix(accent, WHITE, 0.50),
            rounded=True,
        )

    def text(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        value: str,
        *,
        size: float = 10,
        bold: bool = False,
        colour: str = NAVY,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
        line_spacing: float = 1.10,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = anchor
        for index, line in enumerate(str(value).split("\n")):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = align
            paragraph.line_spacing = line_spacing
            paragraph.space_after = Pt(2)
            run = paragraph.add_run()
            run.text = line
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(colour)
        return box

    def rich_text(self, slide, x: float, y: float, w: float, h: float, lines: list[tuple[str, str, bool]], *, size=9.0):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        for index, (value, colour, bold) in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.space_after = Pt(5)
            paragraph.line_spacing = 1.12
            run = paragraph.add_run()
            run.text = value
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(colour)
        return box

    def pill(
        self,
        slide,
        x: float,
        y: float,
        label: str,
        *,
        fill: str,
        colour: str = WHITE,
        w: float | None = None,
        h: float = 0.24,
        size: float = 8,
    ) -> float:
        width = w if w is not None else max(0.78, 0.081 * len(label) + 0.28)
        shape = self.rect(slide, x, y, width, h, fill=fill, rounded=True, radius=h / 2)
        frame = shape.text_frame
        frame.word_wrap = False
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = label
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = rgb(colour)
        return width

    def hline(self, slide, x: float, y: float, w: float, *, colour: str, lw: float = 0.75):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x),
            Inches(y),
            Inches(x + w),
            Inches(y),
        )
        line.line.color.rgb = rgb(colour)
        line.line.width = Pt(lw)
        return line

    def arrow(self, slide, x1: float, y1: float, x2: float, y2: float, *, colour: str = BLUE, lw: float = 1.5):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        line.line.color.rgb = rgb(colour)
        line.line.width = Pt(lw)
        properties = line.line._get_or_add_ln()
        end = properties.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        properties.append(end)
        return line

    def flow_card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        step: str,
        heading: str,
        body: str,
        *,
        accent: str,
    ) -> None:
        self.rect(slide, x, y, w, h, fill=WHITE, edge=mix(accent, WHITE, 0.42), rounded=True)
        self.rect(slide, x, y, w, 0.25, fill=accent, rounded=True)
        self.pill(slide, x + 0.13, y + 0.36, step, fill=accent, w=0.42, h=0.26, size=8)
        self.text(slide, x + 0.65, y + 0.35, w - 0.78, 0.28, heading, size=9.3, bold=True, colour=NAVY)
        self.text(slide, x + 0.16, y + 0.78, w - 0.32, h - 0.88, body, size=8.1, colour=GREY, align=PP_ALIGN.CENTER)

    def info_card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        heading: str,
        body: str,
        *,
        accent: str,
    ) -> None:
        self.panel(slide, x, y, w, h, accent=accent, wash=0.95)
        self.rect(slide, x, y + 0.12, 0.075, h - 0.24, fill=accent)
        self.text(slide, x + 0.23, y + 0.16, w - 0.40, 0.25, heading, size=9.8, bold=True, colour=NAVY)
        self.text(slide, x + 0.23, y + 0.53, w - 0.40, h - 0.65, body, size=8.3, colour=GREY, line_spacing=1.18)

    @staticmethod
    def _cell_border(cell, colour: str = LINE, width: str = "6350") -> None:
        tc_pr = cell._tc.get_or_add_tcPr()
        for edge_name in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            edge = tc_pr.find(qn(edge_name))
            if edge is None:
                edge = OxmlElement(edge_name)
                tc_pr.append(edge)
            edge.set("w", width)
            solid = edge.find(qn("a:solidFill"))
            if solid is None:
                solid = OxmlElement("a:solidFill")
                edge.append(solid)
            for child in list(solid):
                solid.remove(child)
            srgb = OxmlElement("a:srgbClr")
            srgb.set("val", colour)
            solid.append(srgb)

    def style_cell(
        self,
        cell,
        value: str,
        *,
        fill: str,
        colour: str,
        size: float,
        bold: bool = False,
        align=PP_ALIGN.CENTER,
    ) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill)
        cell.margin_left = cell.margin_right = Inches(0.05)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame = cell.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(colour)
        self._cell_border(cell)


SLIDE_NOTES = {
    1: """This is deliberately a measurement plan, not an RL result. The central question is whether the network conditions create a meaningful profile-selection problem at all. If they do not, the correct outcome is a simple rule rather than forcing a learning result.""",
    2: """Keep the scope narrow: one UE, split inference, and one object-map endpoint. We reuse the same sensor samples and hold compute fixed. SKIP, LOCAL, urgency, cooperation, occlusion reasoning, and RL are later stages and do not block this baseline.""",
    3: """These four actions are an experimental shortlist from 72 offline configurations. The rescue action is separate because its pedestrian recall is below the normal floor. Segmentation is still measured, but the current service prioritizes object class and world location.""",
    4: """The chart is the most important reality check. All four action loads are below the central historical poor-link capacity. Only Quality under Poor is close enough to the uncertainty band to justify a direct boundary measurement. We should not assume switching or RL is necessary.""",
    5: """The 16 cells are a planning surface, not 16 experiments. We initially measure Quality under Clear and Poor. We step down only if Poor fails or is borderline. Inferred feasibility never fills unmeasured latency, drop, map-update, or AoI values.""",
    6: """Every direct cell uses the same 10-Hz replay and fixed pipeline. Freshness starts from the source release time of the newest accepted map update, not from send or enqueue. We report multiple latency summaries and derive the acceptable AoI from the measured error trade-off.""",
    7: """The baseline has three scientifically useful outcomes: a fixed/greedy rule, a measured profile boundary, or evidence that temporal history matters. Only the third outcome motivates a sequential learned policy. Please confirm the scope, shortlist, regimes, two-cell start, and AoI approach.""",
}


def build_title(deck: Deck) -> None:
    deck.title_slide(
        "UE Split-Inference Baseline",
        "Does the network create a meaningful split-profile decision?\n"
        "Single UE  •  OAI uplink  •  edge object map  •  supervisor discussion\n"
        f"Abiodun  |  {DECK_DATE}",
        SLIDE_NOTES[1],
    )


def build_scope(deck: Deck) -> None:
    slide = deck.content_slide(
        2,
        "Scope",
        "Start with the smallest question the UE agent must answer",
        "One UE, one object map, SPLIT only; vary the bundle and network while everything else stays fixed.",
        source="Source: UE_SPLIT_ONLY_SUPERVISOR_DELIVERABLE.md §§1–2",
        evidence_status="PLAN",
        notes=SLIDE_NOTES[2],
    )
    cards = [
        ("1", "RETAINED SAMPLE", "same aligned\nRGB + radar input", BLUE),
        ("2", "SPLIT BUNDLE", "registered model +\nquantization + ROI", DEEP),
        ("3", "OAI UPLINK", "measured SNR / MCS\nqueue + delivery", ORANGE),
        ("4", "OBJECT MAP", "accepted update\nclass + world XY", TEAL),
    ]
    card_w = 2.55
    gap = 0.43
    for index, (step, heading, body, accent) in enumerate(cards):
        x = ML + index * (card_w + gap)
        deck.flow_card(slide, x, 2.00, card_w, 1.25, step, heading, body, accent=accent)
        if index < len(cards) - 1:
            deck.arrow(slide, x + card_w + 0.05, 2.62, x + card_w + gap - 0.06, 2.62, colour=BLUE)

    deck.info_card(
        slide,
        ML,
        3.65,
        5.56,
        2.08,
        "IN THIS BASELINE",
        "• same source sequence and fixed compute\n"
        "• four registered split actions\n"
        "• four calibrated OAI regimes\n"
        "• latency, delivery, map freshness, object quality",
        accent=BLUE,
    )
    deck.info_card(
        slide,
        ML + 5.94,
        3.65,
        5.56,
        2.08,
        "DEFERRED UNTIL THE TABLE IS UNDERSTOOD",
        "• SKIP and full LOCAL inference\n"
        "• radar-conditioned action selection\n"
        "• occlusion / cooperation / map sharing\n"
        "• dynamic channels, MPC, DQN, discrete SAC",
        accent=GREY,
    )
    deck.panel(slide, ML + 1.32, 5.95, CONTENT_W - 2.64, 0.48, accent=ORANGE, wash=0.92)
    deck.text(
        slide,
        ML + 1.45,
        6.07,
        CONTENT_W - 2.90,
        0.22,
        "Plan only — no CARLA, OAI, model, or policy-training run is authorized",
        size=9.2,
        bold=True,
        colour=NAVY,
        align=PP_ALIGN.CENTER,
    )


def build_actions(deck: Deck, profiles: dict[str, dict[str, float | str]]) -> None:
    slide = deck.content_slide(
        3,
        "Action space",
        "Three normal profiles plus one explicit rescue span the trade-off",
        "The shortlist comes from 72 offline configurations; the remaining profiles stay available if a missing boundary appears.",
        source="Source: candidate catalog; 2,162-frame offline evaluation; baseline decoder",
        evidence_status="OFFLINE",
        notes=SLIDE_NOTES[3],
    )
    deck.pill(slide, RIGHT - 2.30, 1.88, "OFFLINE MEASURED • 2,162 FRAMES", fill=TEAL, colour=NAVY, w=2.30, h=0.25, size=7.4)

    headers = ("Role", "Registered bundle", "P95 payload", "10-Hz load", "Veh. recall", "Ped. recall", "mIoU")
    table_shape = slide.shapes.add_table(5, len(headers), Inches(ML), Inches(2.22), Inches(CONTENT_W), Inches(2.52))
    table = table_shape.table
    widths = (1.58, 2.35, 1.35, 1.35, 1.40, 1.42, 1.05)
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    table.rows[0].height = Inches(0.43)
    for row_index in range(1, 5):
        table.rows[row_index].height = Inches(0.52)
    for column, value in enumerate(headers):
        deck.style_cell(table.cell(0, column), value, fill=NAVY, colour=WHITE, size=8.4, bold=True)
    for row_index, key in enumerate(PROFILE_ORDER, start=1):
        profile = profiles[key]
        accent = PROFILE_ACCENTS[key]
        values = (
            str(profile["role"]),
            str(profile["label"]),
            f"{profile['p95_kib']:.1f} KiB",
            f"{profile['offer_mbps']:.2f} Mbps",
            f"{profile['vehicle_recall']:.3f}",
            f"{profile['pedestrian_recall']:.3f}",
            f"{profile['miou']:.3f}",
        )
        for column, value in enumerate(values):
            fill = tint(accent, 0.88) if column == 0 else (tint(accent, 0.96) if row_index % 2 else WHITE)
            deck.style_cell(
                table.cell(row_index, column),
                value,
                fill=fill,
                colour=NAVY if column == 0 else GREY,
                size=8.4,
                bold=column in {0, 1},
                align=PP_ALIGN.LEFT if column in {0, 1} else PP_ALIGN.CENTER,
            )

    card_w = (CONTENT_W - 0.36) / 3
    deck.info_card(
        slide,
        ML,
        5.02,
        card_w,
        1.28,
        "OBJECT_MAP_V1 — PRIMARY",
        "vehicle / pedestrian class + confidence\nworld-XY actor-reference location\nsample identity; valid-empty ≠ missing",
        accent=BLUE,
    )
    deck.info_card(
        slide,
        ML + card_w + 0.18,
        5.02,
        card_w,
        1.28,
        "SEGMENTATION — SECONDARY",
        "class IoU and mIoU remain reported\nbut do not define the object-map service\nor become a live agent input",
        accent=TEAL,
    )
    deck.info_card(
        slide,
        ML + 2 * (card_w + 0.18),
        5.02,
        card_w,
        1.28,
        "RESCUE — NOT NORMAL SUCCESS",
        "pedestrian recall is below the normal floor\nuse only if no normal action is feasible\nrecord explicit service debt",
        accent=ORANGE,
    )


def build_network(deck: Deck, profiles: dict[str, dict[str, float | str]], combinations: list[dict[str, str]]) -> None:
    slide = deck.content_slide(
        4,
        "Network evidence",
        "The current regimes may not force profile switching",
        "Only Quality × Poor is near the historical uncertainty boundary; exact 10-Hz map outcomes remain unmeasured.",
        source="Source: historical OAI four-regime sweep + composed 10-Hz planning sheet",
        evidence_status="HISTORICAL",
        notes=SLIDE_NOTES[4],
    )
    regime_order = ("clear", "mild", "mid", "poor")
    capacities = []
    snr_mcs = []
    for regime in regime_order:
        row = next(item for item in combinations if item["network_regime"] == regime)
        capacities.append(float(row["capacity_reference_mbps"]))
        snr_mcs.append(f"{float(row['achieved_snr_db_reference']):.1f} dB / MCS {row['mcs_reference']}")
    quality_load = float(profiles["quality"]["offer_mbps"])

    chart_data = ChartData()
    chart_data.categories = ("Clear", "Mild", "Mid", "Poor")
    chart_data.add_series("Historical capacity reference", capacities)
    chart_data.add_series("Quality offered load @ 10 Hz", (quality_load,) * 4)
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(ML),
        Inches(2.00),
        Inches(7.25),
        Inches(4.22),
        chart_data,
    )
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = FONT
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = rgb(NAVY)
    chart.font.name = FONT
    chart.font.size = Pt(9)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 40
    chart.value_axis.major_unit = 10
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.name = FONT
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.name = FONT
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.plots[0].gap_width = 65

    series_colours = (NAVY, BLUE)
    for series_index, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = rgb(series_colours[series_index])
        series.format.line.fill.background()
        series.has_data_labels = True
        labels = series.data_labels
        labels.position = XL_LABEL_POSITION.OUTSIDE_END
        labels.number_format = "0.0"
        labels.font.name = FONT
        labels.font.size = Pt(8)
        for point_index, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = rgb(ORANGE if point_index == 3 else series_colours[series_index])
            point.format.line.fill.background()

    deck.panel(slide, 8.46, 2.00, 3.96, 4.22, accent=ORANGE, wash=0.95)
    deck.pill(slide, 8.64, 2.17, "REALITY CHECK", fill=ORANGE, colour=NAVY, w=1.45)
    deck.rich_text(
        slide,
        8.67,
        2.66,
        3.58,
        2.62,
        [
            ("Quality offers 8.30 Mbps.", NAVY, True),
            ("Poor-link central reference: 10.39 Mbps.", NAVY, True),
            ("Quality × Poor projected load ratio ≈ 80%.", DEEP, True),
            ("The other 15 cells project below the central capacity references.", GREY, False),
            ("Exact 10-Hz latency, drops, accepted updates, and AoI are not yet measured.", RED, True),
        ],
        size=8.9,
    )
    deck.rect(slide, 8.67, 5.39, 3.54, 0.59, fill=tint(TEAL, 0.86), edge=mix(TEAL, WHITE, 0.35), rounded=True)
    deck.text(
        slide,
        8.84,
        5.52,
        3.20,
        0.30,
        "A valid result may be: no switching is needed here.",
        size=8.9,
        bold=True,
        colour=DEEP,
        align=PP_ALIGN.CENTER,
    )
    deck.text(slide, ML + 0.35, 6.24, 6.70, 0.22, "Achieved SNR / MCS:  " + "   •   ".join(snr_mcs), size=7.5, colour=GREY)


def build_measurement_matrix(deck: Deck) -> None:
    slide = deck.content_slide(
        5,
        "Measurement sequence",
        "Define 16 logical cells; directly measure only two before expanding",
        "The matrix is a planning surface—not a claim that any fixed-10-Hz cell has already run.",
        source="Source: UE_SPLIT_ONLY_SUPERVISOR_COMBINATIONS_V1.csv",
        evidence_status="PROJECTED",
        notes=SLIDE_NOTES[5],
    )
    deck.pill(slide, RIGHT - 1.95, 1.88, "PLANNED • NOT RUN", fill=ORANGE, colour=NAVY, w=1.95, h=0.25, size=7.6)

    headers = ("Action", "Clear", "Mild", "Mid", "Poor")
    matrix = [
        ("Degraded rescue", "Projected", "Projected", "Projected", "Last resort"),
        ("Compact normal", "Projected", "Projected", "Projected", "Conditional"),
        ("Balanced normal", "Projected", "Projected", "Projected", "Conditional"),
        ("Quality normal", "Initial control", "Projected", "Diagnostic", "Initial boundary"),
    ]
    shape = slide.shapes.add_table(5, 5, Inches(ML), Inches(2.25), Inches(7.25), Inches(3.42))
    table = shape.table
    widths = (1.62, 1.40, 1.32, 1.32, 1.59)
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    table.rows[0].height = Inches(0.46)
    for index in range(1, 5):
        table.rows[index].height = Inches(0.69)
    for column, value in enumerate(headers):
        deck.style_cell(table.cell(0, column), value, fill=NAVY, colour=WHITE, size=9, bold=True)
    for row_index, values in enumerate(matrix, start=1):
        role_key = PROFILE_ORDER[row_index - 1]
        for column, value in enumerate(values):
            if column == 0:
                fill, colour, bold = tint(PROFILE_ACCENTS[role_key], 0.88), NAVY, True
            elif (row_index, column) in {(4, 1), (4, 4)}:
                fill, colour, bold = tint(BLUE, 0.80), DEEP, True
            elif (row_index, column) in {(1, 4), (2, 4), (3, 4), (4, 3)}:
                fill, colour, bold = tint(ORANGE, 0.78), NAVY, True
            else:
                fill, colour, bold = tint(GREY, 0.92), GREY, False
            deck.style_cell(
                table.cell(row_index, column),
                value,
                fill=fill,
                colour=colour,
                size=8.5,
                bold=bold,
                align=PP_ALIGN.LEFT if column == 0 else PP_ALIGN.CENTER,
            )

    deck.panel(slide, 8.47, 2.25, 3.95, 3.42, accent=BLUE, wash=0.96)
    deck.pill(slide, 8.66, 2.43, "ADAPTIVE FIRST ROUND", fill=BLUE, w=1.85)
    steps = [
        ("1", "Quality × Clear", "exact-path control", BLUE),
        ("2", "Quality × Poor", "only current boundary", BLUE),
        ("?", "If fail / borderline", "Balanced → Compact", ORANGE),
        ("!", "Only if no normal works", "Degraded rescue", ORANGE),
    ]
    y = 2.91
    for index, (marker, heading, body, accent) in enumerate(steps):
        deck.rect(slide, 8.66, y, 3.56, 0.52, fill=WHITE, edge=mix(accent, WHITE, 0.48), rounded=True)
        deck.pill(slide, 8.79, y + 0.13, marker, fill=accent, colour=NAVY if accent == ORANGE else WHITE, w=0.34, h=0.26, size=8)
        deck.text(slide, 9.27, y + 0.08, 1.72, 0.22, heading, size=8.6, bold=True, colour=NAVY)
        deck.text(slide, 10.94, y + 0.10, 1.12, 0.21, body, size=7.4, colour=GREY, align=PP_ALIGN.RIGHT)
        y += 0.65

    deck.pill(slide, ML, 5.93, "INITIAL", fill=BLUE, w=0.74, h=0.22, size=7)
    deck.text(slide, ML + 0.86, 5.94, 1.45, 0.20, "direct measurement", size=7.5, colour=GREY)
    deck.pill(slide, ML + 2.55, 5.93, "CONDITIONAL", fill=ORANGE, colour=NAVY, w=1.08, h=0.22, size=7)
    deck.text(slide, ML + 3.75, 5.94, 1.55, 0.20, "measure only if needed", size=7.5, colour=GREY)
    deck.pill(slide, ML + 5.55, 5.93, "PROJECTED", fill=GREY, w=0.96, h=0.22, size=7)
    deck.text(slide, ML + 6.63, 5.94, 2.05, 0.20, "feasibility only; outcomes unresolved", size=7.5, colour=GREY)
    deck.panel(slide, ML + 1.25, 6.25, CONTENT_W - 2.50, 0.34, accent=TEAL, wash=0.91)
    deck.text(
        slide,
        ML + 1.40,
        6.32,
        CONTENT_W - 2.80,
        0.18,
        "If Quality × Poor is stable, stop the first round and report that the current regimes do not require switching.",
        size=7.9,
        bold=True,
        colour=DEEP,
        align=PP_ALIGN.CENTER,
    )


def build_measurements(deck: Deck) -> None:
    slide = deck.content_slide(
        6,
        "Measurement contract",
        "Fixed 10-Hz replay links processing, delivery, and freshness",
        "Count only accepted map updates; send and enqueue timestamps are intermediate stages, not the service endpoint.",
        source="Source: supervisor plan §§5–7; full technical contract in UE_SPLIT_ONLY_EXPERIMENT_PLAN.md",
        evidence_status="PLANNED",
        notes=SLIDE_NOTES[6],
    )
    stages = [
        ("RELEASE", "new monotonic\n10-Hz timestamp", BLUE),
        ("UE FRONT", "encode + quantize\nROI + serialize", DEEP),
        ("OAI", "schedule + deliver\nqueue / retx / drops", ORANGE),
        ("EDGE TAIL", "decode + inference\npost-process", NAVY),
        ("MAP UPDATE", "validate + accept\nupdate complete", TEAL),
    ]
    node_w = 1.90
    gap = 0.50
    for index, (heading, body, accent) in enumerate(stages):
        x = ML + index * (node_w + gap)
        deck.rect(slide, x, 2.05, node_w, 1.05, fill=tint(accent, 0.93), edge=mix(accent, WHITE, 0.38), rounded=True)
        deck.text(slide, x + 0.12, 2.21, node_w - 0.24, 0.22, heading, size=8.4, bold=True, colour=accent, align=PP_ALIGN.CENTER)
        deck.text(slide, x + 0.12, 2.53, node_w - 0.24, 0.38, body, size=7.8, colour=GREY, align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            deck.arrow(slide, x + node_w + 0.06, 2.57, x + node_w + gap - 0.07, 2.57, colour=BLUE)

    deck.panel(slide, ML + 1.20, 3.38, CONTENT_W - 2.40, 0.69, accent=TEAL, wash=0.90)
    deck.pill(slide, ML + 1.38, 3.55, "MAP FRESHNESS", fill=TEAL, colour=NAVY, w=1.35)
    deck.text(
        slide,
        ML + 2.95,
        3.49,
        CONTENT_W - 4.45,
        0.31,
        "AoI(t) = t − release_time(newest accepted map update)",
        size=12.5,
        bold=True,
        colour=DEEP,
        align=PP_ALIGN.CENTER,
    )

    card_w = (CONTENT_W - 0.36) / 3
    deck.info_card(
        slide,
        ML,
        4.42,
        card_w,
        1.77,
        "FIXED CONTROLS",
        "• identical retained sequence and order\n"
        "• target 10.00-Hz wall-clock release\n"
        "• fixed hardware, checkpoints, decoder, map\n"
        "• 3 replicates × ≥200 post-warm-up releases",
        accent=BLUE,
    )
    deck.info_card(
        slide,
        ML + card_w + 0.18,
        4.42,
        card_w,
        1.77,
        "MEASURE",
        "• stage latency + accepted-update rate\n"
        "• achieved SNR / MCS + queue / drops\n"
        "• object precision / recall / world XY\n"
        "• segmentation IoU / mIoU as secondary",
        accent=ORANGE,
    )
    deck.info_card(
        slide,
        ML + 2 * (card_w + 0.18),
        4.42,
        card_w,
        1.77,
        "REPORT FAIRLY",
        "• p50 / p90 / p95 / maximum\n"
        "• delivery and deadline-miss rates\n"
        "• keep unexplained slow samples\n"
        "• derive AoI_max from age-versus-error evidence",
        accent=TEAL,
    )


def build_decision(deck: Deck) -> None:
    slide = deck.content_slide(
        7,
        "Decision",
        "The baseline tells us whether learning is necessary",
        "Success is a defensible action rule and feasibility envelope—even if the simplest controller wins.",
        source="Source: supervisor plan §§8–10",
        evidence_status="PLAN",
        notes=SLIDE_NOTES[7],
    )

    # Decision tree (left).
    deck.panel(slide, ML, 2.02, 7.55, 4.10, accent=BLUE, wash=0.97)
    deck.rect(slide, 1.58, 2.24, 6.05, 0.63, fill=NAVY, rounded=True)
    deck.text(slide, 1.78, 2.41, 5.65, 0.25, "Is Quality × Poor stable at true 10 Hz?", size=11.2, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    deck.arrow(slide, 4.61, 2.88, 4.61, 3.13, colour=BLUE)

    deck.rect(slide, 1.30, 3.18, 2.74, 0.82, fill=tint(TEAL, 0.86), edge=mix(TEAL, WHITE, 0.28), rounded=True)
    deck.pill(slide, 1.47, 3.32, "YES", fill=TEAL, colour=NAVY, w=0.55)
    deck.text(slide, 2.15, 3.27, 1.69, 0.50, "Use Quality / greedy baseline\nStop the first round", size=8.6, bold=True, colour=DEEP, align=PP_ALIGN.CENTER)

    deck.rect(slide, 4.35, 3.18, 3.30, 0.82, fill=tint(ORANGE, 0.84), edge=mix(ORANGE, WHITE, 0.25), rounded=True)
    deck.pill(slide, 4.53, 3.32, "NO / BORDER", fill=ORANGE, colour=NAVY, w=1.02)
    deck.text(slide, 5.68, 3.27, 1.74, 0.50, "Balanced → Compact\nRescue only if needed", size=8.4, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)

    deck.arrow(slide, 6.00, 4.02, 6.00, 4.35, colour=ORANGE)
    deck.rect(slide, 4.35, 4.40, 3.30, 0.62, fill=WHITE, edge=mix(DEEP, WHITE, 0.35), rounded=True)
    deck.text(slide, 4.55, 4.58, 2.90, 0.25, "Does the preferred action depend on history?", size=8.9, bold=True, colour=DEEP, align=PP_ALIGN.CENTER)

    deck.arrow(slide, 5.18, 5.03, 5.18, 5.31, colour=DEEP)
    deck.arrow(slide, 6.82, 5.03, 6.82, 5.31, colour=DEEP)
    deck.rect(slide, 4.35, 5.36, 1.55, 0.50, fill=tint(DEEP, 0.91), edge=mix(DEEP, WHITE, 0.35), rounded=True)
    deck.text(slide, 4.49, 5.49, 1.27, 0.22, "NO → deterministic rule", size=7.7, bold=True, colour=DEEP, align=PP_ALIGN.CENTER)
    deck.rect(slide, 6.10, 5.36, 1.55, 0.50, fill=tint(BLUE, 0.86), edge=mix(BLUE, WHITE, 0.30), rounded=True)
    deck.text(slide, 6.24, 5.49, 1.27, 0.22, "YES → policy evidence", size=7.7, bold=True, colour=DEEP, align=PP_ALIGN.CENTER)

    # Discussion decisions (right).
    deck.panel(slide, 8.74, 2.02, 3.68, 4.10, accent=ORANGE, wash=0.96)
    deck.pill(slide, 8.94, 2.22, "SUPERVISOR DISCUSSION", fill=ORANGE, colour=NAVY, w=1.90)
    decisions = [
        ("1", "Single-UE, split-only scope"),
        ("2", "Three normal + one rescue"),
        ("3", "Clear / Mild / Mid / Poor regimes"),
        ("4", "Two-cell initial measurement"),
        ("5", "Segmentation secondary"),
        ("6", "Derive AoI_max from evidence"),
    ]
    y = 2.72
    for number, label in decisions:
        deck.pill(slide, 8.98, y, number, fill=BLUE, w=0.31, h=0.25, size=7.5)
        deck.text(slide, 9.43, y + 0.01, 2.70, 0.25, label, size=8.5, bold=True, colour=NAVY)
        y += 0.48
    deck.rect(slide, 8.96, 5.62, 3.20, 0.32, fill=tint(RED, 0.92), edge=mix(RED, WHITE, 0.45), rounded=True)
    deck.text(slide, 9.10, 5.69, 2.92, 0.18, "No experiment starts from this deck alone", size=7.6, bold=True, colour=RED, align=PP_ALIGN.CENTER)

    deck.panel(slide, ML + 1.38, 6.28, CONTENT_W - 2.76, 0.32, accent=TEAL, wash=0.90)
    deck.text(
        slide,
        ML + 1.55,
        6.34,
        CONTENT_W - 3.10,
        0.18,
        "Measure the decision boundary first; build the UE policy only from evidence that survives the simplest baseline.",
        size=7.9,
        bold=True,
        colour=DEEP,
        align=PP_ALIGN.CENTER,
    )


def write_notes(deck: Deck) -> None:
    lines = [
        "# Presenter notes — UE split-inference supervisor baseline",
        "",
        "The deck content is fully editable. The notes below are also embedded in each slide's speaker notes.",
        "",
    ]
    for record in deck.records:
        lines.extend((f"## Slide {record.number} — {record.title}", "", SLIDE_NOTES[record.number].strip(), ""))
    NOTES.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def write_manifest(deck: Deck, profiles: dict[str, dict[str, float | str]]) -> None:
    source_paths = (TEMPLATE, PLAN, COMBINATIONS, CANDIDATES, DECODER_REPORT)
    payload = {
        "schema": "scenesense.ue_split_supervisor_deck.v1",
        "created_date": "2026-08-20",
        "deck": str(PPTX.relative_to(REPO)),
        "deck_sha256": sha256(PPTX),
        "slide_count": len(deck.records),
        "template": str(TEMPLATE.relative_to(REPO)),
        "template_sha256": sha256(TEMPLATE),
        "classification_expected": "Non-Public",
        "editable_content_contract": {
            "slide_bitmaps": False,
            "native_text": True,
            "native_shapes": True,
            "native_tables": True,
            "native_chart": True,
        },
        "measurement_authority": False,
        "profiles": profiles,
        "slides": [record.__dict__ for record in deck.records],
        "sources": {
            str(path.relative_to(REPO)): {
                "sha256": sha256(path),
                "bytes": path.stat().st_size,
            }
            for path in source_paths
        },
    }
    MANIFEST.write_text(json.dumps(payload, indent=2, sort_keys=True) + "\n", encoding="utf-8")


def validate_output(expected_slides: int = 7) -> None:
    prs = Presentation(str(PPTX))
    if len(prs.slides) != expected_slides:
        raise RuntimeError(f"expected {expected_slides} slides, found {len(prs.slides)}")
    if (round(prs.slide_width / 914400, 4), round(prs.slide_height / 914400, 4)) != (13.3333, 7.5):
        raise RuntimeError("output slide size does not match the 16:9 template")
    chart_count = table_count = picture_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            chart_count += int(getattr(shape, "has_chart", False))
            table_count += int(getattr(shape, "has_table", False))
            picture_count += int(shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    if chart_count != 1:
        raise RuntimeError(f"expected one native chart, found {chart_count}")
    if table_count != 2:
        raise RuntimeError(f"expected two native tables, found {table_count}")
    if picture_count != 0:
        raise RuntimeError(f"content slides unexpectedly contain {picture_count} picture shapes")
    with zipfile.ZipFile(PPTX) as archive:
        custom = archive.read("docProps/custom.xml").decode("utf-8")
        if "Non-Public" not in custom:
            raise RuntimeError("InterDigital Non-Public sensitivity label was not preserved")


def build() -> None:
    check_inputs()
    profiles, combinations = load_evidence()
    OUT.mkdir(parents=True, exist_ok=True)

    deck = Deck()
    build_title(deck)
    build_scope(deck)
    build_actions(deck, profiles)
    build_network(deck, profiles, combinations)
    build_measurement_matrix(deck)
    build_measurements(deck)
    build_decision(deck)

    if len(deck.records) != 7:
        raise RuntimeError("builder did not produce the expected seven slides")
    deck.prs.core_properties.title = "SceneSense UE Split-Inference Baseline"
    deck.prs.core_properties.subject = "Supervisor discussion and bounded measurement plan"
    deck.prs.core_properties.author = "Abiodun"
    deck.prs.core_properties.comments = "Editable native PowerPoint content built from the InterDigital template."
    deck.prs.save(str(PPTX))
    validate_output()
    write_notes(deck)
    write_manifest(deck, profiles)
    print(f"wrote {PPTX}")
    print(f"wrote {NOTES}")
    print(f"wrote {MANIFEST}")


if __name__ == "__main__":
    build()
