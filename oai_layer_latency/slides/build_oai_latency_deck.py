#!/usr/bin/env python3
"""Build the OAI latency investigation team slide deck.

The deck intentionally uses only local, reportable artifacts from the corrected
CARLA/OAI runs and local OAI source/reference files. It does not pull anything
from the network.
"""

from __future__ import annotations

import csv
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# `ROOT` is the project subfolder `abiodun/`, not the repository root.
ROOT = Path(__file__).resolve().parents[2]
SLIDE_DIR = Path(__file__).resolve().parent
OUT = SLIDE_DIR / "oai_latency_investigation_20260723.pptx"
NOTES_OUT = SLIDE_DIR / "oai_latency_investigation_20260723_talk_track.md"

PLOTS = {
    "latency_breakdown": ROOT
    / "downlink_latency_fps/plots/oai_bottleneck/corrected_transport_latency_breakdown.png",
    "reliability_rtt": ROOT
    / "downlink_latency_fps/plots/oai_bottleneck/corrected_transport_reliability_rtt.png",
    "zstd_ab": ROOT
    / "downlink_latency_fps/plots/oai_bottleneck/oai_106prb_drivable_zlib_vs_zstd.png",
    "layer_latency": ROOT / "oai_layer_latency/plots/complementary_latency_summary.png",
    "mcs_prb": ROOT / "oai_layer_latency/plots/complementary_mcs_prb_summary.png",
    "rlc_buffer": ROOT / "oai_layer_latency/plots/complementary_rlc_buffer_timeseries.png",
    "snr": ROOT / "oai_layer_latency/plots/complementary_gnb_snr_timeseries.png",
    "mcs_bottleneck": ROOT / "oai_layer_latency/plots/uplink_mcs_bottleneck_summary.png",
    "advisor": ROOT / "oai_layer_latency/plots/advisor_iperf_vs_carla_bsr_mcs.png",
    "drain_rate": ROOT / "oai_layer_latency/plots/uplink_drain_rate.png",
    "bler_branch": ROOT / "oai_layer_latency/plots/bler_olla_branch_comparison.png",
    "bler_mcs": ROOT / "oai_layer_latency/plots/bler_olla_mcs_timeseries.png",
    "bler_num_sched": ROOT / "oai_layer_latency/plots/bler_olla_num_sched_timeseries.png",
    "mcs_decision": ROOT
    / "metrics_logs/scenesense_ttracer/carla_shape_udp_bw273_mcsdecision_20260723_171153/gnb/analysis/mcs_decision_timeseries.png",
    "phr_drop": ROOT
    / "metrics_logs/scenesense_ttracer/carla_shape_udp_bw273_mcsdecision_20260723_171153/gnb/analysis/mcs_decision_phr_drop_hist.png",
}

SUMMARY_CSV = ROOT / "oai_layer_latency/plots/complementary_experiment_summary.csv"
AWGN_DIR = (
    ROOT
    / "OAI/openairinterface5g/openair1/SIMULATION/NR_PHY/BLER_SIMULATIONS/AWGN/AWGN_results"
)


BG = RGBColor(248, 250, 252)
INK = RGBColor(24, 32, 42)
MUTED = RGBColor(82, 96, 112)
BLUE = RGBColor(38, 112, 196)
TEAL = RGBColor(29, 143, 140)
ORANGE = RGBColor(229, 132, 42)
RED = RGBColor(204, 72, 72)
PURPLE = RGBColor(122, 86, 184)
GREEN = RGBColor(58, 145, 92)
LIGHT_BLUE = RGBColor(224, 237, 252)
LIGHT_ORANGE = RGBColor(255, 238, 217)
LIGHT_RED = RGBColor(252, 230, 230)
LIGHT_GREEN = RGBColor(226, 244, 233)
WHITE = RGBColor(255, 255, 255)
GRID = RGBColor(217, 225, 234)


@dataclass
class ExpRow:
    label: str
    feature_kb: float
    chunks: float
    delivery: float
    rtt_ms: float
    rtt_p95_ms: float
    ran_ms: float
    ran_p95_ms: float
    rlc_queue_ms: float
    mcs_p50: float
    mcs_p95: float
    prb_p50: float
    tbs_p50: float
    snr_p50_db: float


def load_summary() -> list[ExpRow]:
    rows: list[ExpRow] = []
    with SUMMARY_CSV.open(newline="") as f:
        for r in csv.DictReader(f):
            rows.append(
                ExpRow(
                    label=r["label"].replace("\n", " "),
                    feature_kb=float(r["feature_kb"]),
                    chunks=float(r["feature_chunks"]),
                    delivery=float(r["delivery"]),
                    rtt_ms=float(r["rtt_ms"]),
                    rtt_p95_ms=float(r["rtt_p95_ms"]),
                    ran_ms=float(r["ran_p50_ms"]),
                    ran_p95_ms=float(r["ran_p95_ms"]),
                    rlc_queue_ms=float(r["rlc_queue_ms"]),
                    mcs_p50=float(r["mcs_p50"]),
                    mcs_p95=float(r["mcs_p95"]),
                    prb_p50=float(r["prb_p50"]),
                    tbs_p50=float(r["tbs_p50"]),
                    snr_p50_db=float(r["snr_p50_db"]),
                )
            )
    return rows


MCS_TABLE_1 = {
    # 3GPP TS 38.214 Table 5.1.3.1-1 / OAI mcs_table 0.
    # MCS index: (Qm, target_code_rate_x1024, spectral_efficiency)
    0: (2, 120, 0.2344),
    2: (2, 193, 0.3770),
    4: (2, 308, 0.6016),
    5: (2, 379, 0.7402),
    8: (2, 602, 1.1758),
    10: (4, 340, 1.3281),
    16: (4, 658, 2.5703),
    20: (6, 567, 3.3223),
    24: (6, 772, 4.5234),
    28: (6, 948, 5.5547),
}

OAI_SINR_MCS_DB = {
    # gNB_scheduler_primitives.c SINRx10_MCS_mapping; single-layer thresholds.
    # OAI comment: chosen to maintain BLER around 10^-3.
    0: -1.0,
    4: 2.4,
    8: 5.6,
    16: 12.4,
    18: 14.6,
    24: 19.4,
    28: 24.5,
}


def modulation(qm: int) -> str:
    return {2: "QPSK", 4: "16QAM", 6: "64QAM", 8: "256QAM"}.get(qm, f"Qm={qm}")


def awgn_thresholds(mcs_values: Iterable[int], target_bler: float = 0.10) -> list[tuple[int, float | None, float]]:
    """Return first SNR where OAI AWGN reference BLER <= target.

    This is an interpretive reference curve, not a live CQI measurement.
    """

    out: list[tuple[int, float | None, float]] = []
    for mcs in mcs_values:
        path = AWGN_DIR / f"mcs{mcs}_awgn_5G.csv"
        threshold = None
        rate = MCS_TABLE_1[mcs][2]
        if path.exists():
            with path.open(newline="") as f:
                for row in csv.DictReader(f, delimiter=";"):
                    trials = float(row.get("trials0") or 0)
                    err = float(row.get("err0") or 0)
                    if trials > 0 and err / trials <= target_bler:
                        threshold = float(row["SNR"])
                        rate = float(row.get("rate") or rate)
                        break
        out.append((mcs, threshold, rate))
    return out


def set_slide_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    font_size=18,
    color=INK,
    bold=False,
    align=None,
    fill=None,
    line=None,
    radius=True,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(1)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(
        slide,
        Inches(0.35),
        Inches(0.25),
        Inches(12.6),
        Inches(0.48),
        title,
        font_size=25,
        color=INK,
        bold=True,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.39),
            Inches(0.75),
            Inches(12.1),
            Inches(0.32),
            subtitle,
            font_size=10.5,
            color=MUTED,
        )


def add_footer(slide, idx, total):
    add_textbox(
        slide,
        Inches(11.9),
        Inches(7.05),
        Inches(1.0),
        Inches(0.2),
        f"{idx}/{total}",
        font_size=8,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_bullets(slide, x, y, w, h, bullets, font_size=15, color=INK, line_spacing=1.05):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4 * line_spacing)
    return box


def add_caption(slide, x, y, w, text):
    return add_textbox(slide, x, y, w, Inches(0.22), text, font_size=8.5, color=MUTED)


def add_image_fit(slide, path: Path, x, y, w, h, border=True):
    if not path.exists():
        return add_textbox(
            slide,
            x,
            y,
            w,
            h,
            f"Missing plot:\n{path.relative_to(ROOT)}",
            font_size=12,
            color=RED,
            fill=LIGHT_RED,
            line=RED,
        )
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    left = x + (w - pw) / 2
    top = y + (h - ph) / 2
    pic = slide.shapes.add_picture(str(path), left, top, width=pw, height=ph)
    if border:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
        rect.fill.background()
        rect.line.color.rgb = GRID
        rect.line.width = Pt(0.75)
    return pic


def style_table(table, header_fill=BLUE, body_fill=WHITE, font_size=8.5):
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else body_fill
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(font_size if r else font_size + 0.5)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else INK


def add_table(slide, x, y, w, h, rows, col_widths=None, font_size=8.5, header_fill=BLUE):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = shape.table
    for r, vals in enumerate(rows):
        for c, val in enumerate(vals):
            table.cell(r, c).text = str(val)
    if col_widths:
        for c, width in enumerate(col_widths):
            table.columns[c].width = width
    style_table(table, header_fill=header_fill, font_size=font_size)
    return table


def fmt_ms(v: float) -> str:
    return f"{v:.0f} ms"


def fmt_pct(v: float) -> str:
    return f"{100*v:.1f}%"


def build_deck():
    rows = load_summary()
    by_label = {r.label: r for r in rows}
    adaptive_u8 = by_label["273 adaptive uint8"]
    adaptive_u4 = by_label["273 adaptive uint4"]
    fixed273 = by_label["273 fixed MCS28"]
    fixed106 = by_label["106 fixed MCS28"]
    ae = by_label["106 default AE128 u6 r0.5"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 10

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.55), Inches(0.55), Inches(12.1), Inches(0.65),
                "OAI uplink latency investigation", 31, INK, True)
    add_textbox(slide, Inches(0.58), Inches(1.18), Inches(12.0), Inches(0.36),
                "CARLA split-fusion deployment • corrected drivable route • 10 FPS OAI focus", 14, MUTED)
    cards = [
        ("Dominant bottleneck", "Uplink feature transfer / UE RLC queue", RED, LIGHT_RED),
        ("Primary trigger", "Sparse closed-loop ~1 MB bursts interact badly with OAI BLER/OLLA", ORANGE, LIGHT_ORANGE),
        ("Measured symptom", "Adaptive UL MCS stays QPSK-region despite high RFsim gNB PUSCH SNR", BLUE, LIGHT_BLUE),
        ("Best practical relief", "Reduce payload: AE-128/u6/ROI0.5 → 99.8% overall delivery, 100% after startup", GREEN, LIGHT_GREEN),
    ]
    x0, y0 = Inches(0.7), Inches(2.05)
    for i, (h, b, color, fill) in enumerate(cards):
        x = x0 + Inches((i % 2) * 6.05)
        y = y0 + Inches((i // 2) * 1.55)
        add_textbox(slide, x, y, Inches(5.65), Inches(0.36), h, 13, color, True, fill=fill, line=color)
        add_textbox(slide, x, y + Inches(0.42), Inches(5.65), Inches(0.62), b, 14, INK, False, fill=WHITE, line=color)
    add_textbox(slide, Inches(0.75), Inches(5.55), Inches(11.9), Inches(0.55),
                "Headline: fixed MCS28 proves spectral efficiency is the lever (RTT 186→47 ms); the direct BLER/OLLA trace shows sparse closed-loop bursts keep adaptive MCS low, while payload reduction is the deployable mitigation.", 15.2, INK, True, fill=LIGHT_BLUE, line=BLUE)
    add_footer(slide, 1, total)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Experiment setup and measurement path", "What was measured in the corrected live CARLA/OAI runs")
    # Pipeline boxes
    pipeline = [
        ("CARLA front", "camera/radar → split feature tensor"),
        ("UE / OAI uplink", "UDP chunks over oaitun_ue1"),
        ("gNB/CN", "MAC/RLC/PDCP + GTP-U path"),
        ("Edge tail model", "fused detections + loc output"),
        ("Downlink result", "boxes/centroids/scores back to car"),
    ]
    for i, (h, b) in enumerate(pipeline):
        x = Inches(0.55 + i * 2.55)
        add_textbox(slide, x, Inches(1.35), Inches(2.15), Inches(0.42), h, 12, WHITE, True, align=PP_ALIGN.CENTER, fill=BLUE, line=BLUE)
        add_textbox(slide, x, Inches(1.80), Inches(2.15), Inches(0.68), b, 9.5, INK, align=PP_ALIGN.CENTER, fill=WHITE, line=GRID)
        if i < len(pipeline) - 1:
            add_textbox(slide, x + Inches(2.16), Inches(1.65), Inches(0.30), Inches(0.22), "→", 18, MUTED, True)
    add_bullets(slide, Inches(0.75), Inches(3.05), Inches(5.8), Inches(2.3), [
        "Corrected scene: 28 vehicles, 35 pedestrians, seed 31, fixed drivable route.",
        "Main Step-1 baseline: no-AE, ROI 0, per-channel uint8, 200k radar PPS, zstd/zlib where noted.",
        "Layer probe: matched UE PDCP/RLC and gNB MAC/PDCP timestamps under RFsim, plus app RTT/downlink metrics.",
        "Closed-loop behavior: frontend waits for result/timeout before advancing, so missing frames reduce effective rate.",
    ], font_size=12.2)
    setup_rows = [
        ["Run family", "What it isolates"],
        ["Loopback", "App/model overhead and downlink floor without OAI RAN queue"],
        ["106PRB default / UL-heavy", "Default OAI path and TDD uplink-share sensitivity"],
        ["273PRB", "Bandwidth/PRB ceiling without changing app payload"],
        ["Fixed MCS28", "Diagnostic: what if spectral efficiency were high?"],
        ["Reduced payload", "Deployable mitigation: smaller feature bursts"],
    ]
    add_table(slide, Inches(6.85), Inches(3.02), Inches(5.85), Inches(2.35), setup_rows, font_size=8.8, header_fill=TEAL)
    add_textbox(slide, Inches(0.75), Inches(6.05), Inches(11.8), Inches(0.42),
                "Important reporting rule: old 60-vehicle frontend runs were cleaned out and should not be used; this deck uses corrected drivable-route artifacts only.", 11, RED, True, fill=LIGHT_RED, line=RED)
    add_footer(slide, 2, total)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Step-1 latency breakdown: uplink dominates; downlink is tiny", "Downlink carries only detections/boxes/centroids/scores, not the dense feature tensor")
    add_image_fit(slide, PLOTS["latency_breakdown"], Inches(0.35), Inches(1.13), Inches(7.2), Inches(4.65))
    break_rows = [
        ["Condition", "UL payload", "DL payload", "RTT p50", "Capture→result", "Delivery"],
        ["Ideal loopback zstd", "1053.9 KB / 18 chunks", "2.4 KB / 1 chunk", "18.3 ms", "46.1 ms", "100.0%"],
        ["OAI 106PRB zstd", "1055.2 KB / 19 chunks", "2.2 KB / 1 chunk", "162.2 ms", "188.0 ms", "83.6%"],
        ["OAI UL-heavy 106PRB", "1054.0 KB / 18 chunks", "2.5 KB / 1 chunk", "152.1 ms", "177.3 ms", "84.8%"],
        ["OAI 273PRB zstd", "1054.6 KB / 19 chunks", "2.4 KB / 1 chunk", "186.1 ms", "212.7 ms", "85.4%"],
        ["OAI 106PRB AE128/u6/r0.5", "152.7 KB / 3 chunks", "2.3 KB / 1 chunk", "64.2 ms", "86.5 ms", "99.8%"],
    ]
    add_table(slide, Inches(7.78), Inches(1.17), Inches(5.18), Inches(3.18), break_rows, font_size=7.4, header_fill=BLUE)
    add_bullets(slide, Inches(7.9), Inches(4.62), Inches(4.75), Inches(1.25), [
        "The downlink result is only a few KB and one UDP chunk, so it remains cheap.",
        "The heavy direction is car/front → edge/back-half: the split feature tensor is ~1 MB/frame for no-AE.",
        "Report capture→result for end-to-end perception latency; RTT alone starts after feature send begins.",
    ], font_size=10.6)
    add_caption(slide, Inches(0.46), Inches(5.85), Inches(7.0), "Source plot: downlink_latency_fps/plots/oai_bottleneck/corrected_transport_latency_breakdown.png")
    add_footer(slide, 3, total)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Layer localization: the long delay is UE RLC queue-wait", "Per-packet timestamps and Little’s-law queue estimate converge on the same bottleneck")
    add_image_fit(slide, PLOTS["mcs_bottleneck"], Inches(0.42), Inches(1.12), Inches(6.35), Inches(4.75))
    layer_rows = [
        ["Condition", "App RTT", "RAN UL", "RLC queue", "MCS", "Delivery"],
        ["273 adaptive uint8", fmt_ms(adaptive_u8.rtt_ms), fmt_ms(adaptive_u8.ran_ms), fmt_ms(adaptive_u8.rlc_queue_ms), "4 / 8", fmt_pct(adaptive_u8.delivery)],
        ["273 fixed MCS28", fmt_ms(fixed273.rtt_ms), fmt_ms(fixed273.ran_ms), fmt_ms(fixed273.rlc_queue_ms), "28 / 28", fmt_pct(fixed273.delivery)],
        ["106 fixed MCS28", fmt_ms(fixed106.rtt_ms), fmt_ms(fixed106.ran_ms), fmt_ms(fixed106.rlc_queue_ms), "28 / 28", fmt_pct(fixed106.delivery)],
        ["106 AE128 adaptive", fmt_ms(ae.rtt_ms), fmt_ms(ae.ran_ms), fmt_ms(ae.rlc_queue_ms), "2 / 5", fmt_pct(ae.delivery)],
    ]
    add_table(slide, Inches(7.05), Inches(1.18), Inches(5.78), Inches(2.15), layer_rows, font_size=8.2, header_fill=TEAL)
    add_bullets(slide, Inches(7.1), Inches(3.72), Inches(5.58), Inches(1.62), [
        "Adaptive no-AE: ~1 MB frames build a queue; RAN UL p50 ≈112 ms and RLC queue ≈103 ms.",
        "Fixed MCS28: same full payload, same ideal RFsim channel, queue collapses to ~13 ms.",
        "Therefore the latency lever is not edge inference, PDCP handoff, or downlink boxes; it is uplink scheduling/spectral efficiency.",
    ], font_size=11.3)
    add_caption(slide, Inches(0.46), Inches(5.93), Inches(6.6), "Source plot: oai_layer_latency/plots/uplink_mcs_bottleneck_summary.png")
    add_footer(slide, 4, total)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Why MCS matters: low MCS caps spectral efficiency", "OAI mcs_table 0 uses 3GPP TS 38.214 Table 5.1.3.1-1")
    mcs_rows = [["MCS", "Mod.", "Qm", "R×1024", "Eff. bits/RE", "Relative to MCS28"]]
    for idx in [0, 2, 4, 5, 8, 10, 16, 20, 24, 28]:
        qm, r, eff = MCS_TABLE_1[idx]
        mcs_rows.append([idx, modulation(qm), qm, r, f"{eff:.3f}", f"{eff / MCS_TABLE_1[28][2]:.0%}"])
    add_table(slide, Inches(0.55), Inches(1.18), Inches(5.72), Inches(4.95), mcs_rows, font_size=8.4, header_fill=PURPLE)
    add_image_fit(slide, PLOTS["mcs_prb"], Inches(6.55), Inches(1.15), Inches(6.45), Inches(3.85))
    add_textbox(slide, Inches(6.82), Inches(5.25), Inches(5.82), Inches(0.67),
                "Observed adaptive rows sit around MCS 2–5, i.e. QPSK and only ~7–13% of MCS28 spectral efficiency. That is enough to make a 1 MB frame drain over ~100 ms.", 13, INK, True, fill=LIGHT_ORANGE, line=ORANGE)
    add_caption(slide, Inches(0.62), Inches(6.18), Inches(12.0), "MCS lookup source: local OAI scheduler primitive supports mcs_table 0 / TS 38.214 Table 5.1.3.1-1; plot source: complementary_mcs_prb_summary.png")
    add_footer(slide, 5, total)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Advisor check: traffic class is not the visible differentiator", "iperf and CARLA appear on the same data bearer; the big difference is BSR backlog")
    add_image_fit(slide, PLOTS["advisor"], Inches(0.35), Inches(1.08), Inches(8.55), Inches(5.45))
    adv_rows = [
        ["Evidence", "iperf", "CARLA"],
        ["Bearer seen at MAC", "LCID 4 / LCG 1", "LCID 4 / LCG 1"],
        ["BSR backlog p50", "1.2 KB", "685.8 KB"],
        ["UL MCS p50/p95", "28 / 28", "4 / 8"],
        ["gNB PUSCH SNR p50", "50.5 dB", "50.5 dB"],
        ["PHR p50", "65", "65"],
        ["TBS p50", "4992 B", "1633 B"],
    ]
    add_table(slide, Inches(9.13), Inches(1.15), Inches(3.9), Inches(2.63), adv_rows, font_size=7.3, header_fill=RED)
    add_bullets(slide, Inches(9.28), Inches(4.12), Inches(3.48), Inches(1.55), [
        "Current traces do not support a per-application traffic-class explanation at the gNB MAC scheduler.",
        "The observed trigger is backlog shape: CARLA reports ~1 MB BSR bursts, then adaptive MCS backs off.",
        "Pre/post MCS trace ruled out the PHR helper: selected/pre/post/final MCS were already 4/8.",
    ], font_size=9.7)
    add_caption(slide, Inches(0.48), Inches(6.58), Inches(8.5), "Source plot: oai_layer_latency/plots/advisor_iperf_vs_carla_bsr_mcs.png")
    add_footer(slide, 6, total)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "SNR/PER lookup says high MCS should be feasible in RFsim", "OAI’s SNR→MCS table targets BLER/PER around 10^-3; the later scheduler path lowers MCS")
    add_image_fit(slide, PLOTS["snr"], Inches(0.45), Inches(1.10), Inches(6.95), Inches(3.92))
    snr_rows = [["MCS", "Mod.", "Eff.", "OAI SINR threshold", "Margin @50.5 dB"]]
    for mcs in [0, 4, 8, 16, 24, 28]:
        qm = MCS_TABLE_1[mcs][0]
        eff = MCS_TABLE_1[mcs][2]
        thr = OAI_SINR_MCS_DB[mcs]
        snr_rows.append([mcs, modulation(qm), f"{eff:.3f}", f"{thr:.1f} dB", f"+{50.5 - thr:.1f} dB"])
    add_table(slide, Inches(7.48), Inches(1.25), Inches(5.24), Inches(2.55), snr_rows, font_size=7.7, header_fill=BLUE)
    add_bullets(slide, Inches(7.75), Inches(4.08), Inches(4.8), Inches(1.45), [
        "OAI's own SINR→MCS table would permit high MCS at the ~50.5 dB RFsim PUSCH SNR.",
        "So the live low-MCS behavior is not explained by channel quality in this RFsim setup.",
        "Caveat: live CQI/RSRP extraction looked like placeholders, so this slide uses gNB PUSCH SNR + OAI AWGN references.",
    ], font_size=10.7)
    add_textbox(slide, Inches(0.65), Inches(5.45), Inches(12.0), Inches(0.45),
                "Interpretation: adaptive MCS is being held low by scheduler-side BLER/OLLA state/cadence, not by a poor RF channel in this RFsim setup.", 13.4, RED, True, fill=LIGHT_RED, line=RED)
    add_caption(slide, Inches(0.55), Inches(6.13), Inches(12.4), "Source: OAI/openairinterface5g/openair2/LAYER2/NR_MAC_gNB/gNB_scheduler_primitives.c SINRx10_MCS_mapping; OAI comment says target BLER ≈10^-3")
    add_footer(slide, 7, total)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(
        slide,
        "Final mechanism: BLER/OLLA cadence, not PHR-helper clipping",
        "Direct get_mcs_from_bler() trace explains why open-loop traffic ramps high but closed-loop CARLA cadence stays low",
    )
    add_image_fit(slide, PLOTS["bler_mcs"], Inches(0.35), Inches(1.10), Inches(6.45), Inches(3.22))
    add_image_fit(slide, PLOTS["bler_branch"], Inches(6.98), Inches(1.10), Inches(5.95), Inches(3.22))
    branch_rows = [
        ["Active-window result", "Observed pace", "Open-loop 10 FPS"],
        ["MCS p50 / p95", "4 / 8", "23 / 25"],
        ["num_sched p50", "1.0", "3.5"],
        ["increase branch", "21.4%", "50.0%"],
        ["few-samples branch", "78.6%", "50.0%"],
        ["high-BLER branch", "0.0%", "0.0%"],
    ]
    add_table(slide, Inches(0.55), Inches(4.68), Inches(5.85), Inches(1.40), branch_rows, font_size=7.4, header_fill=PURPLE)
    add_bullets(slide, Inches(6.75), Inches(4.62), Inches(5.85), Inches(1.35), [
        "PHR-helper trace: selected/pre/post/final MCS were identical (4/8); `nr_ue_max_mcs_min_rb()` reduced MCS in 0 rows.",
        "BLER/OLLA trace: sparse cadence repeatedly hits the `num_sched <= 3` branch, walking MCS back down between bursts.",
        "Open-loop 10 FPS gives enough scheduled samples for the low-BLER increase branch to keep MCS in the 20s.",
    ], font_size=10.1)
    add_caption(slide, Inches(0.50), Inches(6.18), Inches(12.2), "Source plots: oai_layer_latency/plots/bler_olla_mcs_timeseries.png and bler_olla_branch_comparison.png")
    add_footer(slide, 8, total)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Controls and mitigation: fixed MCS proves the lever; reduced payload fixes reliability", "The deployable answer is not “force 64QAM everywhere”; it is payload-aware policy + scheduler investigation")
    add_image_fit(slide, PLOTS["layer_latency"], Inches(0.48), Inches(1.10), Inches(6.55), Inches(3.35))
    comp_rows = [
        ["Condition", "Payload", "Delivery", "RTT p50", "RLC queue", "Interpretation"],
        ["273 adaptive uint8", "1054 KB", "76.5%", "186 ms", "103 ms", "baseline bottleneck"],
        ["273 adaptive uint4", "395 KB", "99.8%", "113 ms", "54 ms", "payload relief"],
        ["273 fixed MCS28", "1056 KB", "73.8%", "47 ms", "13 ms", "spectral-eff. proof"],
        ["106 AE128/u6/r0.5", "153 KB", "99.8%", "64 ms", "29 ms", "deployable relief"],
    ]
    add_table(slide, Inches(7.22), Inches(1.15), Inches(5.77), Inches(2.35), comp_rows, font_size=7.2, header_fill=GREEN)
    add_bullets(slide, Inches(7.38), Inches(3.80), Inches(5.25), Inches(1.66), [
        "Fixed MCS28 lowers latency for surviving full-payload frames but delivery remains ~73–74%; those rows still send ~1 MB / 18–19 chunks.",
        "Reduced payload lowers chunk count and removes steady-state misses: AE run delivered 1298/1298 after the first two startup misses.",
        "This points to payload/backlog/chunk fragility, not downlink or model-tail compute.",
    ], font_size=10.6)
    add_image_fit(slide, PLOTS["reliability_rtt"], Inches(0.55), Inches(4.83), Inches(6.42), Inches(1.28))
    add_caption(slide, Inches(0.62), Inches(6.17), Inches(12.0), "Source plots: complementary_latency_summary.png and corrected_transport_reliability_rtt.png")
    add_footer(slide, 9, total)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "What to tell the team + next actions", "Current conclusion is strong enough to present; scheduler-side fix is the next research path")
    add_textbox(slide, Inches(0.65), Inches(1.18), Inches(12.0), Inches(0.72),
                "Current conclusion: OAI latency is dominated by uplink split-feature transport. In RFsim, sparse closed-loop feature bursts interact with OAI's BLER/OLLA selector so adaptive UL MCS stays low even with high gNB PUSCH SNR; large frames then wait in UE RLC.", 14.2, INK, True, fill=LIGHT_BLUE, line=BLUE)
    next_rows = [
        ["Priority", "Action", "Why"],
        ["1", "Use reduced-payload profile for near-term runs", "AE128/u6/ROI0.5 removes steady-state misses and cuts RTT."],
        ["2", "Add edge-side per-frame chunk accounting", "Separate missing UDP chunks, frontend timeout, and model/drop behavior."],
        ["3", "Patch/evaluate OAI BLER/OLLA burst handling", "Test whether sparse large bursts can avoid few-sample MCS backoff."],
        ["4", "Sionna / realistic channel pass", "Validate adaptive policy under fading; fixed MCS28 remains diagnostic only."],
    ]
    add_table(slide, Inches(0.75), Inches(2.25), Inches(11.85), Inches(2.05), next_rows, font_size=8.8, header_fill=BLUE)
    add_bullets(slide, Inches(0.95), Inches(4.72), Inches(5.45), Inches(1.25), [
        "Show: latency breakdown, RLC queue plot, MCS/PRB summary, SNR reference plot, MCS lookup table, BLER/OLLA branch plot.",
        "Avoid overclaim: CQI/RSRP live extraction was not reliable; use gNB PUSCH SNR and OAI AWGN references.",
    ], font_size=11.2)
    add_textbox(slide, Inches(6.85), Inches(4.68), Inches(5.55), Inches(1.35),
                "Suggested wording: “We localized the delay to UE RLC queue-wait. Fixed MCS proves spectral efficiency is the lever; the direct BLER/OLLA trace explains why sparse CARLA-like bursts stay low-MCS; payload compression is the practical mitigation.”", 12.1, INK, True, fill=LIGHT_GREEN, line=GREEN)
    add_footer(slide, 10, total)

    prs.save(OUT)
    write_notes(rows)


def write_notes(rows: list[ExpRow]):
    lines = [
        "# OAI latency investigation deck talk track",
        "",
        f"Deck: `{OUT.relative_to(ROOT)}`",
        "",
        "## One-minute presentation version",
        "",
        "We reran the OAI latency study with the corrected CARLA drivable route and focused on the live split-fusion path. The result is asymmetric: the downlink is cheap because it only returns compact detections — boxes, centroids, scores — but the uplink sends a dense split-feature tensor of roughly 1 MB/frame in the no-AE profile. Layer timestamps show the long delay is not edge inference or PDCP handoff; it is UE RLC queue-wait while the burst drains over the OAI uplink. The advisor check is now clean: iperf and CARLA are both observed on LCID 4 / LCG 1 with high RFsim gNB PUSCH SNR around 50.5 dB, so we do not see a special traffic-class explanation. The pre/post scheduler trace also ruled out the PHR helper: selected, pre-PHR, post-PHR, and final MCS were already low. The final direct BLER/OLLA trace shows the mechanism: at observed closed-loop pace, the selector spends about 79% of update decisions in the few-scheduled-samples branch and stays at MCS 4/8; at open-loop 10 FPS, enough scheduled windows arrive to keep ratcheting MCS to 23/25. Fixed MCS28 collapses RTT from about 186 ms to 47 ms, proving spectral efficiency is the lever, but it is only a diagnostic control. The deployable mitigation is payload reduction: the AE-128/u6/ROI0.5 profile cuts the feature burst to about 153 KB, gives 99.8% overall delivery, and after startup it delivered every frame.",
        "",
        "## Main caveats",
        "",
        "- Do not present fixed MCS28 as the deployment fix; it can fail under realistic fading.",
        "- Do not overstate CQI/RSRP from the current extraction. The live CQI/RSRP fields looked unreliable; the deck uses gNB PUSCH SNR plus OAI's SINR→MCS scheduler table.",
        "- Fixed MCS28 is diagnostic, not the deployment policy. The scheduler-side path to investigate is BLER/OLLA behavior in `get_mcs_from_bler()` for sparse low-latency bursts.",
        "",
        "## Selected MCS table",
        "",
        "| MCS | Modulation | Qm | R x1024 | Efficiency bits/RE |",
        "|---:|---|---:|---:|---:|",
    ]
    for idx in [0, 2, 4, 5, 8, 10, 16, 20, 24, 28]:
        qm, r, eff = MCS_TABLE_1[idx]
        lines.append(f"| {idx} | {modulation(qm)} | {qm} | {r} | {eff:.4f} |")
    lines.extend(
        [
            "",
            "## OAI SINR→MCS scheduler thresholds",
            "",
            "Source: `gNB_scheduler_primitives.c` `SINRx10_MCS_mapping`; OAI comment says the table targets BLER around `10^-3`.",
            "",
            "| MCS | SINR threshold | Margin at measured 50.5 dB | Efficiency bits/RE |",
            "|---:|---:|---:|---:|",
        ]
    )
    for mcs in [0, 4, 8, 16, 24, 28]:
        thr = OAI_SINR_MCS_DB[mcs]
        lines.append(f"| {mcs} | {thr:.1f} dB | +{50.5 - thr:.1f} dB | {MCS_TABLE_1[mcs][2]:.4f} |")
    lines.extend(
        [
            "",
            "## Source artifacts used",
            "",
        ]
    )
    for key, path in PLOTS.items():
        lines.append(f"- `{path.relative_to(ROOT)}`")
    lines.extend(
        [
            f"- `{SUMMARY_CSV.relative_to(ROOT)}`",
            "- `OAI/openairinterface5g/openair2/LAYER2/NR_MAC_gNB/gNB_scheduler_ulsch.c`",
            "- `OAI/openairinterface5g/openair2/LAYER2/NR_MAC_gNB/gNB_scheduler_primitives.c`",
            "",
            "## Numeric rows loaded into the deck",
            "",
            "| Condition | Payload KB | Delivery | RTT p50 | RAN UL p50 | RLC queue | MCS p50/p95 | PRB p50 | SNR p50 |",
            "|---|---:|---:|---:|---:|---:|---:|---:|---:|",
        ]
    )
    for r in rows:
        lines.append(
            f"| {r.label} | {r.feature_kb:.1f} | {fmt_pct(r.delivery)} | {r.rtt_ms:.1f} ms | "
            f"{r.ran_ms:.1f} ms | {r.rlc_queue_ms:.1f} ms | {r.mcs_p50:.0f}/{r.mcs_p95:.0f} | "
            f"{r.prb_p50:.0f} | {r.snr_p50_db:.1f} dB |"
        )
    NOTES_OUT.write_text("\n".join(lines) + "\n")


if __name__ == "__main__":
    build_deck()
