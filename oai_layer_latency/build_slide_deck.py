#!/usr/bin/env python3
"""Build the OAI uplink latency investigation slide deck (.pptx) from the
validated findings and existing plots. Editable output for team/advisor sharing."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

AB = Path("/home/shr_aisvcs/workarea/carla_0_10_env/Carla-0.10.0-Linux-Shipping/PythonAPI/neu_collab/abiodun")
PLOTS = AB / "oai_layer_latency" / "plots"
OUT = AB / "oai_layer_latency" / "OAI_uplink_latency_investigation_deck.pptx"

NAVY = RGBColor(0x1B, 0x2A, 0x41)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
RED = RGBColor(0xD1, 0x49, 0x5B)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def title_bar(slide, text, sub=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    if sub:
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = sub
        rr.font.size = Pt(13); rr.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE5)


def bullets(slide, items, l=Inches(0.5), t=Inches(1.3), w=Inches(6.4), h=Inches(5.8), size=16):
    tf = box(slide, l, t, w, h)
    for i, (lvl, txt, *clr) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = clr[0] if clr else RGBColor(0x22, 0x22, 0x22)
        if lvl == 0:
            r.font.bold = True
        p.space_after = Pt(6)
    return tf


def image_fit(slide, path, l, t, maxw, maxh):
    path = Path(path)
    if not path.exists():
        return
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = maxw; h = Emu(int(w / ar))
    if h > maxh:
        h = maxh; w = Emu(int(h * ar))
    slide.shapes.add_picture(str(path), l + Emu(int((maxw - w) / 2)), t, width=w, height=h)


# ---------- 1. Title ----------
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tf = box(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.6))
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Localizing the OAI Uplink Bottleneck for Split-Inference"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = "Per-layer T-tracer instrumentation → UE RLC queue-wait, root-caused to the UL BLER/OLLA MCS selector"
r2.font.size = Pt(18); r2.font.color.rgb = RGBColor(0x9F, 0xB4, 0xCC)
p3 = tf.add_paragraph(); r3 = p3.add_run()
r3.text = "SCAN-AI × OAI  |  CARLA split-inference, 273PRB RFsim  |  2026-07-23"
r3.font.size = Pt(14); r3.font.color.rgb = RGBColor(0x76, 0x8A, 0xA3)

# ---------- 2. Problem & setup ----------
s = add_slide(); title_bar(s, "The problem: uplink dominates split-inference latency",
                           "Fixed deployment profile: no-AE, zstd, per-channel-u8, ROI 0, 200k radar PPS")
bullets(s, [
    (0, "Split inference sends a ~1 MB feature tensor per frame UE→edge over 5G (OAI)", NAVY),
    (1, "Prior study: edge compute ~7 ms, downlink result ~3 ms, but uplink handling ~150 ms — uplink is the bottleneck"),
    (0, "Question from the team/advisor", NAVY),
    (1, "WHICH layer (PHY → MAC → RLC → PDCP) holds the frame too long, and why?"),
    (1, "Is it processing time, queueing, buffer size, MCS, or traffic classification?"),
    (0, "Approach", NAVY),
    (1, "Timestamp the tensor as it enters the stack and break uplink transit down per layer"),
    (1, "Validate with iperf first, then the live CARLA split-inference pipeline"),
    (1, "273PRB, numerology 1, RFsim ideal channel (~50 dB PUSCH SNR)"),
], w=Inches(12.3))

# ---------- 3. Method ----------
s = add_slide(); title_bar(s, "Method: per-layer monotonic-timestamp T-tracer events",
                           "New OAI trace points bracket the RAN uplink transit")
bullets(s, [
    (0, "Added CLOCK_MONOTONIC-stamped T events at each layer boundary:", NAVY),
    (1, "NR_PDCP_TX_SDU  — UE PDCP ingress (t0)"),
    (1, "NR_RLC_TX_SDU / NR_RLC_TX_DEQUEUE  — RLC enqueue / dequeue-to-MAC"),
    (1, "GNB_MAC_RX_SDU / GNB_PDCP_RX_DELIVER  — gNB MAC in / PDCP egress (t_final)"),
    (1, "GNB_MAC_UL_MCS_DECISION, GNB_MAC_BLER_MCS_DECISION — scheduler pre/post MCS"),
    (0, "Key enabler", NAVY),
    (1, "RFsim runs UE+gNB on ONE host → monotonic clocks directly comparable"),
    (1, "DRB SDUs are FIFO → true per-packet transit + Little's-law residency cross-check"),
    (0, "Governing constraint (learned the hard way)", RED),
    (1, "T-tracer byte-compares T_messages.txt vs each softmodem → rebuild BOTH after any edit"),
], w=Inches(12.3))

# ---------- 4. Finding 1: RLC queue-wait ----------
s = add_slide(); title_bar(s, "Finding 1: the bottleneck is the UE RLC queue-wait",
                           "Two independent methods agree; PDCP/MAC/PHY/gNB are small")
bullets(s, [
    (0, "Per-packet RAN transit (PDCP-in → gNB-deliver): mean ~105 ms", NAVY),
    (1, "Little's law (occupancy/drain): ~103 ms — agrees"),
    (0, "Per-layer split:", NAVY),
    (1, "PDCP→RLC handoff: 0.1 ms"),
    (1, "RLC queue-wait: ~100 ms  = ~95%", RED),
    (1, "air (K2 3 ms) + gNB MAC/PHY/reassembly/PDCP: ~5 ms"),
    (0, "It is queueing, not processing", NAVY),
    (1, "A full ~1 MB frame sits in the RLC TX buffer waiting for grants"),
    (1, "RLC buffer = 10 MB/DRB, only ~11% used → buffer size is NOT the constraint"),
], w=Inches(6.1))
image_fit(s, PLOTS / "uplink_mcs_bottleneck_summary.png", Inches(6.7), Inches(1.35), Inches(6.4), Inches(5.7))

# ---------- 5. Finding 2: drain rate / low MCS ----------
s = add_slide(); title_bar(s, "Finding 2: the drain rate is capped by low UL MCS (QPSK)",
                           "Same ideal channel: CARLA sits in QPSK, iperf reaches 64QAM")
bullets(s, [
    (0, "Drain rate = RLC occupancy decay slope", NAVY),
    (1, "1 MB frame drains ~200 ms at QPSK vs near-instant at 64QAM"),
    (0, "QPSK caps spectral efficiency @273 PRB", NAVY),
    (1, "QPSK ~65, 16QAM ~172, 64QAM ~371 Mbps"),
    (0, "CARLA MCS ~4-8 (QPSK); iperf MCS 28 (64QAM)", RED),
    (1, "on the SAME RFsim ideal channel (~50 dB) — so it is not channel quality"),
], w=Inches(6.1))
image_fit(s, PLOTS / "uplink_drain_rate.png", Inches(6.6), Inches(1.6), Inches(6.6), Inches(4.6))

# ---------- 6. What it is NOT ----------
s = add_slide(); title_bar(s, "Root cause: ruling out the usual suspects",
                           "Adversarial checks, including the advisor's hypothesis")
bullets(s, [
    (0, "NOT traffic classification / QoS", NAVY),
    (1, "iperf, CARLA, and synthetic burst all use the same bearer: QFI 1 / 5QI 9 / DRB 1 / LCID 4 / LCG 1"),
    (1, "OAI UL MCS uses a single GLOBAL BLER target (no per-5QI/per-LC PER); SMF QER disabled"),
    (0, "NOT the SNR-PER-MCS table", NAVY),
    (1, "SINRx10_MCS_mapping (target BLER 1e-3): ~50 dB → MCS 16-28, never 8"),
    (0, "NOT power headroom (PHR helper)", NAVY),
    (1, "Direct pre/post trace: nr_ue_max_mcs_min_rb reduced MCS in 0 of 79,806 data rows"),
    (0, "NOT retransmissions", NAVY),
    (1, "Zero HARQ retx; no high-BLER decrease branch observed (ideal channel)"),
], w=Inches(12.3))

# ---------- 7. Root cause: OLLA cadence ----------
s = add_slide(); title_bar(s, "Root cause: the BLER/OLLA selector backs off under sparse cadence",
                           "get_mcs_from_bler(): num_dl_sched ≤ 3 forces MCS down")
bullets(s, [
    (0, "The OLLA update rule (gNB_scheduler_primitives.c):", NAVY),
    (1, "increase MCS  ⇐  low BLER AND num_dl_sched > 3"),
    (1, "DECREASE MCS  ⇐  high BLER OR num_dl_sched ≤ 3", RED),
    (0, "Sparse closed-loop bursts (~1.2 FPS) starve the BLER window", NAVY),
    (1, "num_sched p50 = 1 → 78.6% of updates hit the few-samples DECREASE branch"),
    (1, "→ MCS walks down to QPSK between bursts"),
    (0, "Dense open-loop (10 FPS): num_sched p50 = 3.5", NAVY),
    (1, "50% increase branch → MCS ratchets up to 23-25 (64QAM)"),
], w=Inches(6.3))
image_fit(s, PLOTS / "bler_olla_branch_comparison.png", Inches(6.9), Inches(1.5), Inches(6.2), Inches(4.9))

# ---------- 8. Proof: cadence-sensitive reproduction ----------
s = add_slide(); title_bar(s, "Proof: it is cadence, not payload size",
                           "Synthetic CARLA-shaped UDP burst reproduces it without CARLA/model")
bullets(s, [
    (0, "Same 1 MB burst size, two cadences (273PRB, adaptive MCS):", NAVY),
    (1, "open-loop 10 FPS  → MCS 23/25,  RAN UL p50 ~21 ms"),
    (1, "observed closed-loop pace ~1.2 FPS  → MCS 4/8,  RAN UL p50 ~118 ms", RED),
    (0, "“Large BSR” alone is not the trigger", NAVY),
    (1, "the low MCS appears only with the sparse cadence the CARLA closed loop creates"),
    (0, "Confirmed with the direct scheduler trace", NAVY),
    (1, "low MCS is selected upstream (OLLA), before the PHR/RB helper"),
], w=Inches(12.3))

# ---------- 9. Mitigations ----------
s = add_slide(); title_bar(s, "Mitigations: two independent levers",
                           "Fixed MCS proves the ceiling; payload reduction is deployable")
bullets(s, [
    (0, "Lever A — raise MCS (diagnostic): fixed MCS 28", NAVY),
    (1, "app RTT p50 186 → 47 ms (3.9x);  RAN UL 112 → 17 ms;  RLC queue 103 → 13 ms"),
    (1, "valid only on RFsim ideal channel; real fix = correct OLLA sparse-cadence behavior"),
    (0, "Lever B — reduce payload (deployable): AE / uint4", NAVY),
    (1, "AE-128 (152 KB): RTT p50 64 ms, delivery 99.8% (vs 76% at 1 MB)"),
    (1, "uint4 (394 KB): RTT p50 113 ms"),
    (1, "NOTE: payload reduction does NOT raise MCS (stays 2/5) — it just drains fewer bytes", RED),
    (0, "Best = both: small payload + healthy MCS", NAVY),
], w=Inches(6.1))
image_fit(s, PLOTS / "complementary_latency_summary.png", Inches(6.7), Inches(1.5), Inches(6.4), Inches(5.2))

# ---------- 10. Conclusions ----------
s = add_slide(); title_bar(s, "Conclusions & next steps")
bullets(s, [
    (0, "The uplink bottleneck is the UE RLC queue-wait (~95% of transit)", NAVY),
    (1, "caused by low UL MCS (QPSK), NOT buffer size, PDCP, PHR, retx, or traffic class"),
    (0, "Mechanism: OLLA decreases MCS when a BLER window has ≤3 scheduled TX", RED),
    (1, "sparse closed-loop split-inference cadence starves the selector → QPSK"),
    (0, "Recommended", NAVY),
    (1, "Deployable now: feature compression (AE) — fewer bytes + higher delivery"),
    (1, "Scheduler fix: make OAI OLLA robust to large sparse low-latency bursts, then"),
    (1, "re-test under a realistic Sionna fading channel (RFsim here is ideal/flat)"),
    (0, "Artifacts: abiodun/oai_layer_latency/ (README + analyzer + plots + run scripts)", GREY),
], w=Inches(12.3))

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
