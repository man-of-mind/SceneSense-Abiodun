#!/usr/bin/env python3
"""Create OD-only vs SEG-only loopback/OAI latency plots and deck."""

from __future__ import annotations

import argparse
import csv
import json
import math
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

os.environ.setdefault("MPLCONFIGDIR", "/tmp/scenesense_mplconfig")

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ABIODUN_DIR = Path(__file__).resolve().parents[1]
DEFAULT_ANALYSIS_DIR = ABIODUN_DIR / "metrics_logs" / "scenesense_analysis" / "camera_od_seg_latency_20260604"
DEFAULT_CAMERA_ANALYSIS_DIR = ABIODUN_DIR / "metrics_logs" / "month1_camera_latency_analysis"
GNB_CONF = ABIODUN_DIR / "OAI" / "openairinterface5g" / "targets" / "PROJECTS" / "GENERIC-NR-5GC" / "CONF" / "gnb.sa.band78.fr1.106PRB.usrpb210.conf"
UE_CONF = ABIODUN_DIR / "OAI" / "openairinterface5g" / "targets" / "PROJECTS" / "GENERIC-NR-5GC" / "CONF" / "ue.conf"
CN_CONF = ABIODUN_DIR / "OAI" / "oai-cn5g" / "conf" / "config.yaml"
SCRIPT_CONFIG = ABIODUN_DIR / "scripts" / "config.env"

COLORS = {
    "navy": RGBColor(23, 50, 77),
    "teal": RGBColor(0, 124, 137),
    "blue": RGBColor(47, 111, 173),
    "orange": RGBColor(213, 106, 33),
    "green": RGBColor(61, 122, 70),
    "red": RGBColor(178, 58, 72),
    "purple": RGBColor(112, 87, 163),
    "gray": RGBColor(243, 246, 248),
    "midgray": RGBColor(216, 224, 231),
    "muted": RGBColor(99, 115, 129),
    "dark": RGBColor(32, 42, 51),
    "white": RGBColor(255, 255, 255),
}

PLOT_COLORS = {
    "OD Loopback": "#4C78A8",
    "OD OAI": "#F58518",
    "SEG Loopback": "#54A24B",
    "SEG OAI": "#E45756",
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Read camera-only OD/SEG analyzer JSON outputs and build Month 1 "
            "loopback-vs-OAI latency plots plus a PowerPoint deck."
        )
    )
    parser.add_argument("--od-json", default="", help="OD analyzer JSON path.")
    parser.add_argument("--seg-json", default="", help="SEG analyzer JSON path.")
    parser.add_argument("--output-dir", default=str(DEFAULT_ANALYSIS_DIR), help="Output folder.")
    parser.add_argument("--title", default="Camera-only OD vs SEG latency over loopback and OAI 5G")
    return parser.parse_args()


def latest_json(pattern: str) -> Optional[Path]:
    paths = sorted(DEFAULT_CAMERA_ANALYSIS_DIR.glob(pattern), key=lambda path: path.stat().st_mtime)
    return paths[-1] if paths else None


def read_json_prefix(path: Path) -> Dict[str, object]:
    text = path.read_text(encoding="utf-8")
    start = text.find("{")
    if start < 0:
        raise ValueError(f"{path} does not contain a JSON object")
    decoder = json.JSONDecoder()
    payload, _end = decoder.raw_decode(text[start:])
    if not isinstance(payload, dict):
        raise ValueError(f"{path} did not decode to a JSON object")
    return payload


def safe_float(value: object) -> float:
    try:
        result = float(value)
    except (TypeError, ValueError):
        return float("nan")
    return result if math.isfinite(result) else float("nan")


def safe_int(value: object) -> int:
    value_f = safe_float(value)
    return int(value_f) if math.isfinite(value_f) else 0


def fmt(value: object, digits: int = 1, suffix: str = "") -> str:
    value_f = safe_float(value)
    if not math.isfinite(value_f):
        return "n/a"
    return f"{value_f:.{digits}f}{suffix}"


def classify_source(source: str) -> Optional[Tuple[str, str, str]]:
    lower = source.lower()
    if source == "OVERALL":
        return None
    task = "SEG" if "seg" in lower else "OD" if "od" in lower else ""
    transport = "OAI" if "_oai_" in lower else "Loopback" if "_loopback_" in lower else ""
    if not task or not transport:
        return None
    label = f"{task} {transport}"
    return task, transport, label


def load_rows(payloads: Sequence[Dict[str, object]]) -> List[Dict[str, object]]:
    rows: List[Dict[str, object]] = []
    for payload in payloads:
        for summary in payload.get("summaries", []):  # type: ignore[union-attr]
            if not isinstance(summary, dict):
                continue
            source = str(summary.get("source", ""))
            classified = classify_source(source)
            if classified is None:
                continue
            task, transport, label = classified
            frames = safe_int(summary.get("frames_total"))
            returned = safe_int(summary.get("round_trip_ms_count"))
            if returned <= 0:
                returned = safe_int(summary.get("frames_with_mask"))
            receive_rate = returned / frames if frames else float("nan")
            rtt_median = safe_float(summary.get("round_trip_ms_median"))
            row = {
                "task": task,
                "transport": transport,
                "label": label,
                "source": source,
                "frames_total": frames,
                "returned_frames": returned,
                "receive_rate": receive_rate,
                "missed_frames": max(0, frames - returned),
                "front_ms_median": safe_float(summary.get("front_ms_median")),
                "front_ms_p95": safe_float(summary.get("front_ms_p95")),
                "back_ms_median": safe_float(summary.get("back_ms_median")),
                "back_ms_p95": safe_float(summary.get("back_ms_p95")),
                "round_trip_ms_median": rtt_median,
                "round_trip_ms_p95": safe_float(summary.get("round_trip_ms_p95")),
                "round_trip_ms_mean": safe_float(summary.get("round_trip_ms_mean")),
                "payload_kib_median": safe_float(summary.get("payload_kib_median")),
                "payload_kib_p95": safe_float(summary.get("payload_kib_p95")),
                "payload_kib_mean": safe_float(summary.get("payload_kib_mean")),
                "payload_chunks_median": safe_float(summary.get("payload_chunks_median")),
                "payload_chunks_p95": safe_float(summary.get("payload_chunks_p95")),
                "app_total_ms_median": safe_float(summary.get("front_ms_median")) + rtt_median,
            }
            row["transport_residual_ms_median"] = max(
                0.0,
                safe_float(row["round_trip_ms_median"]) - safe_float(row["back_ms_median"]),
            )
            rows.append(row)
    order = {"OD Loopback": 0, "OD OAI": 1, "SEG Loopback": 2, "SEG OAI": 3}
    rows.sort(key=lambda row: order.get(str(row["label"]), 99))
    return rows


def read_env_value(path: Path, name: str) -> str:
    if not path.exists():
        return ""
    pattern = re.compile(rf"^\s*export\s+{re.escape(name)}=(.*)$")
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        match = pattern.match(line)
        if not match:
            continue
        value = match.group(1).split("#", 1)[0].strip()
        return value.strip('"')
    return ""


def read_conf_number(path: Path, name: str) -> str:
    if not path.exists():
        return ""
    pattern = re.compile(rf"\b{re.escape(name)}\s*=\s*([^;,\n]+)")
    text = path.read_text(encoding="utf-8", errors="ignore")
    match = pattern.search(text)
    if not match:
        return ""
    return match.group(1).strip().strip('"')


def read_yaml_after(path: Path, marker: str, name: str) -> str:
    if not path.exists():
        return ""
    lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    start = 0
    for index, line in enumerate(lines):
        if marker in line:
            start = index
            break
    pattern = re.compile(rf"^\s*{re.escape(name)}:\s*(.+)$")
    for line in lines[start:]:
        match = pattern.match(line)
        if match:
            return match.group(1).strip().strip('"')
    return ""


def radio_config_rows() -> List[Tuple[str, str]]:
    numerology = read_env_value(SCRIPT_CONFIG, "UE_NUMEROLOGY") or read_conf_number(GNB_CONF, "subcarrierSpacing")
    scs = "30 kHz" if numerology == "1" else f"numerology {numerology}"
    prb = read_env_value(SCRIPT_CONFIG, "UE_PRB") or read_conf_number(GNB_CONF, "dl_carrierBandwidth")
    dl_freq_hz = read_env_value(SCRIPT_CONFIG, "UE_DL_FREQ") or ""
    dl_freq = f"{float(dl_freq_hz) / 1e9:.4f} GHz" if dl_freq_hz else "n/a"
    band = read_env_value(SCRIPT_CONFIG, "UE_BAND") or read_conf_number(GNB_CONF, "dl_frequencyBand")
    period = read_conf_number(GNB_CONF, "dl_UL_TransmissionPeriodicity")
    period_label = "5 ms" if period == "6" else f"enum {period}"
    dl_slots = read_conf_number(GNB_CONF, "nrofDownlinkSlots")
    ul_slots = read_conf_number(GNB_CONF, "nrofUplinkSlots")
    dl_symbols = read_conf_number(GNB_CONF, "nrofDownlinkSymbols")
    ul_symbols = read_conf_number(GNB_CONF, "nrofUplinkSymbols")
    five_qi = read_yaml_after(CN_CONF, 'dnn: "oai"', "5qi") or "9"
    ambr_ul = read_yaml_after(CN_CONF, 'dnn: "oai"', "session_ambr_ul") or "10Gbps"
    ambr_dl = read_yaml_after(CN_CONF, 'dnn: "oai"', "session_ambr_dl") or "10Gbps"
    imsi = read_conf_number(UE_CONF, "imsi")
    ue_ip = read_env_value(SCRIPT_CONFIG, "OAI_UE_IP") or "10.0.0.2"
    rx_ip = read_env_value(SCRIPT_CONFIG, "OAI_RX_IP") or "192.168.70.140"
    return [
        ("RAN mode", "NR SA over OAI RFsim, single UE for OD/SEG latency runs"),
        ("Carrier", f"band n{band}, DL {dl_freq}, {prb} PRB at {scs} SCS"),
        ("Bandwidth", f"{prb} PRB = about 40 MHz channel class (38.16 MHz RB span)"),
        ("TDD pattern", f"{period_label}: {dl_slots} DL slots + mixed slot ({dl_symbols} DL sym, {ul_symbols} UL sym) + {ul_slots} UL slots"),
        ("Core QoS", f"DNN oai, S-NSSAI SST=1, 5QI={five_qi}, AMBR UL/DL {ambr_ul}/{ambr_dl}"),
        ("UE path", f"IMSI {imsi}, oaitun_ue1 {ue_ip} -> perception RX {rx_ip}"),
    ]


def write_summary(rows: Sequence[Dict[str, object]], out_dir: Path) -> None:
    fields = [
        "task",
        "transport",
        "source",
        "frames_total",
        "returned_frames",
        "receive_rate",
        "missed_frames",
        "front_ms_median",
        "back_ms_median",
        "round_trip_ms_median",
        "round_trip_ms_p95",
        "app_total_ms_median",
        "payload_kib_median",
        "payload_kib_p95",
        "payload_chunks_median",
    ]
    with (out_dir / "camera_od_seg_latency_summary.csv").open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=fields)
        writer.writeheader()
        for row in rows:
            writer.writerow({field: row.get(field, "") for field in fields})

    lookup = {str(row["label"]): row for row in rows}
    od_loop = lookup.get("OD Loopback", {})
    od_oai = lookup.get("OD OAI", {})
    seg_loop = lookup.get("SEG Loopback", {})
    seg_oai = lookup.get("SEG OAI", {})
    lines = [
        "# Camera-Only OD vs SEG Latency Comparison",
        "",
        f"Generated: {datetime.now().isoformat(timespec='seconds')}",
        "",
        "## Headline",
        "",
        f"- OD median RTT: loopback {fmt(od_loop.get('round_trip_ms_median'))} ms, OAI {fmt(od_oai.get('round_trip_ms_median'))} ms.",
        f"- SEG median RTT: loopback {fmt(seg_loop.get('round_trip_ms_median'))} ms, OAI {fmt(seg_oai.get('round_trip_ms_median'))} ms.",
        f"- OD payload median: {fmt(od_loop.get('payload_kib_median'))} KiB; SEG payload median: {fmt(seg_loop.get('payload_kib_median'))} KiB.",
        f"- SEG/OD payload ratio: {fmt(safe_float(seg_loop.get('payload_kib_median')) / safe_float(od_loop.get('payload_kib_median')), 2)}x on loopback.",
        "",
        "## OAI Config",
        "",
    ]
    lines.extend(f"- {name}: {value}" for name, value in radio_config_rows())
    lines.append("")
    lines.append("## Per-Run Summary")
    lines.append("")
    lines.append("| task | transport | frames | receive | median RTT ms | p95 RTT ms | median payload KiB |")
    lines.append("|---|---:|---:|---:|---:|---:|---:|")
    for row in rows:
        lines.append(
            f"| {row['task']} | {row['transport']} | {row['frames_total']} | "
            f"{safe_float(row['receive_rate']) * 100:.1f}% | "
            f"{fmt(row['round_trip_ms_median'])} | {fmt(row['round_trip_ms_p95'])} | "
            f"{fmt(row['payload_kib_median'])} |"
        )
    (out_dir / "camera_od_seg_latency_summary.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def plot_rtt(rows: Sequence[Dict[str, object]], out_dir: Path) -> Path:
    lookup = {str(row["label"]): row for row in rows}
    task_labels = ["OD", "SEG"]
    loopback = [
        safe_float(lookup[f"{task} Loopback"]["round_trip_ms_median"])
        for task in task_labels
    ]
    oai = [
        safe_float(lookup[f"{task} OAI"]["round_trip_ms_median"])
        for task in task_labels
    ]
    loopback_p95 = [
        safe_float(lookup[f"{task} Loopback"]["round_trip_ms_p95"])
        for task in task_labels
    ]
    oai_p95 = [
        safe_float(lookup[f"{task} OAI"]["round_trip_ms_p95"])
        for task in task_labels
    ]
    x = np.arange(len(task_labels))
    fig, ax = plt.subplots(figsize=(10.6, 5.6))
    width = 0.34
    loop_x = x - width / 2
    oai_x = x + width / 2
    ax.bar(loop_x, loopback, width, label="Loopback median", color="#4C78A8")
    ax.bar(oai_x, oai, width, label="OAI median", color="#F58518")
    ax.scatter(loop_x, loopback_p95, color="#17324D", marker="D", s=48, label="Loopback p95")
    ax.scatter(oai_x, oai_p95, color="#8F3B00", marker="D", s=48, label="OAI p95")
    for xpos, median, p95 in zip(loop_x, loopback, loopback_p95):
        ax.vlines(xpos, median, p95, color="#17324D", linewidth=1.5, alpha=0.75)
    for xpos, median, p95 in zip(oai_x, oai, oai_p95):
        ax.vlines(xpos, median, p95, color="#8F3B00", linewidth=1.5, alpha=0.75)
    ax.set_xticks(x, task_labels)
    ax.set_ylabel("round-trip latency (ms)")
    ax.set_title("Camera-only split inference: Loopback vs OAI by task")
    ax.grid(axis="y", alpha=0.22)
    ax.legend(frameon=False, ncol=2)
    ax.set_ylim(0, max([*loopback_p95, *oai_p95]) * 1.18)
    for xpos, median, p95 in zip(loop_x, loopback, loopback_p95):
        ax.text(xpos, median + 3, f"{median:.1f}", ha="center", va="bottom", fontsize=9)
        ax.text(xpos, p95 + 3, f"p95 {p95:.1f}", ha="center", va="bottom", fontsize=8, color="#17324D")
    for xpos, median, p95 in zip(oai_x, oai, oai_p95):
        ax.text(xpos, median + 3, f"{median:.1f}", ha="center", va="bottom", fontsize=9)
        ax.text(xpos, p95 + 3, f"p95 {p95:.1f}", ha="center", va="bottom", fontsize=8, color="#8F3B00")
    fig.tight_layout()
    path = out_dir / "camera_latency_rtt_median_p95.png"
    fig.savefig(path, dpi=180)
    plt.close(fig)
    return path


def plot_payload(rows: Sequence[Dict[str, object]], out_dir: Path) -> Path:
    labels = [str(row["label"]) for row in rows]
    values = [safe_float(row["payload_kib_median"]) for row in rows]
    chunks = [safe_float(row["payload_chunks_median"]) for row in rows]
    colors = [PLOT_COLORS.get(label, "#4C78A8") for label in labels]
    fig, ax = plt.subplots(figsize=(10.2, 5.2))
    bars = ax.bar(labels, values, color=colors)
    ax.set_ylabel("median feature payload (KiB)")
    ax.set_title("SEG sends about 4.6x more feature payload than OD", pad=18)
    ax.set_ylim(0, max(values) * 1.12)
    ax.grid(axis="y", alpha=0.22)
    ax.tick_params(axis="x", rotation=15)
    for bar, value, chunk in zip(bars, values, chunks):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            value - max(values) * 0.04,
            f"{value:.1f} KiB\n{chunk:.0f} chunks",
            ha="center",
            va="top",
            fontsize=9,
            fontweight="bold",
            color="white",
        )
    fig.tight_layout()
    path = out_dir / "camera_latency_payload_chunks.png"
    fig.savefig(path, dpi=180)
    plt.close(fig)
    return path


def plot_components(rows: Sequence[Dict[str, object]], out_dir: Path) -> Path:
    labels = [str(row["label"]) for row in rows]
    front = np.asarray([safe_float(row["front_ms_median"]) for row in rows])
    residual = np.asarray([safe_float(row["transport_residual_ms_median"]) for row in rows])
    back = np.asarray([safe_float(row["back_ms_median"]) for row in rows])
    x = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(10.8, 5.5))
    ax.bar(x, front, label="front model", color="#4C78A8")
    ax.bar(x, residual, bottom=front, label="transport/wait residual", color="#F58518")
    ax.bar(x, back, bottom=front + residual, label="back model", color="#54A24B")
    ax.set_xticks(x, labels, rotation=15, ha="right")
    ax.set_ylabel("median glass-to-result components (ms)")
    ax.set_title("OAI penalty appears mainly in the split transport/wait segment")
    ax.grid(axis="y", alpha=0.22)
    ax.legend(frameon=False)
    totals = front + residual + back
    ax.set_ylim(0, max(totals) * 1.14)
    for index, total in enumerate(totals):
        ax.text(index, total + 3, f"{total:.1f}", ha="center", va="bottom", fontsize=9)
    fig.tight_layout()
    path = out_dir / "camera_latency_component_stack.png"
    fig.savefig(path, dpi=180)
    plt.close(fig)
    return path


def add_title(slide, title: str, subtitle: str = "") -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.55))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["navy"]
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.88), Inches(11.6), Inches(0.35))
        tf = sub.text_frame
        tf.text = subtitle
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = COLORS["muted"]


def add_textbox(slide, x: float, y: float, w: float, h: float, text: str, size: int = 14, color=COLORS["dark"], bold: bool = False, fill=None) -> None:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold


def add_metric_card(slide, x: float, y: float, w: float, h: float, title: str, value: str, note: str, color) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["gray"]
    shape.line.color.rgb = COLORS["midgray"]
    tf = shape.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.text = title
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = COLORS["muted"]
    p = tf.add_paragraph()
    p.text = value
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    p = tf.add_paragraph()
    p.text = note
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["muted"]


def add_bullets(slide, x: float, y: float, w: float, h: float, bullets: Sequence[str], size: int = 15) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.clear()
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["dark"]
        p.space_after = Pt(7)


def add_table(slide, x: float, y: float, w: float, h: float, rows: Sequence[Tuple[str, str]]) -> None:
    table_shape = slide.shapes.add_table(len(rows), 2, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    table.columns[0].width = Inches(w * 0.27)
    table.columns[1].width = Inches(w * 0.73)
    for row_idx, (name, value) in enumerate(rows):
        for col_idx, text in enumerate((name, value)):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["gray"] if row_idx % 2 == 0 else COLORS["white"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = COLORS["navy"] if col_idx == 0 else COLORS["dark"]
                paragraph.font.bold = col_idx == 0


def image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def create_deck(rows: Sequence[Dict[str, object]], plots: Dict[str, Path], out_dir: Path, title: str) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    lookup = {str(row["label"]): row for row in rows}
    od_loop = lookup["OD Loopback"]
    od_oai = lookup["OD OAI"]
    seg_loop = lookup["SEG Loopback"]
    seg_oai = lookup["SEG OAI"]
    od_penalty = safe_float(od_oai["round_trip_ms_median"]) - safe_float(od_loop["round_trip_ms_median"])
    seg_penalty = safe_float(seg_oai["round_trip_ms_median"]) - safe_float(seg_loop["round_trip_ms_median"])
    payload_ratio = safe_float(seg_loop["payload_kib_median"]) / safe_float(od_loop["payload_kib_median"])

    slide = prs.slides.add_slide(blank)
    add_title(slide, title, "Month 1 latency-only split-inference runs; task GT disabled during these traces.")
    add_metric_card(slide, 0.8, 1.45, 2.65, 1.25, "OD OAI median RTT", fmt(od_oai["round_trip_ms_median"], 1, " ms"), f"+{od_penalty:.1f} ms vs loopback", COLORS["orange"])
    add_metric_card(slide, 3.75, 1.45, 2.65, 1.25, "SEG OAI median RTT", fmt(seg_oai["round_trip_ms_median"], 1, " ms"), f"+{seg_penalty:.1f} ms vs loopback", COLORS["red"])
    add_metric_card(slide, 6.7, 1.45, 2.65, 1.25, "SEG / OD payload", f"{payload_ratio:.1f}x", "median feature payload", COLORS["teal"])
    add_metric_card(slide, 9.65, 1.45, 2.65, 1.25, "Runs", "4 x ~180s", "OD/SEG over loopback/OAI", COLORS["green"])
    add_bullets(
        slide,
        0.9,
        3.05,
        5.7,
        2.8,
        [
            "Same front/back split pattern as the fusion OAI comparison, now isolated to camera-only OD and camera-only SEG.",
            "OD sends compact Faster R-CNN backbone features at 720p; SEG sends denser LR-ASPP features from the 1080p camera route.",
            "OAI increases RTT for both tasks; the larger SEG payload makes the OAI penalty more visible.",
        ],
        size=17,
    )
    add_table(slide, 7.0, 3.0, 5.6, 2.65, radio_config_rows())
    add_textbox(slide, 0.85, 6.55, 11.8, 0.35, "Interpretation: these traces measure transport/runtime behavior. OD/SEG task-quality numbers are taken from separate GT-enabled runs.", size=12, color=COLORS["muted"])

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Round-Trip Latency", "Median and p95 feature-send-to-result-return latency.")
    image(slide, plots["rtt"], 0.75, 1.15, 7.35, 4.65)
    add_bullets(
        slide,
        8.35,
        1.35,
        4.15,
        4.4,
        [
            f"OD median RTT rises from {fmt(od_loop['round_trip_ms_median'])} ms to {fmt(od_oai['round_trip_ms_median'])} ms over OAI.",
            f"SEG median RTT rises from {fmt(seg_loop['round_trip_ms_median'])} ms to {fmt(seg_oai['round_trip_ms_median'])} ms over OAI.",
            f"SEG OAI p95 reaches {fmt(seg_oai['round_trip_ms_p95'])} ms, the largest tail latency in this four-run comparison.",
            "The OAI path is functional but not latency-neutral; it should be modeled as a control constraint.",
        ],
        size=16,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Payload Pressure", "Feature payload size and UDP chunk count explain why SEG stresses the link more.")
    image(slide, plots["payload"], 0.75, 1.15, 7.15, 4.4)
    add_bullets(
        slide,
        8.15,
        1.35,
        4.25,
        4.4,
        [
            f"OD median payload is about {fmt(od_loop['payload_kib_median'])} KiB and 2 chunks.",
            f"SEG median payload is about {fmt(seg_loop['payload_kib_median'])} KiB and 7 chunks.",
            f"That is roughly {payload_ratio:.1f}x larger than OD before any radio scheduling effects.",
            "This gives the controller a clear task-dependent byte/latency tradeoff.",
        ],
        size=16,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Latency Components", "Median front compute plus split RTT components.")
    image(slide, plots["components"], 0.7, 1.05, 7.45, 4.75)
    add_bullets(
        slide,
        8.35,
        1.35,
        4.1,
        4.5,
        [
            f"SEG front-side compute is higher than OD: {fmt(seg_loop['front_ms_median'])} ms vs {fmt(od_loop['front_ms_median'])} ms in loopback.",
            "Back-side segmentation compute is small; most SEG cost is feature extraction plus payload transport.",
            "OAI adds a larger transport/wait residual than loopback for both tasks.",
            "This decomposition supports an RL state with task type, payload size, RTT trend, and radio grant state.",
        ],
        size=16,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Takeaway For Month 1", "Camera-only OD/SEG now have the same latency story as fusion.")
    add_bullets(
        slide,
        0.85,
        1.25,
        6.05,
        4.85,
        [
            "Loopback establishes the compute/runtime floor for each task.",
            "OAI adds a repeatable transport penalty, and the penalty grows with payload pressure.",
            "SEG is not just a different model: it is a different network load class, with larger feature tensors and more UDP chunks.",
            "Current core QoS is default 5QI 9 eMBB; this is a useful baseline before 5QI, scheduling, or action-policy experiments.",
        ],
        size=19,
    )
    add_table(slide, 7.25, 1.25, 5.25, 3.1, [
        ("Evidence", "OD loopback/OAI and SEG loopback/OAI 180s latency traces"),
        ("Plots", "RTT median/p95, payload/chunks, front/back/transport components"),
        ("Quality tie-in", "Use separate GT-enabled OD/SEG summaries for recall/mIoU"),
        ("Next", "Pull raw CSVs for CDF/jitter plots; align with tunnel/RAN traces when available"),
    ])
    add_textbox(slide, 7.25, 4.75, 5.25, 0.9, "Bottom line: task choice changes both model-side compute and network-side pressure, so the future controller should not treat OD, SEG, and fusion as one generic stream.", size=15, color=COLORS["navy"], bold=True, fill=COLORS["gray"])

    out_path = out_dir / "SceneSense_Camera_OD_SEG_Latency_Comparison.pptx"
    prs.save(out_path)
    return out_path


def main() -> None:
    args = parse_args()
    out_dir = Path(args.output_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    od_json = Path(args.od_json).expanduser() if args.od_json else latest_json("month1_latency_od_*.json")
    seg_json = Path(args.seg_json).expanduser() if args.seg_json else latest_json("month1_latency_seg_*.json")
    if od_json is None or seg_json is None:
        raise SystemExit("Provide --od-json and --seg-json, or copy analyzer JSONs into metrics_logs/month1_camera_latency_analysis.")

    payloads = [read_json_prefix(od_json), read_json_prefix(seg_json)]
    rows = load_rows(payloads)
    labels = {str(row["label"]) for row in rows}
    required = {"OD Loopback", "OD OAI", "SEG Loopback", "SEG OAI"}
    missing = sorted(required - labels)
    if missing:
        raise SystemExit(f"Missing required summary rows: {', '.join(missing)}")

    write_summary(rows, out_dir)
    plots = {
        "rtt": plot_rtt(rows, out_dir),
        "payload": plot_payload(rows, out_dir),
        "components": plot_components(rows, out_dir),
    }
    deck = create_deck(rows, plots, out_dir, args.title)
    print(f"Wrote: {out_dir / 'camera_od_seg_latency_summary.csv'}")
    print(f"Wrote: {out_dir / 'camera_od_seg_latency_summary.md'}")
    for path in plots.values():
        print(f"Wrote: {path}")
    print(f"Wrote: {deck}")


if __name__ == "__main__":
    main()
